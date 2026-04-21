"""PDF manipulation tools — merge, split, compress, rotate, protect, unlock, etc."""
from __future__ import annotations

import io
import logging
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Callable, Sequence


logger = logging.getLogger(__name__)


# ===========================================================================
# Universal Conversion Pipeline
# ===========================================================================
#
# Every "convert something → PDF" tool runs through _conversion_pipeline():
#
#   1. Try each engine function in priority order.
#   2. After each successful attempt, run _qa_check_pdf() if Ollama is up.
#      QA renders up to 3 pages and asks Gemma 4 to check for:
#      text, fonts, colors, shapes, tables, images, layout integrity.
#   3. First engine that PASSES QA wins — file is copied to output_path.
#   4. If all fail QA (or Ollama is offline), use the first engine that
#      produced a non-empty file rather than crashing.
#
# ===========================================================================

_QA_PROMPT = (
    "You are performing quality assurance on a PDF page that was converted "
    "from an Office document or image.\n\n"
    "Carefully inspect the page and check ALL of the following:\n"
    "1. TEXT — Is all text visible, readable, and not replaced by boxes or garbled characters?\n"
    "2. FONTS — Do font sizes and weights (bold, italic) look intentional?\n"
    "3. COLORS — Are text and background colors rendered (not just black-on-white if the original had color)?\n"
    "4. SHAPES & GRAPHICS — Are shapes, lines, arrows, or decorative elements present if they should be?\n"
    "5. IMAGES & LOGOS — Are embedded images/logos visible? They may appear pixelated at screen resolution — FAIL only if they are completely blank or missing placeholders.\n"
    "6. TABLES — Are table borders and cell content intact?\n"
    "7. LAYOUT — Is indentation, spacing, and alignment reasonable? Page is not blank or overflowing?\n\n"
    "Reply with exactly:\n"
    "PASS\n"
    "or:\n"
    "FAIL: <reason in ≤20 words>\n\n"
    "Do not add any other text."
)


def _qa_check_pdf(pdf_path: Path) -> tuple[bool, str]:
    """Render up to 3 pages and ask Ollama to QA each one.

    Returns (passed, notes). Returns (True, 'skipped') if Ollama is offline.
    """
    try:
        _ollama_check()
    except Exception:
        return True, "Ollama offline — QA skipped"

    import fitz
    try:
        doc = fitz.open(str(pdf_path))
    except Exception as e:
        return False, f"Cannot open output PDF: {e}"

    n = len(doc)
    if n == 0:
        return False, "Output PDF has 0 pages"

    # Sample: first page, middle page, last page (deduplicated)
    indices = sorted({0, n // 2, n - 1})

    mat = fitz.Matrix(2.0, 2.0)          # ~144 DPI — good enough for images/charts
    issues: list[str] = []

    with tempfile.TemporaryDirectory() as td:
        for idx in indices:
            pix = doc[idx].get_pixmap(matrix=mat)
            img = Path(td) / f"qa_page_{idx}.png"
            pix.save(str(img))
            try:
                resp = _ollama_call_page(img, _QA_PROMPT).strip()
            except Exception as e:
                return True, f"QA call failed ({e}) — skipped"

            first = resp.split("\n")[0].strip().upper()
            if first.startswith("FAIL"):
                reason = resp[5:].strip().lstrip(":").strip()
                issues.append(f"p{idx + 1}: {reason[:80]}")

    if issues:
        return False, "; ".join(issues)
    return True, f"QA passed ({len(indices)} page(s) checked)"


def _conversion_pipeline(
    engines: list[tuple[str, Callable[[Path], None]]],
    output_path: Path,
) -> tuple[str, str]:
    """Try each engine in order. QA after each. Return (engine_name, qa_notes).

    - First engine whose output passes Ollama QA → used immediately.
    - If Ollama is offline → first engine that produces a non-empty file wins.
    - If all engines fail QA → first successful output is used with a warning.
    """
    tmp = Path(tempfile.mkdtemp(prefix="cf_pipe_"))
    first_ok: tuple[str, Path, str] | None = None

    try:
        for name, fn in engines:
            candidate = tmp / f"{name}.pdf"
            try:
                fn(candidate)
            except Exception:
                logger.exception(
                    "PDF conversion engine failed: engine=%s candidate=%s output=%s",
                    name,
                    candidate,
                    output_path,
                )
                continue                                   # engine failed, try next

            if not candidate.exists() or candidate.stat().st_size == 0:
                logger.warning(
                    "PDF conversion engine produced empty output: engine=%s candidate=%s output=%s",
                    name,
                    candidate,
                    output_path,
                )
                continue

            passed, notes = _qa_check_pdf(candidate)

            if first_ok is None:
                first_ok = (name, candidate, notes)

            if passed:
                shutil.copy2(str(candidate), str(output_path))
                return name, notes

        # Nothing passed QA (or Ollama offline) — use first successful output
        if first_ok:
            name, candidate, notes = first_ok
            shutil.copy2(str(candidate), str(output_path))
            return name, f"Best available — QA note: {notes}"

        raise RuntimeError(
            "All conversion engines failed. "
            "Check that the input file is not corrupted."
        )
    finally:
        shutil.rmtree(str(tmp), ignore_errors=True)


# ---------------------------------------------------------------------------
# Merge
# ---------------------------------------------------------------------------

def merge_pdfs(input_paths: Sequence[Path], output_path: Path) -> None:
    from pypdf import PdfWriter
    writer = PdfWriter()
    for path in input_paths:
        writer.append(str(path))
    with output_path.open("wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# Split  (one PDF per page, returned as list of bytes; caller saves them)
# ---------------------------------------------------------------------------

def split_pdf(input_path: Path, output_dir: Path) -> list[Path]:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(str(input_path))
    output_paths: list[Path] = []
    stem = input_path.stem
    for i, page in enumerate(reader.pages, 1):
        writer = PdfWriter()
        writer.add_page(page)
        out = output_dir / f"{stem}_page_{i}.pdf"
        with out.open("wb") as f:
            writer.write(f)
        output_paths.append(out)
    return output_paths


# ---------------------------------------------------------------------------
# Extract specific pages  (1-indexed, comma-separated or range "1,3,5-8")
# ---------------------------------------------------------------------------

def extract_pages(input_path: Path, page_spec: str, output_path: Path) -> None:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    indices = _parse_page_spec(page_spec, total)
    writer = PdfWriter()
    for idx in indices:
        writer.add_page(reader.pages[idx])
    with output_path.open("wb") as f:
        writer.write(f)


def remove_pages(input_path: Path, page_spec: str, output_path: Path) -> None:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    remove_set = set(_parse_page_spec(page_spec, total))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i not in remove_set:
            writer.add_page(page)
    with output_path.open("wb") as f:
        writer.write(f)


def _parse_page_spec(spec: str, total: int) -> list[int]:
    """Parse '1,3,5-8' into 0-indexed list. Raises ValueError on bad input."""
    indices: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            halves = part.split("-", 1)
            if not halves[0].strip().isdigit() or not halves[1].strip().isdigit():
                raise ValueError(f"Invalid page range: '{part}'. Use numbers like '2-5'.")
            a, b = int(halves[0]), int(halves[1])
            if a < 1 or b < a:
                raise ValueError(f"Invalid page range: '{part}'. Start must be >= 1 and <= end.")
            if b > total:
                raise ValueError(
                    f"Page range '{part}' exceeds the document length ({total} page(s))."
                )
            indices.extend(range(a - 1, b))
        elif part.isdigit():
            indices.append(int(part) - 1)
        else:
            raise ValueError(f"Invalid page number: '{part}'. Use digits and ranges like '1,3,5-8'.")
    # clamp and deduplicate preserving order
    seen: set[int] = set()
    result: list[int] = []
    for i in indices:
        if 0 <= i < total and i not in seen:
            seen.add(i)
            result.append(i)
    if not result:
        raise ValueError(f"No valid pages found in '{spec}'. The document has {total} page(s).")
    return result


# ---------------------------------------------------------------------------
# Rotate
# ---------------------------------------------------------------------------

def rotate_pdf(input_path: Path, degrees: int, output_path: Path) -> None:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(str(input_path))
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(degrees)
        writer.add_page(page)
    with output_path.open("wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# Compress  (uses pikepdf for lossless stream compression)
# ---------------------------------------------------------------------------

def compress_pdf(input_path: Path, output_path: Path) -> None:
    import pikepdf
    with pikepdf.open(str(input_path)) as pdf:
        pdf.save(
            str(output_path),
            compress_streams=True,
            object_stream_mode=pikepdf.ObjectStreamMode.generate,
        )


# ---------------------------------------------------------------------------
# Protect (encrypt with user password)
# ---------------------------------------------------------------------------

def protect_pdf(input_path: Path, password: str, output_path: Path) -> None:
    import pikepdf
    with pikepdf.open(str(input_path)) as pdf:
        pdf.save(str(output_path), encryption=pikepdf.Encryption(owner=password, user=password, R=6))


# ---------------------------------------------------------------------------
# Unlock (remove password)
# ---------------------------------------------------------------------------

def unlock_pdf(input_path: Path, password: str, output_path: Path) -> None:
    import pikepdf
    with pikepdf.open(str(input_path), password=password) as pdf:
        pdf.save(str(output_path))


# ---------------------------------------------------------------------------
# Add page numbers
# ---------------------------------------------------------------------------

def add_page_numbers(input_path: Path, output_path: Path, position: str = "bottom-center") -> None:
    from fpdf import FPDF
    from pypdf import PdfReader, PdfWriter
    import io

    reader = PdfReader(str(input_path))
    writer = PdfWriter()

    for i, page in enumerate(reader.pages, 1):
        # Get page dimensions (in points)
        w_pt = float(page.mediabox.width)
        h_pt = float(page.mediabox.height)
        w_mm = w_pt * 25.4 / 72
        h_mm = h_pt * 25.4 / 72

        # Build an overlay PDF with just the page number
        overlay_pdf = FPDF(unit="mm", format=(w_mm, h_mm))
        overlay_pdf.add_page()
        overlay_pdf.set_font("Helvetica", size=9)
        overlay_pdf.set_text_color(100, 100, 100)
        draw_page_number(overlay_pdf, str(i), w_mm, h_mm, position)

        overlay_bytes = io.BytesIO(bytes(overlay_pdf.output()))
        from pypdf import PdfReader as PR2
        overlay_reader = PR2(overlay_bytes)
        overlay_page = overlay_reader.pages[0]

        page.merge_page(overlay_page)
        writer.add_page(page)

    with output_path.open("wb") as f:
        writer.write(f)


def draw_page_number(pdf, text: str, page_width_mm: float, page_height_mm: float, position: str) -> None:
    positions = {
        "top-left": (10, 8, "L"),
        "top-center": (0, 8, "C"),
        "top-right": (-10, 8, "R"),
        "bottom-left": (10, page_height_mm - 10, "L"),
        "bottom-center": (0, page_height_mm - 10, "C"),
        "bottom-right": (-10, page_height_mm - 10, "R"),
    }
    x_offset, y, align = positions.get(position, positions["bottom-center"])
    if align == "L":
        pdf.set_xy(x_offset, y)
        pdf.cell(24, 6, text, align="L")
    elif align == "R":
        pdf.set_xy(page_width_mm - 34, y)
        pdf.cell(24, 6, text, align="R")
    else:
        pdf.set_xy(0, y)
        pdf.cell(page_width_mm, 6, text, align="C")


# ---------------------------------------------------------------------------
# Watermark
# ---------------------------------------------------------------------------

def add_watermark(
    input_path: Path,
    watermark_text: str,
    output_path: Path,
    position: str = "center",
    font_size: int = 40,
) -> None:
    from fpdf import FPDF
    from pypdf import PdfReader, PdfWriter
    import io

    reader = PdfReader(str(input_path))
    writer = PdfWriter()

    for page in reader.pages:
        w_pt = float(page.mediabox.width)
        h_pt = float(page.mediabox.height)
        w_mm = w_pt * 25.4 / 72
        h_mm = h_pt * 25.4 / 72

        wm_pdf = FPDF(unit="mm", format=(w_mm, h_mm))
        wm_pdf.add_page()
        wm_pdf.set_font("Helvetica", size=max(18, min(font_size, 72)))
        wm_pdf.set_text_color(180, 180, 180)
        wm_pdf.set_xy(0, 0)
        wm_pdf.set_auto_page_break(False)
        draw_watermark_text(wm_pdf, watermark_text, w_mm, h_mm, position)

        wm_bytes = io.BytesIO(bytes(wm_pdf.output()))
        from pypdf import PdfReader as PR2
        wm_reader = PR2(wm_bytes)
        wm_page = wm_reader.pages[0]

        page.merge_page(wm_page)
        writer.add_page(page)

    with output_path.open("wb") as f:
        writer.write(f)


def draw_watermark_text(pdf, text: str, page_width_mm: float, page_height_mm: float, position: str) -> None:
    normalized = position.lower()
    if normalized == "top":
        pdf.set_xy(0, 18)
        pdf.cell(page_width_mm, 18, text, align="C")
        return
    if normalized == "bottom":
        pdf.set_xy(0, page_height_mm - 24)
        pdf.cell(page_width_mm, 18, text, align="C")
        return

    with pdf.rotation(45, x=page_width_mm / 2, y=page_height_mm / 2):
        pdf.set_xy(0, page_height_mm / 2 - 10)
        pdf.cell(page_width_mm, 20, text, align="C")


# ---------------------------------------------------------------------------
# PDF → Images  (returns saved paths)
# ---------------------------------------------------------------------------

def pdf_to_images(input_path: Path, output_dir: Path, fmt: str = "jpg") -> list[Path]:
    """Convert each PDF page to an image using pypdf + Pillow (no poppler required).
    Falls back to a best-effort approach via pypdf's page rendering."""
    try:
        # Try pymupdf (fitz) first — best quality, no external deps
        import fitz  # type: ignore
        doc = fitz.open(str(input_path))
        paths: list[Path] = []
        stem = input_path.stem
        for i, page in enumerate(doc, 1):
            mat = fitz.Matrix(2, 2)  # 2x zoom = ~144 DPI
            pix = page.get_pixmap(matrix=mat)
            out = output_dir / f"{stem}_page_{i}.{fmt}"
            pix.save(str(out))
            paths.append(out)
        return paths
    except ImportError:
        pass

    # Fallback: use pypdf to extract embedded images from each page.
    # Note: pages without embedded images cannot be rendered without a PDF rasteriser.
    # Raise a clear error rather than silently producing blank pages.
    from pypdf import PdfReader
    from PIL import Image as PILImage
    reader = PdfReader(str(input_path))
    paths = []
    stem = input_path.stem
    blank_pages: list[int] = []
    for i, page in enumerate(reader.pages, 1):
        images = list(page.images)
        if images:
            img_data = images[0].data
            img = PILImage.open(io.BytesIO(img_data))
        else:
            blank_pages.append(i)
            img = PILImage.new("RGB", (595, 842), color=(255, 255, 255))
        out = output_dir / f"{stem}_page_{i}.{fmt}"
        img.save(str(out))
        paths.append(out)
    if blank_pages:
        total = len(reader.pages)
        if len(blank_pages) == total:
            raise RuntimeError(
                "Could not render this PDF: no embedded images were found and pymupdf is not installed. "
                "Install pymupdf with 'pip install pymupdf' for full PDF rendering support."
            )
    return paths


# ---------------------------------------------------------------------------
# Ollama vision helpers  (shared by pdf_to_text and pdf_to_word)
# ---------------------------------------------------------------------------

def _ollama_check() -> tuple[str, str]:
    """Return (ollama_url, model). Raises RuntimeError if unavailable.

    Deployment note: on hosts without Ollama (e.g. Koyeb free tier), set
    OLLAMA_ENABLED=false to short-circuit the probe and let auto-engine chains
    degrade to tesseract / github / openai immediately.
    """
    import httpx, os
    if os.environ.get("OLLAMA_ENABLED", "true").strip().lower() in ("false", "0", "no"):
        raise RuntimeError("Ollama is disabled (OLLAMA_ENABLED=false).")
    host = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
    model = os.environ.get("OLLAMA_OCR_MODEL", "gemma4:e4b")
    try:
        probe = httpx.get(host + "/api/tags", timeout=5.0)
        probe.raise_for_status()
        available = [m.get("name", "") for m in probe.json().get("models", [])]
        if not any(model in name for name in available):
            raise RuntimeError(
                f"Ollama model '{model}' is not pulled. Run: ollama pull {model}"
            )
    except (httpx.ConnectError, httpx.TimeoutException, httpx.HTTPError) as exc:
        raise RuntimeError(f"Ollama is not reachable at {host}: {exc}")
    return host + "/api/chat", model


def _ollama_post(payload: dict, timeout: float) -> dict:
    """Call Ollama with light retry logic for transient server-side failures."""
    import json
    import time
    import httpx

    url, _ = _ollama_check()
    attempts = 3
    last_error: Exception | None = None

    for attempt in range(1, attempts + 1):
        try:
            resp = httpx.post(url, content=json.dumps(payload), timeout=timeout)
            resp.raise_for_status()
            return resp.json()
        except httpx.HTTPStatusError as exc:
            status = exc.response.status_code
            body = exc.response.text.strip()
            if status < 500 or attempt == attempts:
                detail = body or str(exc)
                raise RuntimeError(f"Ollama request failed ({status}): {detail}") from exc
            last_error = RuntimeError(f"Ollama request failed ({status}): {body or exc}")
        except (httpx.ReadTimeout, httpx.ConnectTimeout, httpx.RemoteProtocolError) as exc:
            if attempt == attempts:
                raise
            last_error = exc

        time.sleep(min(2 * attempt, 5))

    if last_error is not None:
        raise last_error
    raise RuntimeError("Ollama request failed without a response.")


def _is_recoverable_ollama_runtime_error(exc: Exception) -> bool:
    """True when Ollama is up but the selected model cannot serve the request."""
    message = str(exc).lower()
    recoverable_markers = (
        "memory layout cannot be allocated",
        "cuda",
        "out of memory",
    )
    return any(marker in message for marker in recoverable_markers)


def _fallback_text_ai(system_prompt: str, user_content: str) -> str:
    """Best-effort non-Ollama text fallback for local model runtime failures."""
    errors: list[str] = []
    for name, fn in (("github", _github_text), ("openai", _openai_text)):
        try:
            return fn(system_prompt, user_content)
        except Exception as exc:
            errors.append(f"{name}: {exc}")
    raise RuntimeError("No non-Ollama fallback AI engine available.\n" + "\n".join(errors))


def _extract_pdf_text_pages(input_path: Path) -> list[str]:
    """Extract plain text page-by-page via pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(str(input_path))
    return [(page.extract_text() or "").strip() for page in reader.pages]


def _ollama_call_page(image_path: Path, prompt: str) -> str:
    """Send one rendered page image to Ollama. Returns response text."""
    import base64
    _, model = _ollama_check()
    with image_path.open("rb") as f:
        b64 = base64.b64encode(f.read()).decode("ascii")
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt, "images": [b64]}],
        "stream": False,
    }
    resp = _ollama_post(payload, timeout=180.0)
    return resp.get("message", {}).get("content", "").strip()


def _render_pdf_pages(input_path: Path, tmp_dir: Path, dpi: int = 144) -> list[Path]:
    """Render every PDF page to a PNG in tmp_dir using pymupdf."""
    import fitz  # type: ignore
    doc = fitz.open(str(input_path))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    paths: list[Path] = []
    for i, page in enumerate(doc, 1):
        out = tmp_dir / f"page_{i:04d}.png"
        pix = page.get_pixmap(matrix=mat)
        pix.save(str(out))
        paths.append(out)
    return paths


def _render_pdf_pages_with_sizes(
    input_path: Path,
    tmp_dir: Path,
    dpi: int = 144,
    image_format: str = "png",
    jpeg_quality: int = 82,
) -> list[tuple[Path, float, float]]:
    """Render PDF pages and return image path plus original page size in inches."""
    import fitz  # type: ignore
    from PIL import Image as PILImage

    normalized_format = image_format.lower()
    if normalized_format not in {"png", "jpg", "jpeg"}:
        raise ValueError(f"Unsupported render image format: {image_format}")

    extension = "jpg" if normalized_format in {"jpg", "jpeg"} else "png"
    doc = fitz.open(str(input_path))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pages: list[tuple[Path, float, float]] = []
    for i, page in enumerate(doc, 1):
        out = tmp_dir / f"page_{i:04d}.{extension}"
        pix = page.get_pixmap(matrix=mat, alpha=False)
        if extension == "jpg":
            image = PILImage.frombytes("RGB", (pix.width, pix.height), pix.samples)
            try:
                image.save(
                    str(out),
                    format="JPEG",
                    quality=max(35, min(95, int(jpeg_quality))),
                    optimize=True,
                )
            finally:
                image.close()
        else:
            pix.save(str(out))
        pages.append((out, float(page.rect.width) / 72.0, float(page.rect.height) / 72.0))
    return pages


def _fidelity_render_settings(pdf_doc) -> dict:
    """Pick a render profile that balances visual fidelity against output size."""
    page_count = len(pdf_doc)
    if page_count == 0:
        return {"dpi": 144, "image_format": "png", "jpeg_quality": 82}

    sample_indices = sorted({0, page_count // 2, page_count - 1})
    drawing_like_samples = 0
    for index in sample_indices:
        page = pdf_doc[index]
        if _page_is_drawing_like(page.get_text("dict"), float(page.rect.width), float(page.rect.height)):
            drawing_like_samples += 1

    small_format_only = True
    for page in pdf_doc:
        if not _is_small_format_page(float(page.rect.width) / 72.0, float(page.rect.height) / 72.0):
            small_format_only = False
            break

    if small_format_only or drawing_like_samples >= max(1, len(sample_indices) // 2 + 1):
        return {"dpi": 144, "image_format": "png", "jpeg_quality": 82}
    if page_count >= 40:
        return {"dpi": 96, "image_format": "jpg", "jpeg_quality": 68}
    if page_count >= 20:
        return {"dpi": 110, "image_format": "jpg", "jpeg_quality": 74}
    if page_count >= 8:
        return {"dpi": 128, "image_format": "jpg", "jpeg_quality": 80}
    return {"dpi": 144, "image_format": "png", "jpeg_quality": 82}


def _is_small_format_page(page_width_in: float, page_height_in: float) -> bool:
    """Return True for business-card/postcard-like pages rather than document pages."""
    long_edge = max(page_width_in, page_height_in)
    short_edge = min(page_width_in, page_height_in)
    return long_edge <= 6.0 and short_edge <= 4.5


def _small_page_pairing_spec(
    rendered: list[tuple[Path, float, float]],
    gap_in: float = 0.25,
) -> tuple[bool, float, float, float]:
    """Return pairing info for two similarly sized small-format pages."""
    if len(rendered) != 2:
        return False, 0.0, 0.0, 0.0
    first_width_in, first_height_in = rendered[0][1], rendered[0][2]
    second_width_in, second_height_in = rendered[1][1], rendered[1][2]
    if not (_is_small_format_page(first_width_in, first_height_in) and _is_small_format_page(second_width_in, second_height_in)):
        return False, 0.0, 0.0, 0.0
    if abs(first_width_in - second_width_in) > 0.05 or abs(first_height_in - second_height_in) > 0.05:
        return False, 0.0, 0.0, 0.0
    return True, first_width_in * 2.0 + gap_in, first_height_in, gap_in


def _group_small_format_pages(
    rendered: list[tuple[Path, float, float]],
    gap_in: float = 0.25,
) -> list[dict]:
    """Group consecutive small-format pages into front/back pairs when possible."""
    groups: list[dict] = []
    index = 0
    while index < len(rendered):
        image_path, page_width_in, page_height_in = rendered[index]
        if index + 1 < len(rendered):
            next_image_path, next_width_in, next_height_in = rendered[index + 1]
            if (
                _is_small_format_page(page_width_in, page_height_in)
                and _is_small_format_page(next_width_in, next_height_in)
                and abs(page_width_in - next_width_in) <= 0.05
                and abs(page_height_in - next_height_in) <= 0.05
            ):
                groups.append(
                    {
                        "kind": "pair",
                        "indices": [index, index + 1],
                        "width_in": page_width_in * 2.0 + gap_in,
                        "height_in": page_height_in,
                        "gap_in": gap_in,
                    }
                )
                index += 2
                continue

        groups.append(
            {
                "kind": "single",
                "indices": [index],
                "width_in": page_width_in,
                "height_in": page_height_in,
                "gap_in": 0.0,
            }
        )
        index += 1
    return groups


def _combine_rendered_images_side_by_side(
    image_paths: list[Path],
    output_path: Path,
    gap_px: int,
    image_format: str = "png",
    jpeg_quality: int = 82,
) -> None:
    """Combine rendered page images into one side-by-side canvas."""
    from PIL import Image

    images = [Image.open(path).convert("RGB") for path in image_paths]
    try:
        width = sum(image.width for image in images) + gap_px * max(0, len(images) - 1)
        height = max(image.height for image in images)
        combined = Image.new("RGB", (width, height), color=(255, 255, 255))
        x = 0
        for image in images:
            combined.paste(image, (x, 0))
            x += image.width + gap_px
        if image_format.lower() in {"jpg", "jpeg"}:
            combined.save(
                str(output_path),
                format="JPEG",
                quality=max(35, min(95, int(jpeg_quality))),
                optimize=True,
            )
        else:
            combined.save(str(output_path))
    finally:
        for image in images:
            image.close()


def _pdf_pages_to_docx_images(input_path: Path, output_path: Path, dpi: int = 144) -> None:
    """Preserve PDF page appearance by embedding each page as a full-page image in DOCX."""
    import tempfile
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.enum.section import WD_ORIENT, WD_SECTION
    from docx.shared import Inches

    doc = Document()
    first_section = doc.sections[0]

    with tempfile.TemporaryDirectory() as td:
        rendered = _render_pdf_pages_with_sizes(input_path, Path(td), dpi=dpi)
        for index, (image_path, page_width_in, page_height_in) in enumerate(rendered):
            section = first_section if index == 0 else doc.add_section(WD_SECTION.NEW_PAGE)
            section.page_width = Inches(page_width_in)
            section.page_height = Inches(page_height_in)
            section.orientation = WD_ORIENT.LANDSCAPE if page_width_in >= page_height_in else WD_ORIENT.PORTRAIT
            # Fidelity mode should align the placed image to the true PDF page box.
            section.left_margin = Inches(0)
            section.right_margin = Inches(0)
            section.top_margin = Inches(0)
            section.bottom_margin = Inches(0)
            paragraph = doc.add_paragraph()
            paragraph.paragraph_format.space_before = 0
            paragraph.paragraph_format.space_after = 0
            run = paragraph.add_run()
            run.add_picture(
                str(image_path),
                width=Inches(page_width_in),
                height=Inches(page_height_in),
            )

    doc.save(str(output_path))


def _add_docx_textbox(
    paragraph,
    left_pt: float,
    top_pt: float,
    width_pt: float,
    height_pt: float,
    block_lines: list[list[dict]],
    alignment: str,
) -> None:
    """Add a floating textbox anchored to `paragraph` using WordprocessingML/VML."""
    from docx.oxml import OxmlElement, parse_xml

    safe_width = max(8.0, width_pt)
    safe_height = max(8.0, height_pt)
    style = (
        "position:absolute;"
        f"margin-left:{left_pt:.2f}pt;"
        f"margin-top:{top_pt:.2f}pt;"
        f"width:{safe_width:.2f}pt;"
        f"height:{safe_height:.2f}pt;"
        "z-index:251660288;"
        "mso-wrap-style:none;"
        "visibility:visible"
    )
    pict = parse_xml(
        '<w:pict xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
        'xmlns:v="urn:schemas-microsoft-com:vml">'
        f'<v:shape id="cfTextBox" stroked="f" filled="f" style="{style}">'
        '<v:textbox inset="0,0,0,0">'
        "<w:txbxContent/>"
        "</v:textbox>"
        "</v:shape>"
        "</w:pict>"
    )
    txbx_content = next(
        el for el in pict.iter()
        if el.tag == "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}txbxContent"
    )
    p = OxmlElement("w:p")
    p_pr = OxmlElement("w:pPr")
    spacing = OxmlElement("w:spacing")
    spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}before", "0")
    spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}after", "0")
    spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}line", "240")
    p_pr.append(spacing)
    p.append(p_pr)

    jc = OxmlElement("w:jc")
    jc.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", alignment)
    p_pr.append(jc)

    paragraph_elements: list = []
    for line_index, spans in enumerate(block_lines):
        target_p = p if line_index == 0 else OxmlElement("w:p")
        if line_index > 0:
            target_p_pr = OxmlElement("w:pPr")
            target_spacing = OxmlElement("w:spacing")
            target_spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}before", "0")
            target_spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}after", "0")
            target_spacing.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}line", "240")
            target_p_pr.append(target_spacing)
            target_jc = OxmlElement("w:jc")
            target_jc.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", alignment)
            target_p_pr.append(target_jc)
            target_p.append(target_p_pr)

        for span in spans:
            text = str(span.get("text", ""))
            if not text:
                continue
            r = OxmlElement("w:r")
            r_pr = OxmlElement("w:rPr")

            size_half_points = str(int(round(max(6.0, min(72.0, float(span.get("size", 12.0)))) * 2)))
            sz = OxmlElement("w:sz")
            sz.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", size_half_points)
            r_pr.append(sz)

            r_fonts = OxmlElement("w:rFonts")
            font_name = _fit_font_name(span)
            for key in ("ascii", "hAnsi", "cs"):
                r_fonts.set(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{key}", font_name)
            r_pr.append(r_fonts)

            if _span_is_bold(span):
                r_pr.append(OxmlElement("w:b"))
            if _span_is_italic(span):
                r_pr.append(OxmlElement("w:i"))

            color = OxmlElement("w:color")
            r_rgb = _int_to_rgb(int(span.get("color", 0)))
            color.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "%02X%02X%02X" % r_rgb)
            r_pr.append(color)

            r.append(r_pr)
            t = OxmlElement("w:t")
            if text != text.strip():
                t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            t.text = text
            r.append(t)
            target_p.append(r)

        if len(target_p) > 1:
            paragraph_elements.append(target_p)

    for paragraph_element in paragraph_elements:
        txbx_content.append(paragraph_element)

    if paragraph_elements:
        run = paragraph.add_run()
        run._r.append(pict)


def _pdf_pages_to_docx_hybrid(
    input_path: Path,
    output_path: Path,
    dpi: int = 144,
    auto_pair_small_pages: bool = True,
) -> None:
    """Embed each PDF page as an image and overlay editable text boxes in DOCX."""
    import tempfile
    import fitz  # type: ignore
    from docx import Document
    from docx.enum.section import WD_ORIENT, WD_SECTION
    from docx.shared import Inches

    doc = Document()
    first_section = doc.sections[0]
    pdf_doc = fitz.open(str(input_path))
    render_settings = _fidelity_render_settings(pdf_doc)
    dpi = int(render_settings["dpi"])

    with tempfile.TemporaryDirectory() as td:
        rendered = _render_pdf_pages_with_sizes(
            input_path,
            Path(td),
            dpi=dpi,
            image_format=str(render_settings["image_format"]),
            jpeg_quality=int(render_settings["jpeg_quality"]),
        )
        suppressed_overlay_pages = 0
        small_format_only = bool(rendered) and all(
            _is_small_format_page(page_width_in, page_height_in)
            for _image_path, page_width_in, page_height_in in rendered
        )
        groups = _group_small_format_pages(rendered) if auto_pair_small_pages else []
        if any(group["kind"] == "pair" for group in groups):
            for group_index, group in enumerate(groups):
                section = first_section if group_index == 0 else doc.add_section(WD_SECTION.NEW_PAGE)
                section.page_width = Inches(group["width_in"])
                section.page_height = Inches(group["height_in"])
                section.orientation = WD_ORIENT.LANDSCAPE if group["width_in"] >= group["height_in"] else WD_ORIENT.PORTRAIT
                section.left_margin = Inches(0)
                section.right_margin = Inches(0)
                section.top_margin = Inches(0)
                section.bottom_margin = Inches(0)

                image_paragraph = doc.add_paragraph()
                image_paragraph.paragraph_format.space_before = 0
                image_paragraph.paragraph_format.space_after = 0
                image_run = image_paragraph.add_run()

                overlay_anchor = doc.add_paragraph()
                overlay_anchor.paragraph_format.space_before = 0
                overlay_anchor.paragraph_format.space_after = 0

                if group["kind"] == "pair":
                    combined_ext = "jpg" if str(render_settings["image_format"]).lower() in {"jpg", "jpeg"} else "png"
                    combined_image = Path(td) / f"paired_small_pages_{group_index}.{combined_ext}"
                    gap_px = max(12, int(round(group["gap_in"] * dpi)))
                    _combine_rendered_images_side_by_side(
                        [rendered[group["indices"][0]][0], rendered[group["indices"][1]][0]],
                        combined_image,
                        gap_px,
                        image_format=str(render_settings["image_format"]),
                        jpeg_quality=int(render_settings["jpeg_quality"]),
                    )
                    image_run.add_picture(
                        str(combined_image),
                        width=Inches(group["width_in"]),
                        height=Inches(group["height_in"]),
                    )
                    gap_pt = group["gap_in"] * 72.0
                    first_page_width_pt = float(pdf_doc[group["indices"][0]].rect.width)
                    for local_index, page_index in enumerate(group["indices"]):
                        page = pdf_doc[page_index]
                        page_dict = page.get_text("dict")
                        if not _page_supports_live_text_overlays(page, page_dict):
                            suppressed_overlay_pages += 1
                            continue
                        x_offset = local_index * (first_page_width_pt + gap_pt)
                        for lines, (left, top, right, bottom) in _sorted_text_blocks(page_dict, float(page.rect.width)):
                            _add_docx_textbox(
                                overlay_anchor,
                                left + x_offset,
                                top,
                                right - left,
                                bottom - top,
                                lines,
                                _block_alignment(left, right, float(page.rect.width)),
                            )
                else:
                    page_index = group["indices"][0]
                    image_path, page_width_in, page_height_in = rendered[page_index]
                    image_run.add_picture(
                        str(image_path),
                        width=Inches(page_width_in),
                        height=Inches(page_height_in),
                    )
                    page = pdf_doc[page_index]
                    page_dict = page.get_text("dict")
                    if not _page_supports_live_text_overlays(page, page_dict):
                        suppressed_overlay_pages += 1
                        continue
                    for lines, (left, top, right, bottom) in _sorted_text_blocks(page_dict, float(page.rect.width)):
                        _add_docx_textbox(
                            overlay_anchor,
                            left,
                            top,
                            right - left,
                            bottom - top,
                            lines,
                            _block_alignment(left, right, float(page.rect.width)),
                        )

            if suppressed_overlay_pages and not small_format_only:
                _append_word_editable_text_appendix(doc, pdf_doc)
            doc.save(str(output_path))
            return

        for index, ((image_path, page_width_in, page_height_in), page) in enumerate(zip(rendered, pdf_doc)):
            section = first_section if index == 0 else doc.add_section(WD_SECTION.NEW_PAGE)
            section.page_width = Inches(page_width_in)
            section.page_height = Inches(page_height_in)
            section.orientation = WD_ORIENT.LANDSCAPE if page_width_in >= page_height_in else WD_ORIENT.PORTRAIT
            # Keep the PDF image and overlay coordinate systems identical.
            section.left_margin = Inches(0)
            section.right_margin = Inches(0)
            section.top_margin = Inches(0)
            section.bottom_margin = Inches(0)
            image_paragraph = doc.add_paragraph()
            image_paragraph.paragraph_format.space_before = 0
            image_paragraph.paragraph_format.space_after = 0
            image_run = image_paragraph.add_run()
            image_run.add_picture(
                str(image_path),
                width=Inches(page_width_in),
                height=Inches(page_height_in),
            )
            if small_format_only:
                continue

            overlay_anchor = doc.add_paragraph()
            overlay_anchor.paragraph_format.space_before = 0
            overlay_anchor.paragraph_format.space_after = 0
            page_dict = page.get_text("dict")
            if not _page_supports_live_text_overlays(page, page_dict):
                suppressed_overlay_pages += 1
                continue
            for lines, (left, top, right, bottom) in _sorted_text_blocks(page_dict, float(page.rect.width)):
                _add_docx_textbox(
                    overlay_anchor,
                    left,
                    top,
                    right - left,
                    bottom - top,
                    lines,
                    _block_alignment(left, right, float(page.rect.width)),
                )

        if (len(pdf_doc) >= 8 or suppressed_overlay_pages > 0) and not small_format_only:
            _append_word_editable_text_appendix(doc, pdf_doc)

    doc.save(str(output_path))


def _pdf_pages_to_pptx_images(input_path: Path, output_path: Path, dpi: int = 144) -> None:
    """Preserve PDF page appearance by turning each page into a full-bleed slide image."""
    import tempfile
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as td:
        rendered = _render_pdf_pages_with_sizes(input_path, Path(td), dpi=dpi)
        if not rendered:
            raise ValueError("The PDF has no pages to convert.")

        first_width_in, first_height_in = rendered[0][1], rendered[0][2]
        prs = Presentation()
        prs.slide_width = Inches(first_width_in)
        prs.slide_height = Inches(first_height_in)
        blank_layout = prs.slide_layouts[6]

        for image_path, _, _ in rendered:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        prs.save(str(output_path))


def _int_to_rgb(color_int: int) -> tuple[int, int, int]:
    """Convert PyMuPDF packed integer color into an RGB tuple."""
    return ((color_int >> 16) & 255, (color_int >> 8) & 255, color_int & 255)


def _span_is_bold(span: dict) -> bool:
    name = str(span.get("font", "")).lower()
    return "bold" in name or "black" in name or "heavy" in name


def _span_is_italic(span: dict) -> bool:
    name = str(span.get("font", "")).lower()
    return "italic" in name or "oblique" in name


def _fit_font_name(span: dict) -> str:
    """Map PDF font names to a small safe PowerPoint font set."""
    name = str(span.get("font", "")).lower()
    if "courier" in name or "mono" in name:
        return "Courier New"
    if "times" in name or "serif" in name:
        return "Times New Roman"
    if "calibri" in name:
        return "Calibri"
    if "arial" in name:
        return "Arial"
    return "Helvetica"


def _block_lines(block: dict) -> list[list[dict]]:
    """Return non-empty span lines for a text block."""
    output: list[list[dict]] = []
    for line in block.get("lines") or []:
        spans = [span for span in line.get("spans", []) if str(span.get("text", "")).strip()]
        if spans:
            output.append(spans)
    return output


def _block_bbox(lines: list[list[dict]]) -> tuple[float, float, float, float]:
    """Return left, top, right, bottom for a grouped text block."""
    left = min(float(span["bbox"][0]) for line in lines for span in line)
    top = min(float(span["bbox"][1]) for line in lines for span in line)
    right = max(float(span["bbox"][2]) for line in lines for span in line)
    bottom = max(float(span["bbox"][3]) for line in lines for span in line)
    return left, top, right, bottom


def _block_alignment(left: float, right: float, page_width_pt: float) -> str:
    """Infer rough alignment from block geometry."""
    center = (left + right) / 2.0
    page_center = page_width_pt / 2.0
    width = right - left
    if abs(center - page_center) < max(18.0, page_width_pt * 0.05):
        return "center"
    if right > page_width_pt * 0.7 and width < page_width_pt * 0.45:
        return "right"
    return "left"


def _sorted_text_blocks(page_dict: dict, page_width_pt: float) -> list[tuple[list[list[dict]], tuple[float, float, float, float]]]:
    """Return text blocks in an approximate human reading order across columns."""
    entries: list[tuple[list[list[dict]], tuple[float, float, float, float]]] = []
    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        lines = _block_lines(block)
        if not lines:
            continue
        entries.append((lines, _block_bbox(lines)))

    if not entries:
        return []

    # Full-width blocks like titles should stay near the top regardless of columns.
    full_width_cutoff = page_width_pt * 0.6
    full_width_entries: list[tuple[list[list[dict]], tuple[float, float, float, float]]] = []
    column_entries: list[tuple[list[list[dict]], tuple[float, float, float, float]]] = []
    for item in entries:
        _lines, (left, _top, right, _bottom) = item
        if (right - left) >= full_width_cutoff:
            full_width_entries.append(item)
        else:
            column_entries.append(item)

    full_width_entries.sort(key=lambda item: (item[1][1], item[1][0]))
    if not column_entries:
        return full_width_entries

    # Cluster blocks into reading columns by left edge. This handles common
    # left-to-right multi-column layouts better than a pure top-first sort.
    column_entries.sort(key=lambda item: item[1][0])
    columns: list[list[tuple[list[list[dict]], tuple[float, float, float, float]]]] = []
    gap_threshold = max(36.0, page_width_pt * 0.08)

    for item in column_entries:
        left = item[1][0]
        if not columns:
            columns.append([item])
            continue
        current_column = columns[-1]
        column_lefts = [entry[1][0] for entry in current_column]
        avg_left = sum(column_lefts) / len(column_lefts)
        if abs(left - avg_left) <= gap_threshold:
            current_column.append(item)
        else:
            columns.append([item])

    for column in columns:
        column.sort(key=lambda item: (item[1][1], item[1][0]))

    ordered_columns = sorted(
        columns,
        key=lambda col: (sum(entry[1][0] for entry in col) / len(col), min(entry[1][1] for entry in col)),
    )

    ordered: list[tuple[list[list[dict]], tuple[float, float, float, float]]] = []
    ordered.extend(full_width_entries)
    for column in ordered_columns:
        ordered.extend(column)
    return ordered


def _block_plain_text(lines: list[list[dict]]) -> list[str]:
    """Return plain text lines for a grouped text block."""
    output: list[str] = []
    for spans in lines:
        text = "".join(str(span.get("text", "")) for span in spans).strip()
        if text:
            output.append(text)
    return output


def _page_plain_text_blocks(page_dict: dict, page_width_pt: float) -> list[list[str]]:
    """Return grouped plain-text blocks in reading order for a page."""
    output: list[list[str]] = []
    for lines, _bbox in _sorted_text_blocks(page_dict, page_width_pt):
        block_lines = _block_plain_text(lines)
        if block_lines:
            output.append(block_lines)
    return output


def _cluster_positions(values: list[float], threshold: float) -> list[float]:
    """Cluster sorted positions into anchor points."""
    if not values:
        return []
    anchors: list[list[float]] = [[values[0]]]
    for value in values[1:]:
        if abs(value - anchors[-1][-1]) <= threshold:
            anchors[-1].append(value)
        else:
            anchors.append([value])
    return [sum(group) / len(group) for group in anchors]


def _extract_grid_like_table(page_dict: dict, page_width_pt: float) -> list[list[str]]:
    """Infer a simple editable table from text alignment when line borders are absent."""
    line_entries: list[dict] = []
    heights: list[float] = []

    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [span for span in line.get("spans", []) if str(span.get("text", "")).strip()]
            if not spans:
                continue
            text = "".join(str(span.get("text", "")) for span in spans).strip()
            if not text:
                continue
            bbox = line.get("bbox")
            if not bbox:
                continue
            left, top, right, bottom = (float(v) for v in bbox)
            heights.append(bottom - top)
            line_entries.append(
                {
                    "text": text,
                    "bbox": (left, top, right, bottom),
                    "y_center": (top + bottom) / 2.0,
                    "x_left": left,
                }
            )

    if len(line_entries) < 6:
        return []

    typical_height = sorted(heights)[len(heights) // 2] if heights else 12.0
    row_threshold = max(8.0, typical_height * 0.9)
    line_entries.sort(key=lambda entry: (entry["y_center"], entry["x_left"]))

    raw_rows: list[list[dict]] = []
    for entry in line_entries:
        if not raw_rows or abs(entry["y_center"] - raw_rows[-1][0]["y_center"]) > row_threshold:
            raw_rows.append([entry])
        else:
            raw_rows[-1].append(entry)

    candidate_rows = []
    for row in raw_rows:
        row.sort(key=lambda entry: entry["x_left"])
        if len(row) >= 2:
            candidate_rows.append(row)

    if len(candidate_rows) < 3:
        return []

    column_positions = sorted(entry["x_left"] for row in candidate_rows for entry in row)
    column_threshold = max(20.0, page_width_pt * 0.03)
    anchors = _cluster_positions(column_positions, column_threshold)
    if len(anchors) < 2:
        return []

    grid_rows: list[list[str]] = []
    populated_rows = 0
    for row in candidate_rows:
        cells = [""] * len(anchors)
        for entry in row:
            nearest_index = min(range(len(anchors)), key=lambda idx: abs(entry["x_left"] - anchors[idx]))
            if cells[nearest_index]:
                cells[nearest_index] = f"{cells[nearest_index]} {entry['text']}".strip()
            else:
                cells[nearest_index] = entry["text"]
        filled = sum(1 for value in cells if value.strip())
        if filled >= 2:
            populated_rows += 1
            grid_rows.append(cells)

    if populated_rows < 3:
        return []

    non_empty_counts = [sum(1 for row in grid_rows if idx < len(row) and str(row[idx]).strip()) for idx in range(len(anchors))]
    keep_indices = [idx for idx, count in enumerate(non_empty_counts) if count >= 2]
    if len(keep_indices) < 2:
        return []

    compact_rows = [[row[idx].strip() for idx in keep_indices] for row in grid_rows]
    if len(compact_rows) < 3:
        return []
    return compact_rows


def _merge_wrapped_table_rows(rows: list[list[str]]) -> list[list[str]]:
    """Merge continuation rows back into the previous logical row."""
    if len(rows) <= 2:
        return rows

    header = [str(cell).strip() for cell in rows[0]]
    merged: list[list[str]] = [header]

    for row in rows[1:]:
        clean = [str(cell).strip() for cell in row]
        filled = sum(1 for value in clean if value)
        if (
            merged
            and not clean[0]
            and filled > 0
            and any(str(value).strip() for value in merged[-1])
        ):
            previous = merged[-1]
            for index, value in enumerate(clean):
                if not value:
                    continue
                if index >= len(previous):
                    previous.extend([""] * (index + 1 - len(previous)))
                previous[index] = f"{previous[index]} {value}".strip() if previous[index] else value
        else:
            merged.append(clean)

    return merged


def _is_useful_structured_table(rows: list[list[str]]) -> bool:
    """Reject malformed table detections that collapse into one giant text cell."""
    if len(rows) < 2:
        return False
    width = max(len(row) for row in rows)
    if width < 2:
        return False
    non_empty_per_row = [sum(1 for cell in row if str(cell).strip()) for row in rows]
    if max(non_empty_per_row, default=0) < 2:
        return False
    populated_columns = sum(
        1
        for col_index in range(width)
        if sum(1 for row in rows if col_index < len(row) and str(row[col_index]).strip()) >= 2
    )
    if populated_columns < 2:
        return False
    giant_single_cell_rows = sum(
        1
        for row, count in zip(rows, non_empty_per_row)
        if count == 1 and max((len(str(cell).strip()) for cell in row if str(cell).strip()), default=0) >= 160
    )
    return giant_single_cell_rows < len(rows) / 2


def _normalize_table_cell_text(value: str) -> str:
    """Collapse PDF line breaks into cleaner single-line Excel cell values."""
    text = str(value or "").replace("\r", "\n")
    text = re.sub(r"\s*\n\s*", " ", text)
    text = re.sub(r"\s{2,}", " ", text)
    text = re.sub(r"\s+([,.;:])", r"\1", text)
    return text.strip()


def _extract_blob_table_amount_hints(rows: list[list[str]]) -> dict[str, str]:
    """Recover item->amount hints from malformed one-cell table blobs."""
    hints: dict[str, str] = {}
    for row in rows:
        blob = " ".join(str(cell or "").strip() for cell in row if cell is not None).strip()
        if not blob or "Contract Amount" not in blob:
            continue
        blob = blob.replace("\r", "\n")
        item_matches = list(re.finditer(r"(?m)^\s*(\d{1,3})\s+", blob))
        if not item_matches:
            continue
        for index, match in enumerate(item_matches):
            item = match.group(1)
            start = match.start()
            end = item_matches[index + 1].start() if index + 1 < len(item_matches) else len(blob)
            segment = blob[start:end]
            date_match = re.search(r"\b\d{2}-[A-Za-z]{3}-\d{2}\b", segment)
            if not date_match:
                continue
            prefix = segment[:date_match.start()]
            amount_match = re.search(r"\b\d{1,3}(?:,\d{3})+\.\d{2}\b", prefix)
            if amount_match:
                hints[item] = amount_match.group(0)
    return hints


def _coerce_excel_cell_value(value: str):
    """Convert common table strings into typed Excel values when it is reasonably safe."""
    from datetime import datetime

    text = _normalize_table_cell_text(value)
    if not text:
        return ""

    compact = text.replace(",", "")

    date_formats = (
        "%m/%d/%Y",
        "%m/%d/%y",
        "%d/%m/%Y",
        "%d/%m/%y",
        "%Y-%m-%d",
        "%d-%b-%y",
        "%d-%b-%Y",
        "%d-%B-%Y",
        "%b %d, %Y",
        "%B %d, %Y",
    )
    for fmt in date_formats:
        try:
            return datetime.strptime(text, fmt)
        except ValueError:
            continue

    # Keep IDs and mixed alphanumeric codes as text after date parsing.
    if re.search(r"[A-Za-z]", compact) and re.search(r"\d", compact):
        return text

    if re.fullmatch(r"\(?-?\$?\d[\d,]*\.\d+\)?", text) or re.fullmatch(r"\(?-?\$?\d[\d,]*\)?", text):
        negative = text.startswith("(") and text.endswith(")")
        numeric = text.strip("()").replace("$", "").replace(",", "")
        try:
            value = float(numeric) if "." in numeric else int(numeric)
            return -value if negative else value
        except ValueError:
            return text

    if re.fullmatch(r"-?\d+(\.\d+)?%", text):
        try:
            return float(text[:-1]) / 100.0
        except ValueError:
            return text

    return text


def _header_semantic(header_value: str) -> str:
    """Infer a basic semantic type from a table header."""
    header = str(header_value or "").strip().lower()
    if not header:
        return "text"
    if any(token in header for token in ("date", "completed", "start", "end", "deadline", "due")):
        return "date"
    if any(token in header for token in ("amount", "price", "cost", "total", "budget", "value")):
        return "amount"
    if any(token in header for token in ("percent", "margin", "rate", "%")):
        return "percent"
    if any(token in header for token in ("qty", "quantity", "count", "no.", "number", "item")):
        return "count"
    return "text"


def _apply_excel_table_header_style(cell):
    """Make extracted table headers read like native spreadsheet headers."""
    from openpyxl.styles import Alignment, Font, PatternFill

    cell.font = Font(bold=True)
    cell.alignment = Alignment(vertical="top", horizontal="center", wrap_text=True)
    cell.fill = PatternFill(fill_type="solid", fgColor="E8EEF7")


def _mark_excel_recovered_cell(cell, note: str):
    """Visually flag a recovered fallback value in the workbook."""
    from openpyxl.comments import Comment
    from openpyxl.styles import PatternFill

    cell.fill = PatternFill(fill_type="solid", fgColor="FFF4CC")
    cell.comment = Comment(note, "ConvertFlow")


def _excel_column_name(index: int) -> str:
    """Convert a 1-based column index to an Excel column name."""
    result = ""
    current = max(1, int(index))
    while current:
        current, remainder = divmod(current - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _add_excel_table(ws, start_row: int, column_count: int, row_count: int, name_prefix: str, table_index: int):
    """Add an Excel table object over an extracted structured range."""
    from openpyxl.worksheet.table import Table, TableStyleInfo

    if column_count < 1 or row_count < 2:
        return
    seen_headers: set[str] = set()
    for column_index in range(1, column_count + 1):
        cell = ws.cell(start_row, column_index)
        header = str(cell.value).strip() if cell.value is not None else ""
        if not header:
            header = f"Column {column_index}"
        base_header = header
        suffix = 2
        while header in seen_headers:
            header = f"{base_header} {suffix}"
            suffix += 1
        cell.value = header
        seen_headers.add(header)
    end_col = _excel_column_name(column_count)
    end_row = start_row + row_count - 1
    table = Table(displayName=f"{name_prefix}_{table_index}", ref=f"A{start_row}:{end_col}{end_row}")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    ws.add_table(table)


def _write_excel_table_cell(cell, value, header_value: str = ""):
    """Write a structured table cell with basic value typing and formatting."""
    from datetime import datetime
    from openpyxl.styles import Alignment

    coerced = _coerce_excel_cell_value(value)
    cell.value = coerced
    semantic = _header_semantic(header_value)
    horizontal = "left"
    if isinstance(coerced, datetime):
        horizontal = "center"
    elif isinstance(coerced, (int, float)):
        horizontal = "right"
    elif semantic == "date":
        horizontal = "center"
    cell.alignment = Alignment(vertical="top", horizontal=horizontal, wrap_text=True)
    if isinstance(coerced, datetime):
        cell.number_format = "yyyy-mm-dd"
    elif isinstance(coerced, float):
        if str(value).strip().endswith("%"):
            cell.number_format = "0.00%"
        elif "." in str(value):
            cell.number_format = "#,##0.00"
    elif isinstance(coerced, int):
        cell.number_format = "#,##0"


def _append_word_editable_text_appendix(doc, pdf_doc) -> None:
    """Append native editable text to the end of a DOCX for easier copy/edit workflows."""
    added_any = False
    total_pages = len(pdf_doc)
    for page_index, page in enumerate(pdf_doc, 1):
        page_dict = page.get_text("dict")
        if not _page_supports_notes_extract(page, page_dict, total_pages):
            continue
        blocks = _page_plain_text_blocks(page_dict, float(page.rect.width))
        if not blocks:
            continue
        if not added_any:
            doc.add_page_break()
            doc.add_heading("Editable Text Extract", level=1)
            doc.add_paragraph(
                "Use this section to copy or revise text without disturbing the visual page reconstruction above."
            )
            added_any = True
        doc.add_heading(f"Page {page_index}", level=2)
        for block_lines in blocks:
            doc.add_paragraph(" ".join(block_lines))


def _set_slide_notes_from_blocks(slide, sections: list[tuple[str, list[list[str]]]]) -> None:
    """Populate PowerPoint speaker notes with editable text extracted from the PDF page(s)."""
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    first_paragraph = True
    for heading, blocks in sections:
        if not blocks:
            continue
        heading_paragraph = notes_frame.paragraphs[0] if first_paragraph else notes_frame.add_paragraph()
        heading_paragraph.text = heading
        first_paragraph = False
        for block_lines in blocks:
            paragraph = notes_frame.add_paragraph()
            paragraph.text = " ".join(block_lines)


def _page_is_drawing_like(page_dict: dict, page_width_pt: float, page_height_pt: float) -> bool:
    """Heuristic: CAD/drawing sheets often contain many tiny or vertical annotations."""
    total_lines = 0
    vertical_lines = 0
    tiny_lines = 0

    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [span for span in line.get("spans", []) if str(span.get("text", "")).strip()]
            if not spans:
                continue
            total_lines += 1
            dir_x, dir_y = line.get("dir", (1.0, 0.0))
            if abs(dir_y) > abs(dir_x):
                vertical_lines += 1
            bbox = line.get("bbox")
            if bbox:
                width = float(bbox[2]) - float(bbox[0])
                height = float(bbox[3]) - float(bbox[1])
                if width < max(28.0, page_width_pt * 0.025) or height < max(8.0, page_height_pt * 0.008):
                    tiny_lines += 1

    if total_lines == 0:
        return False

    vertical_ratio = vertical_lines / total_lines
    tiny_ratio = tiny_lines / total_lines
    large_sheet = page_width_pt >= 1000 or page_height_pt >= 1000
    return large_sheet and (vertical_ratio >= 0.35 or tiny_ratio >= 0.55 or (vertical_lines + tiny_lines) >= 30)


def _page_supports_live_text_overlays(page, page_dict: dict) -> bool:
    """Return False when live overlays are likely to hurt fidelity more than help."""
    page_width_pt = float(page.rect.width)
    page_height_pt = float(page.rect.height)
    if _page_is_drawing_like(page_dict, page_width_pt, page_height_pt):
        return False

    text_blocks = _sorted_text_blocks(page_dict, page_width_pt)
    if len(text_blocks) >= 18:
        return False

    image_count = len(page.get_images(full=True))
    if image_count >= 6 and len(text_blocks) >= 10:
        return False

    try:
        tables = page.find_tables()
        if len(tables.tables) > 0:
            return False
    except Exception:
        pass

    return True


def _page_supports_notes_extract(page, page_dict: dict, total_pages: int) -> bool:
    """Return True when notes/appendix text is likely to help more than clutter."""
    if total_pages >= 8:
        return True
    if _page_is_drawing_like(page_dict, float(page.rect.width), float(page.rect.height)):
        return False
    width_in = float(page.rect.width) / 72.0
    height_in = float(page.rect.height) / 72.0
    if _is_small_format_page(width_in, height_in):
        return False
    try:
        tables = page.find_tables()
        if len(tables.tables) > 0:
            return True
    except Exception:
        pass
    return not _page_supports_live_text_overlays(page, page_dict)


def _bbox_overlaps(a: tuple[float, float, float, float], b: tuple[float, float, float, float], padding: float = 4.0) -> bool:
    """Return True when two rectangles overlap, allowing a small padding."""
    a_left, a_top, a_right, a_bottom = a
    b_left, b_top, b_right, b_bottom = b
    return not (
        a_right < b_left - padding
        or b_right < a_left - padding
        or a_bottom < b_top - padding
        or b_bottom < a_top - padding
    )


def _add_pdf_text_overlays_to_slide(slide, page_dict: dict, page_width_pt: float, page_height_pt: float) -> None:
    """Overlay editable text boxes onto a slide using grouped PyMuPDF text blocks."""
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    for lines, (left, top, right, bottom) in _sorted_text_blocks(page_dict, page_width_pt):
        width = max(8.0, right - left)
        height = max(8.0, bottom - top)

        textbox = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        alignment = _block_alignment(left, right, page_width_pt)
        ppt_align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[alignment]

        for line_index, spans in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.alignment = ppt_align
            for span in spans:
                run = paragraph.add_run()
                run.text = str(span.get("text", ""))
                font = run.font
                font.name = _fit_font_name(span)
                font.size = Pt(max(6.0, min(72.0, float(span.get("size", 12.0)))))
                font.bold = _span_is_bold(span)
                font.italic = _span_is_italic(span)
                font.color.rgb = RGBColor(*_int_to_rgb(int(span.get("color", 0))))

        if str(text_frame.text).strip():
            fill = textbox.fill
            fill.background()
            textbox.line.fill.background()


def _pdf_pages_to_pptx_hybrid(
    input_path: Path,
    output_path: Path,
    dpi: int = 144,
    auto_pair_small_pages: bool = True,
) -> None:
    """Create slides with a full-page background image plus editable text overlays."""
    import tempfile
    import fitz  # type: ignore
    from pptx import Presentation
    from pptx.util import Inches

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to convert.")
    render_settings = _fidelity_render_settings(doc)
    dpi = int(render_settings["dpi"])

    with tempfile.TemporaryDirectory() as td:
        rendered = _render_pdf_pages_with_sizes(
            input_path,
            Path(td),
            dpi=dpi,
            image_format=str(render_settings["image_format"]),
            jpeg_quality=int(render_settings["jpeg_quality"]),
        )
        groups = _group_small_format_pages(rendered) if auto_pair_small_pages else []
        if any(group["kind"] == "pair" for group in groups):
            deck_width_in = max(float(group["width_in"]) for group in groups)
            deck_height_in = max(float(group["height_in"]) for group in groups)
            prs = Presentation()
            prs.slide_width = Inches(deck_width_in)
            prs.slide_height = Inches(deck_height_in)
            blank_layout = prs.slide_layouts[6]
            deck_width_pt = deck_width_in * 72.0
            deck_height_pt = deck_height_in * 72.0

            for group_index, group in enumerate(groups):
                slide = prs.slides.add_slide(blank_layout)
                if group["kind"] == "pair":
                    combined_ext = "jpg" if str(render_settings["image_format"]).lower() in {"jpg", "jpeg"} else "png"
                    combined_image = Path(td) / f"paired_small_pages_{group_index}.{combined_ext}"
                    gap_px = max(12, int(round(group["gap_in"] * dpi)))
                    _combine_rendered_images_side_by_side(
                        [rendered[group["indices"][0]][0], rendered[group["indices"][1]][0]],
                        combined_image,
                        gap_px,
                        image_format=str(render_settings["image_format"]),
                        jpeg_quality=int(render_settings["jpeg_quality"]),
                    )
                    slide.shapes.add_picture(
                        str(combined_image),
                        0,
                        0,
                        width=Inches(group["width_in"]),
                        height=Inches(group["height_in"]),
                    )
                    gap_pt = group["gap_in"] * 72.0
                    first_page_width_pt = float(doc[group["indices"][0]].rect.width)
                    group_width_pt = first_page_width_pt * 2.0 + gap_pt
                    for local_index, page_index in enumerate(group["indices"]):
                        page = doc[page_index]
                        page_dict = page.get_text("dict")
                        if not _page_supports_live_text_overlays(page, page_dict):
                            continue
                        x_offset = local_index * (first_page_width_pt + gap_pt)
                        translated_dict = {"blocks": []}
                        for block in page_dict.get("blocks", []):
                            if block.get("type") != 0:
                                continue
                            new_block = {"type": 0, "lines": []}
                            for line in block.get("lines", []):
                                new_line = dict(line)
                                new_spans = []
                                for span in line.get("spans", []):
                                    new_span = dict(span)
                                    bbox = list(new_span.get("bbox", (0, 0, 0, 0)))
                                    bbox[0] += x_offset
                                    bbox[2] += x_offset
                                    new_span["bbox"] = tuple(bbox)
                                    new_spans.append(new_span)
                                new_line["spans"] = new_spans
                                new_block["lines"].append(new_line)
                            translated_dict["blocks"].append(new_block)
                        _add_pdf_text_overlays_to_slide(
                            slide,
                            translated_dict,
                            group_width_pt,
                            float(page.rect.height),
                        )
                    note_sections: list[tuple[str, list[list[str]]]] = []
                    for page_index in group["indices"]:
                        page = doc[page_index]
                        page_dict = page.get_text("dict")
                        blocks = _page_plain_text_blocks(page_dict, float(page.rect.width))
                        if blocks and _page_supports_notes_extract(page, page_dict, len(doc)):
                            note_sections.append((f"Page {page_index + 1}", blocks))
                    if note_sections:
                        _set_slide_notes_from_blocks(slide, note_sections)
                else:
                    page_index = group["indices"][0]
                    image_path, page_width_in, page_height_in = rendered[page_index]
                    x_offset_in = max(0.0, (deck_width_in - page_width_in) / 2.0)
                    y_offset_in = max(0.0, (deck_height_in - page_height_in) / 2.0)
                    slide.shapes.add_picture(
                        str(image_path),
                        Inches(x_offset_in),
                        Inches(y_offset_in),
                        width=Inches(page_width_in),
                        height=Inches(page_height_in),
                    )
                    page = doc[page_index]
                    page_dict = page.get_text("dict")
                    supports_live_overlay = _page_supports_live_text_overlays(page, page_dict)
                    if supports_live_overlay:
                        translated_dict = {"blocks": []}
                        x_offset_pt = x_offset_in * 72.0
                        y_offset_pt = y_offset_in * 72.0
                        for block in page_dict.get("blocks", []):
                            if block.get("type") != 0:
                                continue
                            new_block = {"type": 0, "lines": []}
                            for line in block.get("lines", []):
                                new_line = dict(line)
                                new_spans = []
                                for span in line.get("spans", []):
                                    new_span = dict(span)
                                    bbox = list(new_span.get("bbox", (0, 0, 0, 0)))
                                    bbox[0] += x_offset_pt
                                    bbox[1] += y_offset_pt
                                    bbox[2] += x_offset_pt
                                    bbox[3] += y_offset_pt
                                    new_span["bbox"] = tuple(bbox)
                                    new_spans.append(new_span)
                                new_line["spans"] = new_spans
                                new_block["lines"].append(new_line)
                            translated_dict["blocks"].append(new_block)
                        _add_pdf_text_overlays_to_slide(
                            slide,
                            translated_dict,
                            deck_width_pt,
                            deck_height_pt,
                        )
                    blocks = _page_plain_text_blocks(page_dict, float(page.rect.width))
                    if blocks and _page_supports_notes_extract(page, page_dict, len(doc)):
                        _set_slide_notes_from_blocks(slide, [(f"Page {page_index + 1}", blocks)])
            prs.save(str(output_path))
            return

        first_width_in, first_height_in = rendered[0][1], rendered[0][2]
        prs = Presentation()
        prs.slide_width = Inches(first_width_in)
        prs.slide_height = Inches(first_height_in)
        blank_layout = prs.slide_layouts[6]

        for page_index, ((image_path, _, _), page) in enumerate(zip(rendered, doc), 1):
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            page_dict = page.get_text("dict")
            supports_live_overlay = _page_supports_live_text_overlays(page, page_dict)
            if supports_live_overlay:
                _add_pdf_text_overlays_to_slide(
                    slide,
                    page_dict,
                    float(page.rect.width),
                    float(page.rect.height),
                )
            blocks = _page_plain_text_blocks(page_dict, float(page.rect.width))
            if blocks and _page_supports_notes_extract(page, page_dict, len(doc)):
                _set_slide_notes_from_blocks(slide, [(f"Page {page_index}", blocks)])

        prs.save(str(output_path))


def _pdf_pages_to_excel_hybrid(input_path: Path, output_path: Path, dpi: int = 144) -> None:
    """Preserve the page snapshot and add extracted tables/text into editable cells."""
    import tempfile
    import fitz  # type: ignore
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Alignment, Font

    wb = Workbook()
    wb.remove(wb.active)
    pdf_doc = fitz.open(str(input_path))
    render_settings = _fidelity_render_settings(pdf_doc)
    dpi = int(render_settings["dpi"])

    with tempfile.TemporaryDirectory() as td:
        rendered = _render_pdf_pages_with_sizes(
            input_path,
            Path(td),
            dpi=dpi,
            image_format=str(render_settings["image_format"]),
            jpeg_quality=int(render_settings["jpeg_quality"]),
        )
        if not rendered:
            wb.create_sheet("Page 1")
        for idx, ((image_path, _, _), page) in enumerate(zip(rendered, pdf_doc), 1):
            ws = wb.create_sheet(title=f"Page {idx}")
            first_table_header_row = None
            ws["A1"] = f"PDF page {idx} snapshot"
            ws["A1"].font = Font(bold=True)
            ws["A2"] = "Hybrid fidelity mode preserves the page image and adds extracted content below."
            ws["A3"] = "Use the rows below for editable data; use the snapshot for original visual context."
            image = XLImage(str(image_path))
            image.width = max(640, int(image.width * 0.9))
            image.height = max(900, int(image.height * 0.9))
            ws.add_image(image, "A5")
            ws.column_dimensions["A"].width = 18
            ws.column_dimensions["B"].width = 18
            ws.column_dimensions["C"].width = 18
            ws.column_dimensions["D"].width = 18
            ws.row_dimensions[5].height = 24

            current_row = max(55, int(image.height / 18) + 8)
            ws.cell(current_row, 1, "Extracted content")
            ws.cell(current_row, 1).font = Font(bold=True, size=12)
            current_row += 2

            table_bboxes: list[tuple[float, float, float, float]] = []
            wrote_structured = False
            page_dict = page.get_text("dict")
            blob_amount_hints: dict[str, str] = {}
            try:
                tables = page.find_tables()
                for table_index, table in enumerate(tables.tables, 1):
                    rows = table.extract()
                    if not rows:
                        continue
                    blob_amount_hints.update(_extract_blob_table_amount_hints(rows))
                    rows = _merge_wrapped_table_rows(
                        [[str(cell).strip() if cell is not None else "" for cell in row] for row in rows]
                    )
                    rows = [[_normalize_table_cell_text(cell) for cell in row] for row in rows]
                    if not _is_useful_structured_table(rows):
                        continue
                    table_bboxes.append(tuple(float(v) for v in table.bbox))
                    ws.cell(current_row, 1, f"Table {table_index}")
                    ws.cell(current_row, 1).font = Font(bold=True)
                    current_row += 1
                    table_start_row = current_row
                    header_row = [str(cell).strip() if cell is not None else "" for cell in rows[0]] if rows else []
                    for row_index, row in enumerate(rows):
                        for col_index, value in enumerate(row, 1):
                            cell = ws.cell(current_row, col_index)
                            recovered_from_blob = False
                            if row_index > 0 and not str(value or "").strip():
                                header_value = header_row[col_index - 1] if col_index - 1 < len(header_row) else ""
                                item_value = str(row[0]).strip() if row and len(row) >= 1 else ""
                                if _header_semantic(header_value) == "amount" and item_value in blob_amount_hints:
                                    value = blob_amount_hints[item_value]
                                    recovered_from_blob = True
                            header_value = header_row[col_index - 1] if row_index > 0 and col_index - 1 < len(header_row) else ""
                            _write_excel_table_cell(cell, value, header_value=header_value)
                            if recovered_from_blob:
                                _mark_excel_recovered_cell(
                                    cell,
                                    f"Recovered from malformed blob table fallback for item {item_value}.",
                                )
                            if row_index == 0:
                                _apply_excel_table_header_style(cell)
                        current_row += 1
                    _add_excel_table(
                        ws,
                        table_start_row,
                        max(len(row) for row in rows) if rows else 0,
                        len(rows),
                        f"ExtractedTable{idx}",
                        table_index,
                    )
                    if first_table_header_row is None:
                        first_table_header_row = table_start_row + 1
                    current_row += 2
                    wrote_structured = True
            except Exception:
                pass

            if not wrote_structured:
                inferred_rows = _merge_wrapped_table_rows(
                    _extract_grid_like_table(page_dict, float(page.rect.width))
                )
                inferred_rows = [[_normalize_table_cell_text(cell) for cell in row] for row in inferred_rows]
                if inferred_rows and _is_useful_structured_table(inferred_rows):
                    ws.cell(current_row, 1, "Inferred table")
                    ws.cell(current_row, 1).font = Font(bold=True)
                    current_row += 1
                    table_start_row = current_row
                    header_row = [str(cell).strip() if cell is not None else "" for cell in inferred_rows[0]] if inferred_rows else []
                    for row_index, row in enumerate(inferred_rows):
                        for col_index, value in enumerate(row, 1):
                            cell = ws.cell(current_row, col_index)
                            header_value = header_row[col_index - 1] if row_index > 0 and col_index - 1 < len(header_row) else ""
                            _write_excel_table_cell(cell, value, header_value=header_value)
                            if row_index == 0:
                                _apply_excel_table_header_style(cell)
                        current_row += 1
                    _add_excel_table(
                        ws,
                        table_start_row,
                        max(len(row) for row in inferred_rows) if inferred_rows else 0,
                        len(inferred_rows),
                        f"InferredTable{idx}",
                        1,
                    )
                    if first_table_header_row is None:
                        first_table_header_row = table_start_row + 1
                    current_row += 2
                    wrote_structured = True

            narrative_blocks: list[list[str]] = []
            for lines, bbox in _sorted_text_blocks(page_dict, float(page.rect.width)):
                if any(_bbox_overlaps(bbox, table_bbox) for table_bbox in table_bboxes):
                    continue
                block_lines = _block_plain_text(lines)
                if wrote_structured:
                    dense_table_like_block = (
                        len(block_lines) >= 4
                        and sum(len(line) for line in block_lines) >= 180
                        and (bbox[2] - bbox[0]) >= float(page.rect.width) * 0.55
                    )
                    if dense_table_like_block:
                        continue
                if block_lines:
                    narrative_blocks.append(block_lines)

            if narrative_blocks:
                ws.cell(current_row, 1, "Text lines")
                ws.cell(current_row, 1).font = Font(bold=True)
                current_row += 1
                for block_lines in narrative_blocks:
                    for text in block_lines:
                        cell = ws.cell(current_row, 1, text)
                        cell.alignment = Alignment(vertical="top", wrap_text=True)
                        current_row += 1
                    current_row += 1
            elif not wrote_structured:
                ws.cell(current_row, 1, "No extractable table or text blocks found.")
                ws.cell(current_row, 1).alignment = Alignment(vertical="top", wrap_text=True)
                current_row += 1

            max_col = min(ws.max_column, 8)
            for col_idx in range(1, max_col + 1):
                column_letter = ws.cell(1, col_idx).column_letter
                max_len = 14
                for row in ws.iter_rows(min_col=col_idx, max_col=col_idx, max_row=ws.max_row):
                    value = row[0].value
                    if value is not None:
                        max_len = min(60, max(max_len, len(str(value)) + 2))
                ws.column_dimensions[column_letter].width = max_len
            if first_table_header_row:
                ws.freeze_panes = f"A{first_table_header_row}"
        wb.save(str(output_path))


def _strip_bold(s: str) -> str:
    """Remove surrounding **...** that models sometimes wrap headings in."""
    if s.startswith("**") and s.endswith("**") and len(s) > 4:
        return s[2:-2].strip()
    return s


def _markdown_to_docx(text: str, doc) -> None:
    """Convert simple markdown lines into python-docx content."""
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            doc.add_heading(_strip_bold(s[4:]), level=3)
        elif s.startswith("## "):
            doc.add_heading(_strip_bold(s[3:]), level=2)
        elif s.startswith("# "):
            doc.add_heading(_strip_bold(s[2:]), level=1)
        elif s.startswith(("- ", "* ", "• ")):
            doc.add_paragraph(s[2:], style="List Bullet")
        elif len(s) > 1 and s[0].isdigit() and s[1] in ".)" and s[2:3] == " ":
            doc.add_paragraph(s[s.index(" ") + 1:], style="List Number")
        else:
            # Strip inline bold/italic markers for plain paragraphs
            clean = s.replace("**", "").replace("__", "")
            doc.add_paragraph(clean)


# ---------------------------------------------------------------------------
# AI text helpers  (text in → text out — Ollama, GitHub Models, OpenAI)
# ---------------------------------------------------------------------------

def _ollama_text(system_prompt: str, user_content: str) -> str:
    """Send plain text to Ollama (no image). Returns response."""
    _, model = _ollama_check()
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_content},
        ],
        "stream": False,
    }
    resp = _ollama_post(payload, timeout=300.0)
    return resp.get("message", {}).get("content", "").strip()


def _github_text(system_prompt: str, user_content: str) -> str:
    """Send plain text to GitHub Models. Returns response."""
    import os
    try:
        from openai import OpenAI
    except ImportError:
        raise RuntimeError("openai package not installed.")
    token = os.environ.get("GITHUB_TOKEN")
    if not token:
        raise RuntimeError("GITHUB_TOKEN is not set. Configure it as an environment variable.")
    model = os.environ.get("GITHUB_OCR_MODEL", "openai/gpt-4.1-mini")
    client = OpenAI(base_url="https://models.inference.ai.azure.com", api_key=token)
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_content},
        ],
        max_tokens=4096,
    )
    return resp.choices[0].message.content.strip()


def _openai_text(system_prompt: str, user_content: str) -> str:
    """Send plain text to OpenAI. Returns response."""
    import os
    try:
        from openai import OpenAI
    except ImportError:
        raise RuntimeError("openai package not installed.")
    model = os.environ.get("OPENAI_OCR_MODEL", "gpt-4.1-mini")
    client = OpenAI()
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_content},
        ],
        max_tokens=4096,
    )
    return resp.choices[0].message.content.strip()


def _ai_text(system_prompt: str, user_content: str, engine: str) -> str:
    """Route to the correct AI engine. engine: auto | ollama | github | openai"""
    if engine == "auto":
        errors: list[str] = []
        for name, fn in [("ollama", _ollama_text), ("github", _github_text)]:
            try:
                return fn(system_prompt, user_content)
            except Exception as e:
                errors.append(f"{name}: {e}")
        raise RuntimeError(
            "No AI engine available. Run 'ollama serve' or set GITHUB_TOKEN as an environment variable.\n"
            + "\n".join(errors)
        )
    dispatch = {"ollama": _ollama_text, "github": _github_text, "openai": _openai_text}
    if engine not in dispatch:
        raise ValueError(f"Unknown engine: {engine!r}. Choose: auto, ollama, github, openai.")
    return dispatch[engine](system_prompt, user_content)


# ---------------------------------------------------------------------------
# Edit PDF — AI helpers
# ---------------------------------------------------------------------------

_AI_FIX_ACTIONS = {
    "rephrase":  "Rephrase the following text to sound more natural and fluent. Return only the rewritten text, nothing else.",
    "grammar":   "Fix all spelling and grammar errors in the following text. Return only the corrected text, nothing else.",
    "formal":    "Rewrite the following text in a professional, formal tone. Return only the rewritten text, nothing else.",
    "shorter":   "Shorten the following text while preserving its meaning. Return only the shortened text, nothing else.",
    "bullet":    "Rewrite the following text as a concise bulleted list. Return only the bullets, nothing else.",
}

def ai_fix_text(text: str, action: str, engine: str = "auto") -> str:
    """Apply an AI rewrite action to a text overlay string."""
    system = _AI_FIX_ACTIONS.get(action)
    if not system:
        raise ValueError(f"Unknown AI fix action: {action!r}. Valid: {list(_AI_FIX_ACTIONS)}")
    if not text.strip():
        raise ValueError("Text is empty — nothing to fix.")
    result = _ai_text(system, text.strip(), engine)
    # Trim surrounding quotes/markdown the model sometimes adds
    result = result.strip().strip('"').strip("'")
    if result.startswith("```") and result.endswith("```"):
        result = result[3:-3].strip()
    return result or text


def ai_suggest_stamp(text: str, engine: str = "auto") -> str:
    """Given text extracted from a PDF page, suggest the best stamp label."""
    system = (
        "You are a document classification assistant. "
        "Given text from a PDF document, reply with exactly ONE word or short phrase "
        "that best describes the document status. "
        "Choose from: DRAFT, APPROVED, REJECTED, CONFIDENTIAL, COPY, FOR REVIEW, VOID, PAID. "
        "Reply with only that single word or phrase — no explanation, no punctuation."
    )
    if not text.strip():
        return "DRAFT"
    result = _ai_text(system, text[:2000], engine).strip().upper()
    allowed = {"DRAFT", "APPROVED", "REJECTED", "CONFIDENTIAL", "COPY", "FOR REVIEW", "VOID", "PAID"}
    # Return the first allowed token found in the response, fallback to DRAFT
    for token in allowed:
        if token in result:
            return token
    return result.split("\n")[0][:20] or "DRAFT"


_WORD_PROMPT = (
    "Extract all visible text from this document page as clean markdown.\n"
    "Rules:\n"
    "- Use # for the largest/title text, ## for section headings, ### for sub-headings.\n"
    "- Use - for bullet points, 1. 2. 3. for numbered lists.\n"
    "- For tables use pipe rows: | Col1 | Col2 |\n"
    "- Skip decorative borders, watermarks, and page numbers.\n"
    "- Return ONLY the extracted content — no commentary, no explanations."
)

_TEXT_PROMPT = (
    "Extract all text from this document page in reading order. "
    "Preserve paragraph breaks. "
    "Return only the text content with no commentary or extra formatting."
)


# ---------------------------------------------------------------------------
# PDF → Text
# ---------------------------------------------------------------------------

def pdf_to_text(input_path: Path, output_path: Path, engine: str = "basic") -> None:
    if engine in ("ollama", "auto"):
        try:
            import tempfile
            with tempfile.TemporaryDirectory() as td:
                pages = _render_pdf_pages(input_path, Path(td))
                lines: list[str] = []
                for i, img in enumerate(pages, 1):
                    lines.append(f"--- Page {i} ---")
                    lines.append(_ollama_call_page(img, _TEXT_PROMPT))
                    lines.append("")
            output_path.write_text("\n".join(lines), encoding="utf-8")
            return
        except Exception as exc:
            if engine == "ollama" and not _is_recoverable_ollama_runtime_error(exc):
                raise
            # auto: fall through to basic

    # basic: pypdf text extraction
    page_texts = _extract_pdf_text_pages(input_path)
    lines = []
    for i, text in enumerate(page_texts, 1):
        lines.append(f"--- Page {i} ---")
        lines.append(text)
        lines.append("")
    output_path.write_text("\n".join(lines), encoding="utf-8")


# ---------------------------------------------------------------------------
# PDF → Word  (semantic extraction or visual-fidelity page snapshots)
# ---------------------------------------------------------------------------

def pdf_to_word(
    input_path: Path,
    output_path: Path,
    engine: str = "basic",
    auto_pair_small_pages: bool = True,
) -> None:
    if engine == "fidelity":
        _pdf_pages_to_docx_hybrid(input_path, output_path, auto_pair_small_pages=auto_pair_small_pages)
        return
    if engine == "auto":
        try:
            _pdf_pages_to_docx_hybrid(input_path, output_path, auto_pair_small_pages=auto_pair_small_pages)
            return
        except Exception:
            pass

    if engine in ("ollama", "auto"):
        try:
            from docx import Document
            import tempfile
            doc = Document()
            with tempfile.TemporaryDirectory() as td:
                pages = _render_pdf_pages(input_path, Path(td))
                for i, img in enumerate(pages, 1):
                    if i > 1:
                        doc.add_page_break()
                    text = _ollama_call_page(img, _WORD_PROMPT)
                    _markdown_to_docx(text, doc)
            doc.save(str(output_path))
            return
        except Exception as exc:
            if engine == "ollama" and not _is_recoverable_ollama_runtime_error(exc):
                raise
            # auto: fall through to basic

    # basic: pypdf text extraction with heading/list detection
    import re as _re
    from docx import Document
    page_texts = _extract_pdf_text_pages(input_path)
    doc = Document()

    _bullet_re = _re.compile(r"^[\-\*\u2022\u2023\u25e6\u2013\u2014]\s+")
    _numbered_re = _re.compile(r"^\d+[\.\)]\s+")

    def _classify_line(line: str, next_line: str) -> str:
        """Return 'heading', 'bullet', 'numbered', or 'body'."""
        if not line:
            return "body"
        if _bullet_re.match(line):
            return "bullet"
        if _numbered_re.match(line):
            return "numbered"
        # Short ALL-CAPS line with no trailing punctuation → heading
        stripped = line.rstrip(".,:;")
        if stripped == stripped.upper() and len(stripped) <= 80 and len(stripped) > 2 and stripped.replace(" ", "").isalpha():
            return "heading"
        # Short line followed by empty or shorter continuation → heading
        if len(line) <= 60 and (not next_line or len(next_line) > len(line) * 1.5):
            if not line[-1] in ".,:;?!":
                return "heading"
        return "body"

    for i, text in enumerate(page_texts, 1):
        if i > 1:
            doc.add_page_break()
        lines = text.split("\n")
        for j, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            next_line = lines[j + 1].strip() if j + 1 < len(lines) else ""
            kind = _classify_line(line, next_line)
            if kind == "heading":
                doc.add_paragraph(line, style="Heading 2")
            elif kind == "bullet":
                clean = _bullet_re.sub("", line)
                doc.add_paragraph(clean, style="List Bullet")
            elif kind == "numbered":
                clean = _numbered_re.sub("", line)
                doc.add_paragraph(clean, style="List Number")
            else:
                doc.add_paragraph(line)
    doc.save(str(output_path))


# ---------------------------------------------------------------------------
# PDF → Excel
# ---------------------------------------------------------------------------

_EXCEL_PROMPT = (
    "Extract all data from this document page into pipe-delimited format for a spreadsheet.\n"
    "Rules:\n"
    "- Output each row as pipe-separated values using | as the delimiter.\n"
    "- For tables: keep column headers in the first row.\n"
    "- For multiple tables: separate them with one blank line.\n"
    "- For non-tabular content: output each line as a single-column row.\n"
    "- Preserve numbers exactly as shown (e.g. 1,250 stays as 1,250 — do NOT split on commas).\n"
    "- Do NOT quote values.\n"
    "- Return ONLY the pipe-delimited data — no commentary, no markdown fences."
)


def pdf_to_excel(input_path: Path, output_path: Path, engine: str = "basic") -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment

    if engine == "fidelity":
        _pdf_pages_to_excel_hybrid(input_path, output_path)
        return
    if engine == "auto":
        try:
            _pdf_pages_to_excel_hybrid(input_path, output_path)
            return
        except Exception:
            pass

    wb = Workbook()
    wb.remove(wb.active)  # drop default blank sheet

    if engine in ("ollama", "auto"):
        try:
            import csv, io, tempfile
            with tempfile.TemporaryDirectory() as td:
                pages = _render_pdf_pages(input_path, Path(td))
                for i, img in enumerate(pages, 1):
                    raw = _ollama_call_page(img, _EXCEL_PROMPT)
                    ws = wb.create_sheet(title=f"Page {i}")
                    reader = csv.reader(io.StringIO(raw), delimiter="|")
                    first = True
                    for row in reader:
                        ws.append(row)
                        if first:
                            for cell in ws[ws.max_row]:
                                cell.font = Font(bold=True)
                            first = False
            if not wb.sheetnames:
                wb.create_sheet("Sheet1")
            wb.save(str(output_path))
            return
        except Exception as exc:
            if engine == "ollama" and not _is_recoverable_ollama_runtime_error(exc):
                raise
            # auto: fall through to basic

    # basic: pymupdf table detection → openpyxl
    import fitz  # type: ignore
    doc = fitz.open(str(input_path))
    try:
        for page_num, page in enumerate(doc, 1):
            ws = wb.create_sheet(title=f"Page {page_num}")
            written = False
            try:
                tabs = page.find_tables()
                for t_idx, table in enumerate(tabs.tables):
                    if t_idx > 0:
                        ws.append([])  # blank row between tables
                    rows = table.extract()
                    for r_idx, row in enumerate(rows):
                        clean = [str(c).strip() if c is not None else "" for c in row]
                        ws.append(clean)
                        if r_idx == 0:
                            for cell in ws[ws.max_row]:
                                cell.font = Font(bold=True)
                    written = bool(tabs.tables)
            except Exception:
                pass

            if not written:
                # fallback: dump text line by line
                for line in page.get_text().split("\n"):
                    line = line.strip()
                    if line:
                        ws.append([line])

            # auto-width columns (cap at 50)
            for col in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col), default=8)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)
    finally:
        doc.close()

    if not wb.sheetnames:
        wb.create_sheet("Sheet1")
    wb.save(str(output_path))


# ---------------------------------------------------------------------------
# Images → PDF
# ---------------------------------------------------------------------------

def images_to_pdf(input_paths: Sequence[Path], output_path: Path) -> None:
    def _build(out: Path) -> None:
        from PIL import Image as PILImage
        images = []
        image_details: list[dict[str, object]] = []
        try:
            for p in input_paths:
                details: dict[str, object] = {
                    "filename": p.name,
                    "path": str(p),
                }
                try:
                    details["size_bytes"] = p.stat().st_size
                except OSError:
                    pass
                try:
                    with PILImage.open(p) as img:
                        details.update(
                            {
                                "format": img.format,
                                "mode": img.mode,
                                "size": img.size,
                            }
                        )
                        converted = img.convert("RGB")
                except Exception as exc:
                    raise RuntimeError(
                        f"Could not load image for PDF conversion: {details}"
                    ) from exc

                images.append(converted)
                image_details.append(details)

            if not images:
                raise ValueError("No images provided.")

            try:
                images[0].save(str(out), save_all=True, append_images=images[1:])
            except Exception as exc:
                raise RuntimeError(
                    f"Could not save PDF from images: {image_details}"
                ) from exc
        finally:
            for image in images:
                image.close()

    _conversion_pipeline([("pil", _build)], output_path)


# ---------------------------------------------------------------------------
# COM helpers — Word, Excel  (Windows / Office 2016+ required)
# ---------------------------------------------------------------------------

def _com_office_export(app_progid: str, open_fn: Callable, output_path: Path) -> None:
    """Generic COM export helper. Raises RuntimeError on failure."""
    try:
        import win32com.client
        import pythoncom
    except ImportError:
        raise RuntimeError("pywin32 not installed.")

    pythoncom.CoInitialize()
    app = None
    doc = None
    try:
        app = win32com.client.Dispatch(app_progid)
        app.DisplayAlerts = False
        doc = open_fn(app)
        # 17 = wdFormatPDF (Word), 57 = xlTypePDF (Excel) — caller passes correct constant
        doc._com_export(str(output_path.resolve()))
    finally:
        if doc is not None:
            try: doc._com_close()
            except Exception: pass
        if app is not None:
            try: app.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


def _word_to_pdf_via_com(input_path: Path, output_path: Path) -> None:
    """Export .docx → PDF using the real Word app (full fidelity)."""
    try:
        import win32com.client, pythoncom
    except ImportError:
        raise RuntimeError("pywin32 not installed.")

    src = str(input_path.resolve())
    dst = str(output_path.resolve())
    pythoncom.CoInitialize()
    app = doc = None
    try:
        app = win32com.client.Dispatch("Word.Application")
        app.Visible = False
        app.DisplayAlerts = False
        doc = app.Documents.Open(src, ReadOnly=True)
        doc.SaveAs(dst, FileFormat=17)   # 17 = wdFormatPDF
    except Exception as exc:
        raise RuntimeError(f"Word COM export failed: {exc}") from exc
    finally:
        if doc:
            try: doc.Close(False)
            except Exception: pass
        if app:
            try: app.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


_EXCEL_PAPER_SIZE = {
    "a4": 9,       # xlPaperA4
    "a3": 8,       # xlPaperA3
    "letter": 1,   # xlPaperLetter
    "legal": 5,    # xlPaperLegal
}


def _excel_to_pdf_via_com(
    input_path: Path,
    output_path: Path,
    paper_size: str = "a4",
    orientation: str = "landscape",
    fit_pages_wide: int = 0,
) -> None:
    """Export .xlsx → PDF using the real Excel app (full fidelity).

    fit_pages_wide:
      0  = no scaling (use workbook's own zoom/page breaks)
      1  = fit entire sheet to 1 page wide (tall auto)
      2+ = fit to N pages wide (tall auto)
      -1 = fit entire sheet onto a single page
    """
    try:
        import win32com.client, pythoncom
    except ImportError:
        raise RuntimeError("pywin32 not installed.")

    src = str(input_path.resolve())
    dst = str(output_path.resolve())
    xl_paper = _EXCEL_PAPER_SIZE.get(paper_size, 9)
    xl_orient = 2 if orientation == "landscape" else 1  # xlLandscape=2, xlPortrait=1

    pythoncom.CoInitialize()
    app = wb = None
    try:
        app = win32com.client.Dispatch("Excel.Application")
        app.Visible = False
        app.DisplayAlerts = False
        wb = app.Workbooks.Open(src, ReadOnly=False)
        for ws in wb.Worksheets:
            ps = ws.PageSetup
            ps.PaperSize = xl_paper
            ps.Orientation = xl_orient
            # Always clamp print area to used range to eliminate blank pages
            used = ws.UsedRange
            if used is not None:
                ps.PrintArea = used.Address
            if fit_pages_wide == -1:
                # Fit entire sheet to one page
                ps.Zoom = False
                ps.FitToPagesWide = 1
                ps.FitToPagesTall = 1
            elif fit_pages_wide > 0:
                # Fit to N pages wide, as many rows tall as needed
                # FitToPagesTall=9999 means effectively unlimited (0 is rejected by COM)
                ps.Zoom = False
                ps.FitToPagesWide = fit_pages_wide
                ps.FitToPagesTall = 9999
            # else: leave zoom/scaling untouched (workbook defaults)
        wb.ExportAsFixedFormat(0, dst)   # 0 = xlTypePDF
    except Exception as exc:
        raise RuntimeError(f"Excel COM export failed: {exc}") from exc
    finally:
        if wb:
            try: wb.Close(False)
            except Exception: pass
        if app:
            try: app.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


# ---------------------------------------------------------------------------
# Word → PDF
# ---------------------------------------------------------------------------

def _word_to_pdf_basic(input_path: Path, output_path: Path) -> None:
    """Fallback: python-docx text extraction + fpdf2."""
    from docx import Document
    from fpdf import FPDF

    doc = Document(str(input_path))
    pdf = FPDF()
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    col_w = pdf.w - pdf.l_margin - pdf.r_margin
    pdf.set_font("Helvetica", size=11)

    for para in doc.paragraphs:
        text = _sanitize_latin1(para.text.strip())
        if not text:
            pdf.ln(5)
            continue
        style = para.style.name.lower()
        if "heading 1" in style:
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(col_w, 10, text)
            pdf.ln(2)
            pdf.set_font("Helvetica", size=11)
        elif "heading 2" in style:
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(col_w, 8, text)
            pdf.ln(1)
            pdf.set_font("Helvetica", size=11)
        else:
            pdf.multi_cell(col_w, 7, text)
    pdf.output(str(output_path))


def word_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert Word → PDF.

    Engine chain (tried in order, each QA-checked by Ollama):
      1. Word COM    — full fidelity: fonts, colors, images, tables, shapes (Windows + Office only)
      2. LibreOffice — high fidelity headless fallback (Linux / any platform with soffice on PATH)
      3. Basic       — text-only fallback via python-docx + fpdf2
    """
    engines: list[tuple[str, Callable[[Path], None]]] = []
    if _powerpoint_available():
        engines.append(("word_com", lambda o: _word_to_pdf_via_com(input_path, o)))
    if _libreoffice_available():
        engines.append(("libreoffice", lambda o: _office_to_pdf_via_libreoffice(input_path, o)))
    engines.append(("basic", lambda o: _word_to_pdf_basic(input_path, o)))
    _conversion_pipeline(engines, output_path)


# ---------------------------------------------------------------------------
# Repair PDF  (re-save through pikepdf which fixes common issues)
# ---------------------------------------------------------------------------

def repair_pdf(input_path: Path, output_path: Path) -> None:
    import pikepdf
    with pikepdf.open(str(input_path), suppress_warnings=True) as pdf:
        pdf.save(str(output_path))


# ---------------------------------------------------------------------------
# AI Summarizer
# ---------------------------------------------------------------------------

_SUMMARY_PAGE_PROMPT = (
    "Extract the key points from this document page as a concise bulleted list.\n"
    "Include specific names, numbers, dates, and decisions.\n"
    "Return only the bullet points — no preamble, no commentary."
)

_SUMMARY_FINAL_PROMPTS = {
    "brief": (
        "Write a brief 2–3 sentence summary of this document based on the key points below.\n"
        "Be direct and factual.\n\nKey points:\n{notes}\n\nReturn only the summary."
    ),
    "standard": (
        "Write a structured summary of this document based on the key points below.\n"
        "Include: a short overview, the main topics, and key facts or conclusions.\n\n"
        "Key points:\n{notes}\n\nReturn only the summary."
    ),
    "detailed": (
        "Write a detailed structured summary of this document based on the key points below.\n"
        "Include: overview, all major sections, key data and figures, conclusions, and any action items.\n"
        "Use headings and bullet points to organise the output.\n\n"
        "Key points:\n{notes}\n\nReturn only the summary."
    ),
}


def summarize_pdf(
    input_path: Path,
    output_path: Path,
    engine: str = "auto",
    length: str = "standard",
    output_format: str = "txt",
) -> None:
    import tempfile

    # Step 1: extract key points per page
    # Ollama uses vision (works on scans); GitHub/OpenAI use text extraction (faster)
    page_notes: list[str] = []
    if engine in ("ollama", "auto"):
        try:
            with tempfile.TemporaryDirectory() as td:
                pages = _render_pdf_pages(input_path, Path(td))
                for i, img in enumerate(pages, 1):
                    notes = _ollama_call_page(img, _SUMMARY_PAGE_PROMPT)
                    page_notes.append(f"[Page {i}]\n{notes}")
        except Exception as exc:
            if engine == "ollama" and not _is_recoverable_ollama_runtime_error(exc):
                raise
            page_notes = []  # fall through to text extraction

    if not page_notes:
        # Text extraction path (github / openai / auto fallback)
        for i, text in enumerate(_extract_pdf_text_pages(input_path), 1):
            if text:
                page_notes.append(f"[Page {i}]\n{text}")

    if not page_notes:
        raise ValueError("No extractable text found in this PDF.")

    # Step 2: final summary — always text-only (fast for all engines)
    combined = "\n\n".join(page_notes)
    final_prompt = _SUMMARY_FINAL_PROMPTS.get(length, _SUMMARY_FINAL_PROMPTS["standard"]).format(notes=combined)
    system = "You are a professional document summarizer. Return only the requested summary."

    effective_engine = engine if engine not in ("auto",) else "auto"
    # For auto after ollama vision succeeded, use ollama for the final pass too
    if page_notes and engine == "auto":
        try:
            summary = _ollama_text(system, final_prompt)
        except Exception:
            try:
                summary = _github_text(system, final_prompt)
            except Exception:
                summary = _openai_text(system, final_prompt)
    elif engine == "ollama":
        try:
            summary = _ollama_text(system, final_prompt)
        except Exception as exc:
            if not _is_recoverable_ollama_runtime_error(exc):
                raise
            summary = _fallback_text_ai(system, final_prompt)
    else:
        summary = _ai_text(system, final_prompt, effective_engine)

    if output_format == "pdf":
        from fpdf import FPDF
        pdf = FPDF()
        pdf.set_margins(18, 18, 18)
        pdf.set_auto_page_break(auto=True, margin=18)
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 15)
        pdf.cell(0, 10, "Document Summary", ln=True)
        pdf.ln(3)
        pdf.set_font("Helvetica", size=11)
        col_w = pdf.w - pdf.l_margin - pdf.r_margin
        for line in summary.split("\n"):
            line = line.strip()
            if not line:
                pdf.ln(4)
            elif line.startswith("# "):
                pdf.set_font("Helvetica", "B", 13)
                pdf.multi_cell(col_w, 8, line[2:])
                pdf.set_font("Helvetica", size=11)
            elif line.startswith("## "):
                pdf.set_font("Helvetica", "B", 11)
                pdf.multi_cell(col_w, 7, line[3:])
                pdf.set_font("Helvetica", size=11)
            else:
                pdf.multi_cell(col_w, 7, line)
        pdf.output(str(output_path))
    elif output_format == "docx":
        from docx import Document
        doc = Document()
        doc.add_heading("Document Summary", 0)
        _markdown_to_docx(summary, doc)
        doc.save(str(output_path))
    else:
        output_path.write_text(summary, encoding="utf-8")


# ---------------------------------------------------------------------------
# Translate PDF
# ---------------------------------------------------------------------------

_TRANSLATE_PROMPT_TPL = (
    "Translate all text on this document page to {language}.\n"
    "Rules:\n"
    "- Use # for the largest/title text, ## for section headings, ### for sub-headings.\n"
    "- Use - for bullet points, 1. 2. 3. for numbered lists.\n"
    "- For tables use pipe rows: | Col1 | Col2 |\n"
    "- Preserve every line of text — do not skip or summarise anything.\n"
    "- Return ONLY the translated content in clean markdown — no commentary."
)


def translate_pdf(
    input_path: Path,
    output_path: Path,
    engine: str = "auto",
    language: str = "Spanish",
    output_format: str = "docx",
) -> None:
    import tempfile
    vision_prompt = _TRANSLATE_PROMPT_TPL.format(language=language)
    system_text = (
        f"You are a professional translator. Translate the text to {language}. "
        "Rules: use # for the largest/title text, ## for section headings, ### for sub-headings, "
        "- for bullets, 1. 2. 3. for numbered lists. "
        "Preserve every line — do not skip or summarise. "
        "Return only the translated content in clean markdown."
    )

    translated_pages: list[str] = []

    if engine in ("ollama", "auto"):
        try:
            with tempfile.TemporaryDirectory() as td:
                pages = _render_pdf_pages(input_path, Path(td))
                for img in pages:
                    translated_pages.append(_ollama_call_page(img, vision_prompt))
        except Exception as exc:
            if engine == "ollama" and not _is_recoverable_ollama_runtime_error(exc):
                raise
            translated_pages = []

    if not translated_pages:
        # Text extraction path for github / openai / auto fallback
        for text in _extract_pdf_text_pages(input_path):
            if not text:
                translated_pages.append("")
                continue
            if engine == "ollama":
                try:
                    translated_pages.append(_ollama_text(system_text, text))
                except Exception as exc:
                    if not _is_recoverable_ollama_runtime_error(exc):
                        raise
                    translated_pages.append(_fallback_text_ai(system_text, text))
            else:
                translated_pages.append(_ai_text(system_text, text, engine))

    if output_format == "pdf":
        from fpdf import FPDF
        pdf = FPDF()
        pdf.set_margins(15, 15, 15)
        pdf.set_auto_page_break(auto=True, margin=15)
        col_w = 0.0
        for i, text in enumerate(translated_pages, 1):
            pdf.add_page()
            if col_w == 0:
                col_w = pdf.w - pdf.l_margin - pdf.r_margin
            pdf.set_font("Helvetica", "I", 8)
            pdf.set_text_color(150, 150, 150)
            pdf.cell(col_w, 5, f"Page {i}", align="R")
            pdf.ln(6)
            pdf.set_text_color(0, 0, 0)
            pdf.set_font("Helvetica", size=11)
            for line in text.split("\n"):
                line = line.strip()
                if not line:
                    pdf.ln(4)
                elif line.startswith("# "):
                    pdf.set_font("Helvetica", "B", 14)
                    pdf.multi_cell(col_w, 9, line[2:])
                    pdf.set_font("Helvetica", size=11)
                elif line.startswith("## "):
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.multi_cell(col_w, 8, line[3:])
                    pdf.set_font("Helvetica", size=11)
                else:
                    pdf.multi_cell(col_w, 7, line)
        pdf.output(str(output_path))
    elif output_format == "txt":
        output_path.write_text("\n\n".join(translated_pages), encoding="utf-8")
    else:
        from docx import Document
        doc = Document()
        for i, text in enumerate(translated_pages, 1):
            if i > 1:
                doc.add_page_break()
            _markdown_to_docx(text, doc)
        doc.save(str(output_path))


# ---------------------------------------------------------------------------
# PDF → PowerPoint
# ---------------------------------------------------------------------------

_PPTX_SLIDE_PROMPT = (
    "Extract the slide content from this document page.\n"
    "Return in this exact format:\n"
    "TITLE: <the main heading or title>\n"
    "BULLETS:\n"
    "- <key point 1>\n"
    "- <key point 2>\n"
    "(up to 6 bullets maximum)\n\n"
    "If there is no clear title, infer one from the content.\n"
    "Return ONLY this structure — no other text."
)


def _parse_slide_content(text: str) -> tuple[str, list[str]]:
    """Parse TITLE:/BULLETS: format from Ollama response."""
    title = "Slide"
    bullets: list[str] = []
    in_bullets = False
    for line in text.split("\n"):
        s = line.strip()
        if s.upper().startswith("TITLE:"):
            title = s[6:].strip()
        elif s.upper().startswith("BULLETS:"):
            in_bullets = True
        elif in_bullets and s.startswith("- "):
            bullets.append(s[2:])
        elif in_bullets and s.startswith("• "):
            bullets.append(s[2:])
    return title, bullets


def pdf_to_pptx(
    input_path: Path,
    output_path: Path,
    engine: str = "basic",
    auto_pair_small_pages: bool = True,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if engine == "fidelity":
        _pdf_pages_to_pptx_hybrid(input_path, output_path, auto_pair_small_pages=auto_pair_small_pages)
        return
    if engine == "auto":
        try:
            _pdf_pages_to_pptx_hybrid(input_path, output_path, auto_pair_small_pages=auto_pair_small_pages)
            return
        except Exception:
            pass

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[1]  # title + content
    blank_layout = prs.slide_layouts[6]  # blank

    if engine in ("ollama", "github", "openai", "auto"):
        # AI path: for ollama use vision; for github/openai use text extraction
        try:
            import tempfile
            if engine in ("ollama", "auto"):
                with tempfile.TemporaryDirectory() as td:
                    pages = _render_pdf_pages(input_path, Path(td))
                    for img in pages:
                        raw = _ollama_call_page(img, _PPTX_SLIDE_PROMPT)
                        title_text, bullets = _parse_slide_content(raw)
                        slide = prs.slides.add_slide(title_layout)
                        slide.shapes.title.text = title_text
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        for i, bullet in enumerate(bullets):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
            else:
                # github / openai: text extraction per page → AI structuring
                import fitz  # type: ignore
                doc_fitz = fitz.open(str(input_path))
                system = "You extract structured content for presentation slides."
                try:
                    for page_num, page in enumerate(doc_fitz, 1):
                        text = page.get_text().strip()
                        if not text:
                            continue
                        raw = _ai_text(system, f"{_PPTX_SLIDE_PROMPT}\n\nPage content:\n{text}", engine)
                        title_text, bullets = _parse_slide_content(raw)
                        slide = prs.slides.add_slide(title_layout)
                        slide.shapes.title.text = title_text or f"Slide {page_num}"
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        for i, bullet in enumerate(bullets):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                finally:
                    doc_fitz.close()
            prs.save(str(output_path))
            return
        except Exception:
            if engine not in ("auto",):
                raise

    # basic: text extraction, each page = one slide
    import fitz  # type: ignore
    doc = fitz.open(str(input_path))
    try:
        for page_num, page in enumerate(doc, 1):
            text = page.get_text().strip()
            lines = [l.strip() for l in text.split("\n") if l.strip()]

            # Skip trivial leading lines (page numbers, "Section X", single digits)
            title_text = f"Slide {page_num}"
            body_start = 0
            for idx, candidate in enumerate(lines):
                stripped = candidate.strip("0123456789. ").lower()
                if len(candidate) > 3 and stripped not in ("page", "section", "slide", "chapter"):
                    title_text = candidate[:80]
                    body_start = idx + 1
                    break

            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title_text

            body_lines = lines[body_start:body_start + 8]  # up to 8 bullets
            if body_lines:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, line in enumerate(body_lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line[:150]
                    p.level = 0
    finally:
        doc.close()

    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Crop PDF  (trim margins by mm from each side)
# ---------------------------------------------------------------------------

def crop_pdf(
    input_path: Path,
    output_path: Path,
    left_mm: float = 0,
    top_mm: float = 0,
    right_mm: float = 0,
    bottom_mm: float = 0,
) -> None:
    from pypdf import PdfReader, PdfWriter
    MM_TO_PT = 72 / 25.4
    reader = PdfReader(str(input_path))
    writer = PdfWriter()
    for page in reader.pages:
        w_pt = float(page.mediabox.width)
        h_pt = float(page.mediabox.height)
        l = left_mm * MM_TO_PT
        r = right_mm * MM_TO_PT
        t = top_mm * MM_TO_PT
        b = bottom_mm * MM_TO_PT
        if l + r >= w_pt or t + b >= h_pt:
            raise ValueError("Crop margins exceed page dimensions.")
        page.cropbox.lower_left = (l, b)
        page.cropbox.upper_right = (w_pt - r, h_pt - t)
        writer.add_page(page)
    with output_path.open("wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# PDF → PDF/A  (best-effort: re-saves with PDF/A-1b XMP metadata)
# ---------------------------------------------------------------------------

def pdf_to_pdfa(input_path: Path, output_path: Path) -> None:
    import pikepdf
    with pikepdf.open(str(input_path), suppress_warnings=True) as pdf:
        with pdf.open_metadata(set_pikepdf_as_editor=False) as meta:
            meta["pdfaid:part"] = "1"
            meta["pdfaid:conformance"] = "B"
        pdf.save(str(output_path))


# ---------------------------------------------------------------------------
# Shared print-setup helpers
# ---------------------------------------------------------------------------

_PAPER_MM = {
    "a4":     (210, 297),
    "a3":     (297, 420),
    "letter": (216, 279),
    "legal":  (216, 356),
}

_PAPER_CSS_NAME = {"a4": "A4", "a3": "A3", "letter": "letter", "legal": "legal"}


def _make_fpdf(paper_size: str, orientation: str, margin_mm: int):
    """Return (pdf, content_w_mm, content_h_mm) for the given page setup."""
    from fpdf import FPDF
    w, h = _PAPER_MM.get(paper_size.lower(), (210, 297))
    if orientation == "landscape":
        w, h = h, w
    pdf = FPDF(unit="mm", format=(w, h))
    pdf.set_margins(margin_mm, margin_mm, margin_mm)
    pdf.set_auto_page_break(auto=True, margin=margin_mm)
    cw = w - 2 * margin_mm
    ch = h - 2 * margin_mm
    return pdf, cw, ch


def _page_css(paper_size: str, orientation: str, margin_mm: int) -> str:
    """CSS @page rule for xhtml2pdf."""
    name = _PAPER_CSS_NAME.get(paper_size.lower(), "A4")
    orient = "landscape" if orientation == "landscape" else "portrait"
    return (
        f"@page{{size:{name} {orient};margin:{margin_mm}mm}}"
        "body{font-family:Arial,sans-serif;margin:0}"
    )


# ---------------------------------------------------------------------------
# PowerPoint → PDF
# ---------------------------------------------------------------------------

# ppPrintOutputSlides=1, TwoSlides=2, ThreeSlides=3, SixSlides=4, FourSlides=8
_PPTX_SLIDES_PER_PAGE_OUTPUT_TYPE = {1: 1, 2: 2, 3: 3, 4: 8, 6: 4, 9: 9}


def _pptx_to_pdf_via_com(
    input_path: Path,
    output_path: Path,
    slides_per_page: int = 1,
    orientation: str = "landscape",
) -> None:
    """Export PPTX → PDF using the real PowerPoint app via COM (Windows only).

    Preserves 100 % of the original: colors, fonts, shapes, images, logos.
    Raises RuntimeError if PowerPoint is not installed or COM fails.
    """
    try:
        import win32com.client
        import pythoncom
    except ImportError:
        raise RuntimeError("pywin32 is not installed. Run: pip install pywin32")

    src = str(input_path.resolve())
    dst = str(output_path.resolve())
    output_type = _PPTX_SLIDES_PER_PAGE_OUTPUT_TYPE.get(slides_per_page, 1)
    # ppPrintHandoutHorizontalFirst=2 (landscape), ppPrintHandoutVerticalFirst=1 (portrait)
    handout_order = 2 if orientation == "landscape" else 1

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(src, ReadOnly=True, WithWindow=False)
        if slides_per_page == 1:
            # SaveAs is simplest and most reliable for single-slide-per-page
            presentation.SaveAs(dst, 32)  # 32 = ppSaveAsPDF
        else:
            # ExportAsFixedFormat supports handout layouts (multiple slides per page)
            # ppFixedFormatTypePDF=2, ppFixedFormatIntentPrint=1
            presentation.ExportAsFixedFormat(
                dst, 2, 1, False, handout_order, output_type, False, None, 1, "", True, True, True, True, False
            )
    except Exception as exc:
        raise RuntimeError(f"PowerPoint COM export failed: {exc}") from exc
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def _powerpoint_available() -> bool:
    """Return True if PowerPoint COM automation is usable."""
    try:
        import win32com.client, pythoncom  # noqa: F401
        import winreg
        winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                       r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE")
        return True
    except Exception:
        return False


def _libreoffice_available() -> bool:
    """Return True if LibreOffice (soffice) is on PATH."""
    return shutil.which("soffice") is not None or shutil.which("libreoffice") is not None


def _office_to_pdf_via_libreoffice(input_path: Path, output_path: Path) -> None:
    """Convert any Office file → PDF using LibreOffice headless.

    Works on Linux, macOS, and Windows wherever soffice is on PATH.
    Does NOT honour slides_per_page, paper_size, or orientation —
    uses the document's own settings. The _conversion_pipeline QA check
    will fall through to the AI engine if layout is wrong.
    """
    with tempfile.TemporaryDirectory(prefix="cf_libre_") as tmp:
        tmp_path = Path(tmp)
        src = tmp_path / input_path.name
        shutil.copy2(str(input_path), str(src))

        bin_name = "soffice" if shutil.which("soffice") else "libreoffice"
        result = subprocess.run(
            [bin_name, "--headless", "--convert-to", "pdf",
             str(src), "--outdir", str(tmp_path)],
            capture_output=True,
            timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: "
                f"{result.stderr.decode(errors='replace')[:200]}"
            )

        out_pdf = tmp_path / (src.stem + ".pdf")
        if not out_pdf.exists() or out_pdf.stat().st_size == 0:
            raise RuntimeError("LibreOffice produced no output.")

        shutil.copy2(str(out_pdf), str(output_path))


_PPTX_TO_PDF_SLIDE_PROMPT = (
    "You are given the raw text content of one presentation slide.\n"
    "Format it as clean inner HTML for a slide:\n"
    "- The first line (slide title) → <h1>\n"
    "- Lines starting with - or • → group into <ul><li> items\n"
    "- Short ALL-CAPS lines or lines ending with : → <h2>\n"
    "- All other lines → <p>\n"
    "Return ONLY the inner HTML — no <html>, <head>, <body>, or <div> wrapper, no commentary."
)


def _sanitize_latin1(text: str) -> str:
    """Replace common non-Latin-1 chars with ASCII equivalents, then strip the rest.

    fpdf2 built-in fonts (Helvetica, Times, Courier) only cover Latin-1 / cp1252.
    PPTXs frequently contain private-use Wingdings bullets (U+F0B7, U+F0A7),
    curly quotes, em-dashes, and other Unicode that causes a hard crash.
    """
    _MAP = {
        # Curly quotes
        '\u2018': "'", '\u2019': "'",
        '\u201c': '"', '\u201d': '"',
        # Dashes
        '\u2013': '-', '\u2014': '-',
        # Ellipsis
        '\u2026': '...',
        # Common bullets (including private-use Wingdings/Symbol used by PowerPoint)
        '\u2022': '-', '\u25cf': '-', '\u25cb': '-',
        '\u25a0': '-', '\u25a1': '-', '\u25aa': '-', '\u25ab': '-',
        '\uf0b7': '-', '\uf0a7': '-', '\uf0d8': '-', '\uf0fc': '-',
        # Arrows
        '\u2192': '->', '\u2190': '<-', '\u21d2': '=>',
        # Spaces & misc
        '\u00a0': ' ', '\u200b': '',
        # Symbols
        '\u00ae': '(R)', '\u00a9': '(C)', '\u2122': '(TM)',
        '\u00b7': '-',
    }
    for src, dst in _MAP.items():
        text = text.replace(src, dst)
    # Final pass: drop anything still outside cp1252
    return text.encode('cp1252', errors='replace').decode('cp1252')


def _pptx_slide_texts(slide) -> list[tuple[str, bool]]:
    """Return [(text, is_title), ...] for all paragraphs in a slide."""
    from pptx.util import Pt
    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = _sanitize_latin1(para.text.strip())
            if not text:
                continue
            is_title = any(
                (run.font.bold or (run.font.size and run.font.size >= Pt(18)))
                for run in para.runs if run.font
            )
            result.append((text, is_title))
    return result


def _draw_slide_box(pdf, items, x, y, w, h, slide_num, title_fs, body_fs):
    """Render slide text inside a bordered box at (x, y, w, h)."""
    pdf.set_draw_color(210, 210, 210)
    pdf.rect(x, y, w, h)
    # slide badge
    pdf.set_font("Helvetica", "I", max(5, body_fs - 1))
    pdf.set_text_color(170, 170, 170)
    pdf.set_xy(x + w - 22, y + 1.5)
    pdf.cell(20, 3, f"Slide {slide_num}", align="R")
    pdf.set_text_color(0, 0, 0)

    cy = y + 6.5
    bottom = y + h - 2.5
    pad = 3.0
    iw = w - pad * 2

    for text, is_title in items:
        if cy >= bottom:
            break
        if is_title:
            pdf.set_font("Helvetica", "B", title_fs)
            lh = title_fs * 0.45
        else:
            pdf.set_font("Helvetica", "", body_fs)
            lh = body_fs * 0.42
        # Truncate to fit one line
        max_chars = max(10, int(iw / (body_fs * 0.175)))
        line = text if len(text) <= max_chars else text[:max_chars - 2] + ".."
        pdf.set_xy(x + pad, cy)
        pdf.cell(iw, lh, line)
        cy += lh + 0.8


def _pptx_to_pdf_ai(
    input_path: Path, output_path: Path,
    engine: str, paper_size: str, orientation: str, margin_mm: int, slides_per_page: int,
) -> None:
    """AI path: extract slide text → AI formats as HTML → xhtml2pdf renders."""
    from pptx import Presentation
    from xhtml2pdf import pisa
    prs = Presentation(str(input_path))
    system = "You format presentation slide text as clean HTML."
    page_style = _page_css(paper_size, orientation, margin_mm)
    slides_html: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines = [t for t, _ in _pptx_slide_texts(slide)]
        if not lines:
            continue
        inner = _ai_text(
            system,
            f"{_PPTX_TO_PDF_SLIDE_PROMPT}\n\nSlide {slide_num}:\n" + "\n".join(lines),
            engine,
        )
        inner = _sanitize_latin1(inner.strip().lstrip("```html").lstrip("```").rstrip("```").strip())
        slides_html.append(
            '<div style="padding:0 4mm">'
            f'<p style="color:#aaa;font-size:8px;text-align:right;margin:0 0 4px">Slide {slide_num}</p>'
            '<style>h1{color:#2F5496;font-size:20px;border-bottom:2px solid #2F5496;'
            'padding-bottom:5px;margin-bottom:12px}'
            'h2{color:#404040;font-size:13px;margin-top:12px}'
            'ul{margin:6px 0 6px 18px}li{margin:4px 0;font-size:11px}'
            'p{font-size:11px;line-height:1.5;margin:5px 0}</style>'
            + inner + '</div>'
        )
    cols = 2 if slides_per_page > 1 else 1
    rows = max(1, slides_per_page // cols)
    pages_html: list[str] = []
    for i in range(0, max(1, len(slides_html)), slides_per_page):
        group = slides_html[i:i + slides_per_page]
        if cols == 1:
            pages_html.append("".join(group))
        else:
            trs = ""
            for r in range(rows):
                row_cells = slides_html[i + r * cols: i + r * cols + cols]
                if not row_cells:
                    break
                row_html = "".join(
                    f'<td style="width:{100//cols}%;vertical-align:top;'
                    f'border:0.5pt solid #ddd;padding:4px">{s}</td>'
                    for s in row_cells
                )
                while len(row_cells) < cols:
                    row_html += f'<td style="width:{100//cols}%;border:0.5pt solid #ddd"></td>'
                    row_cells.append("")
                trs += f"<tr>{row_html}</tr>"
            pages_html.append(
                f'<table style="width:100%;border-collapse:collapse;height:100%">{trs}</table>'
            )
    sep = '<div style="page-break-after:always"></div>'
    full_html = (
        f'<html><head><meta charset="utf-8"/><style>{page_style}</style></head>'
        f'<body>{sep.join(pages_html)}</body></html>'
    )
    with output_path.open("wb") as f:
        result = pisa.CreatePDF(full_html, dest=f, encoding="utf-8")
    if result.err:
        raise RuntimeError(f"xhtml2pdf rendering failed ({result.err} error(s)).")


def _pptx_to_pdf_basic(
    input_path: Path, output_path: Path,
    paper_size: str, orientation: str, margin_mm: int, slides_per_page: int,
) -> None:
    """Basic path: python-pptx text extraction + fpdf2."""
    from pptx import Presentation
    prs = Presentation(str(input_path))
    pdf, cw, ch = _make_fpdf(paper_size, orientation, margin_mm)
    title_fs = {1: 15, 2: 11, 4: 9, 6: 8}.get(slides_per_page, 15)
    body_fs  = {1: 11, 2:  8, 4: 7, 6: 6}.get(slides_per_page, 11)
    if slides_per_page == 1:
        for slide_num, slide in enumerate(prs.slides, 1):
            pdf.add_page()
            pdf.set_font("Helvetica", "I", 8)
            pdf.set_text_color(150, 150, 150)
            pdf.cell(cw, 5, f"Slide {slide_num}", align="R")
            pdf.ln(6)
            pdf.set_text_color(0, 0, 0)
            for text, is_title in _pptx_slide_texts(slide):
                if is_title:
                    pdf.set_font("Helvetica", "B", title_fs)
                    pdf.multi_cell(cw, title_fs * 0.5, text)
                    pdf.ln(2)
                else:
                    pdf.set_font("Helvetica", "", body_fs)
                    pdf.multi_cell(cw, body_fs * 0.5, text)
    else:
        cols = 2
        rows_per_page = slides_per_page // cols
        gap = 3.0
        box_w = (cw - gap * (cols - 1)) / cols
        box_h = (ch - gap * (rows_per_page - 1)) / rows_per_page
        slides = list(enumerate(prs.slides, 1))
        for page_start in range(0, max(1, len(slides)), slides_per_page):
            pdf.add_page()
            for idx, (slide_num, slide) in enumerate(slides[page_start: page_start + slides_per_page]):
                col, row = idx % cols, idx // cols
                x = margin_mm + col * (box_w + gap)
                y = margin_mm + row * (box_h + gap)
                _draw_slide_box(pdf, _pptx_slide_texts(slide), x, y, box_w, box_h, slide_num, title_fs, body_fs)
    pdf.output(str(output_path))


def pptx_to_pdf(
    input_path: Path,
    output_path: Path,
    engine: str = "powerpoint",
    paper_size: str = "a4",
    orientation: str = "portrait",
    margin_mm: int = 15,
    slides_per_page: int = 1,
) -> None:
    """Convert PowerPoint → PDF.

    Engine chain (tried in order, each QA-checked by Ollama):
      1. PowerPoint COM — full fidelity: colors, images, fonts, shapes, logos (Windows + Office only)
      2. LibreOffice    — high fidelity headless fallback; ignores slides_per_page/orientation (Linux)
      3. AI (ollama/github) — styled HTML per slide via AI
      4. Basic          — plain text extraction via python-pptx + fpdf2
    """
    engines: list[tuple[str, Callable[[Path], None]]] = []

    if engine in ("powerpoint", "auto") and _powerpoint_available():
        engines.append(("powerpoint_com", lambda o: _pptx_to_pdf_via_com(input_path, o, slides_per_page=slides_per_page, orientation=orientation)))

    if engine in ("powerpoint", "auto") and _libreoffice_available():
        engines.append(("libreoffice", lambda o: _office_to_pdf_via_libreoffice(input_path, o)))

    ai_engine = engine if engine in ("ollama", "github", "openai") else "auto"
    if engine not in ("basic",):
        engines.append((
            f"ai_{ai_engine}",
            lambda o: _pptx_to_pdf_ai(input_path, o, ai_engine, paper_size, orientation, margin_mm, slides_per_page),
        ))

    engines.append((
        "basic",
        lambda o: _pptx_to_pdf_basic(input_path, o, paper_size, orientation, margin_mm, slides_per_page),
    ))

    _conversion_pipeline(engines, output_path)


# ---------------------------------------------------------------------------
# Excel → PDF  (tabular layout via openpyxl + fpdf2)
# ---------------------------------------------------------------------------

_EXCEL_TO_PDF_PROMPT = (
    "You are given raw spreadsheet data extracted from an Excel sheet.\n"
    "Convert it into a clean HTML <table> element:\n"
    "- Wrap the header row in <thead><tr><th>...</th></tr></thead>\n"
    "- Wrap all data rows in <tbody><tr><td>...</td></tr></tbody>\n"
    "- Right-align cells that contain only numbers (add style=\"text-align:right\")\n"
    "- If a row looks like a totals/subtotal row (contains words like Total, Sum, Grand Total, Subtotal), "
    "wrap its cells in <strong>\n"
    "Return ONLY valid HTML from <table to </table> — no surrounding tags, no commentary, no code fences."
)


def _excel_to_pdf_ai(
    input_path: Path, output_path: Path,
    engine: str, paper_size: str, orientation: str, margin_mm: int,
) -> None:
    """AI path: extract cell data → AI builds HTML table → xhtml2pdf renders."""
    from openpyxl import load_workbook
    from xhtml2pdf import pisa
    wb = load_workbook(str(input_path), data_only=True)
    system = "You convert spreadsheet data to clean, well-structured HTML tables."
    page_style = _page_css(paper_size, orientation, margin_mm)
    tbl_style = (
        "table{border-collapse:collapse;width:100%;font-family:Arial,sans-serif;font-size:10px}"
        "th{background:#2F5496;color:white;padding:5px 8px;text-align:left}"
        "td{border:1px solid #ccc;padding:4px 8px}"
        "tbody tr:nth-child(even){background:#f0f4fa}"
    )
    sheets_html: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.max_row == 0 or ws.max_column == 0:
            continue
        rows_text = [" | ".join(str(c) if c is not None else "" for c in row)
                     for row in ws.iter_rows(values_only=True)]
        sheet_data = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
        table_html = _ai_text(system, f"{_EXCEL_TO_PDF_PROMPT}\n\n{sheet_data}", engine)
        table_html = table_html.strip().lstrip("```html").lstrip("```").rstrip("```").strip()
        sheets_html.append(
            f'<h2 style="font-family:Arial,sans-serif;color:#2F5496;'
            f'font-size:14px;margin-bottom:8px">{sheet_name}</h2>'
            f'<style>{tbl_style}</style>' + table_html
        )
    sep = '<div style="page-break-after:always;margin-bottom:16px"></div>'
    full_html = (
        f'<html><head><meta charset="utf-8"/><style>{page_style}</style></head>'
        f'<body>{sep.join(sheets_html)}</body></html>'
    )
    with output_path.open("wb") as f:
        result = pisa.CreatePDF(full_html, dest=f, encoding="utf-8")
    if result.err:
        raise RuntimeError(f"xhtml2pdf rendering failed ({result.err} error(s)).")


def _excel_to_pdf_basic(
    input_path: Path, output_path: Path,
    paper_size: str, orientation: str, margin_mm: int, fit_columns: bool,
) -> None:
    """Basic path: openpyxl + fpdf2 fixed-width table."""
    from openpyxl import load_workbook
    wb = load_workbook(str(input_path), data_only=True)
    pdf, cw, _ = _make_fpdf(paper_size, orientation, margin_mm)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.max_row == 0 or ws.max_column == 0:
            continue
        pdf.add_page()
        col_count = ws.max_column
        col_w = cw / col_count if fit_columns else min(cw / col_count, 45)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 7, sheet_name, ln=True)
        pdf.ln(2)
        for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
            pdf.set_font("Helvetica", "B" if row_idx == 0 else "", 7)
            for cell in row:
                value = str(cell) if cell is not None else ""
                max_len = 100 if fit_columns else 22
                if len(value) > max_len:
                    value = value[:max_len - 2] + ".."
                try:
                    pdf.cell(col_w, 5, value, border=1)
                except Exception:
                    pdf.cell(col_w, 5, "", border=1)
            pdf.ln()
    pdf.output(str(output_path))


def excel_to_pdf(
    input_path: Path,
    output_path: Path,
    engine: str = "excel",
    paper_size: str = "a4",
    orientation: str = "landscape",
    margin_mm: int = 10,
    fit_columns: bool = False,
    fit_pages_wide: int = 0,
) -> None:
    """Convert Excel → PDF.

    Engine chain (tried in order, each QA-checked by Ollama):
      1. Excel COM   — full fidelity: colors, charts, merged cells, formatting (Windows + Office only)
      2. LibreOffice — high fidelity headless fallback; ignores paper_size/orientation (Linux)
      3. AI (ollama/github) — HTML table via AI
      4. Basic       — plain fpdf2 grid
    """
    engines: list[tuple[str, Callable[[Path], None]]] = []

    if engine in ("excel", "auto") and _powerpoint_available():
        engines.append(("excel_com", lambda o: _excel_to_pdf_via_com(
            input_path, o,
            paper_size=paper_size,
            orientation=orientation,
            fit_pages_wide=fit_pages_wide,
        )))

    if engine in ("excel", "auto") and _libreoffice_available():
        engines.append(("libreoffice", lambda o: _office_to_pdf_via_libreoffice(input_path, o)))

    ai_engine = engine if engine in ("ollama", "github", "openai") else "auto"
    if engine not in ("basic",):
        engines.append((
            f"ai_{ai_engine}",
            lambda o: _excel_to_pdf_ai(input_path, o, ai_engine, paper_size, orientation, margin_mm),
        ))

    engines.append((
        "basic",
        lambda o: _excel_to_pdf_basic(input_path, o, paper_size, orientation, margin_mm, fit_columns),
    ))

    _conversion_pipeline(engines, output_path)


# ---------------------------------------------------------------------------
# HTML → PDF  (via xhtml2pdf)
# ---------------------------------------------------------------------------

def html_to_pdf(input_path: Path, output_path: Path) -> None:
    def _build(out: Path) -> None:
        try:
            from xhtml2pdf import pisa
        except ImportError:
            raise RuntimeError("xhtml2pdf is not installed. Run: pip install xhtml2pdf")
        html_content = input_path.read_text(encoding="utf-8", errors="replace")
        with out.open("wb") as f:
            result = pisa.CreatePDF(html_content, dest=f, encoding="utf-8")
        if result.err:
            raise RuntimeError(f"HTML to PDF failed with {result.err} error(s).")

    _conversion_pipeline([("xhtml2pdf", _build)], output_path)


# ---------------------------------------------------------------------------
# PDF → SVG  (best-effort vector-preserving export for drawings)
# ---------------------------------------------------------------------------

def _strip_svg_preamble(svg: str) -> str:
    """Remove XML/doctype preamble so the SVG can be embedded in a parent SVG."""
    svg = re.sub(r"^\s*<\?xml[^>]*>\s*", "", svg)
    svg = re.sub(r"^\s*<!DOCTYPE[^>]*>\s*", "", svg, flags=re.IGNORECASE)
    return svg.strip()


def _offset_svg(svg: str, x_pt: float, y_pt: float) -> str:
    """Inject x/y offsets into an SVG root so it can be composed with others."""
    svg = _strip_svg_preamble(svg)
    return re.sub(
        r"<svg\b",
        f'<svg x="{x_pt:.2f}pt" y="{y_pt:.2f}pt"',
        svg,
        count=1,
    )


def _safe_output_stem(value: str) -> str:
    cleaned = "".join(c.lower() if c.isalnum() else "_" for c in value.strip())
    cleaned = "_".join(part for part in cleaned.split("_") if part)
    return cleaned or "converted"


def pdf_page_to_svg(
    input_path: Path,
    output_path: Path,
    page_number: int = 0,
    text_as_path: bool = False,
) -> None:
    """Export a single PDF page to SVG, preserving vectors and text where possible."""
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to export.")
    if page_number < 0 or page_number >= len(doc):
        raise ValueError(f"Page number {page_number} is out of range for a {len(doc)}-page PDF.")

    page = doc[page_number]
    svg = page.get_svg_image(text_as_path=text_as_path)
    output_path.write_text(svg, encoding="utf-8")


def pdf_to_svg(
    input_path: Path,
    output_dir: Path,
    export_mode: str = "single",
    page_number: int = 0,
    text_as_path: bool = False,
    auto_pair_small_pages: bool = True,
) -> list[Path]:
    """Export a PDF to one or more SVG files."""
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to export.")
    output_dir.mkdir(parents=True, exist_ok=True)

    stem = _safe_output_stem(input_path.stem)
    outputs: list[Path] = []

    if export_mode == "single":
        if page_number < 0 or page_number >= len(doc):
            raise ValueError(f"Page number {page_number + 1} is out of range for a {len(doc)}-page PDF.")
        out = output_dir / f"{stem}_page_{page_number + 1:02d}.svg"
        pdf_page_to_svg(input_path, out, page_number=page_number, text_as_path=text_as_path)
        return [out]

    if export_mode == "all":
        for index in range(len(doc)):
            out = output_dir / f"{stem}_page_{index + 1:02d}.svg"
            pdf_page_to_svg(input_path, out, page_number=index, text_as_path=text_as_path)
            outputs.append(out)
        return outputs

    if export_mode != "paired":
        raise ValueError("SVG export mode must be one of: single, all, paired.")

    rendered = [(Path(f"page_{i+1:04d}.png"), float(page.rect.width) / 72.0, float(page.rect.height) / 72.0) for i, page in enumerate(doc)]
    groups = _group_small_format_pages(rendered) if auto_pair_small_pages else [
        {"kind": "single", "indices": [index], "width_in": rendered[index][1], "height_in": rendered[index][2], "gap_in": 0.0}
        for index in range(len(rendered))
    ]

    for group_index, group in enumerate(groups, 1):
        if group["kind"] == "single":
            page_index = group["indices"][0]
            out = output_dir / f"{stem}_page_{page_index + 1:02d}.svg"
            pdf_page_to_svg(input_path, out, page_number=page_index, text_as_path=text_as_path)
            outputs.append(out)
            continue

        first_index, second_index = group["indices"]
        first_page = doc[first_index]
        second_page = doc[second_index]
        gap_pt = float(group["gap_in"]) * 72.0
        first_svg = _offset_svg(first_page.get_svg_image(text_as_path=text_as_path), 0.0, 0.0)
        second_svg = _offset_svg(second_page.get_svg_image(text_as_path=text_as_path), float(first_page.rect.width) + gap_pt, 0.0)
        total_width_pt = float(first_page.rect.width) + float(second_page.rect.width) + gap_pt
        total_height_pt = max(float(first_page.rect.height), float(second_page.rect.height))
        combined = (
            '<?xml version="1.0" encoding="UTF-8" standalone="no"?>\n'
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'xmlns:xlink="http://www.w3.org/1999/xlink" '
            f'width="{total_width_pt:.2f}pt" height="{total_height_pt:.2f}pt" '
            f'viewBox="0 0 {total_width_pt:.2f} {total_height_pt:.2f}">\n'
            f'<rect x="0" y="0" width="{total_width_pt:.2f}" height="{total_height_pt:.2f}" fill="white"/>\n'
            f"{first_svg}\n{second_svg}\n"
            "</svg>\n"
        )
        out = output_dir / f"{stem}_pages_{first_index + 1:02d}_{second_index + 1:02d}.svg"
        out.write_text(combined, encoding="utf-8")
        outputs.append(out)

    return outputs


def pdf_to_svg_plan(
    input_path: Path,
    export_mode: str = "single",
    page_number: int = 0,
    auto_pair_small_pages: bool = True,
) -> dict:
    """Return a lightweight export summary for the SVG tool UI."""
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to export.")

    page_sizes = [
        {
            "page": index + 1,
            "width_in": round(float(page.rect.width) / 72.0, 2),
            "height_in": round(float(page.rect.height) / 72.0, 2),
        }
        for index, page in enumerate(doc)
    ]

    if export_mode == "single":
        if page_number < 0 or page_number >= len(doc):
            raise ValueError(f"Page number {page_number + 1} is out of range for a {len(doc)}-page PDF.")
        selected = page_sizes[page_number]
        return {
            "mode": "single",
            "page_count": len(doc),
            "output_count": 1,
            "download_type": "single",
            "summary": f"Export page {selected['page']} as one SVG.",
            "outputs": [
                {
                    "kind": "single",
                    "pages": [selected["page"]],
                    "label": f"Page {selected['page']}",
                    "size_in": [selected["width_in"], selected["height_in"]],
                }
            ],
        }

    if export_mode == "all":
        return {
            "mode": "all",
            "page_count": len(doc),
            "output_count": len(doc),
            "download_type": "zip" if len(doc) > 1 else "single",
            "summary": f"Export all {len(doc)} pages as separate SVG files.",
            "outputs": [
                {
                    "kind": "single",
                    "pages": [page["page"]],
                    "label": f"Page {page['page']}",
                    "size_in": [page["width_in"], page["height_in"]],
                }
                for page in page_sizes
            ],
        }

    if export_mode != "paired":
        raise ValueError("SVG export mode must be one of: single, all, paired.")

    rendered = [(Path(f"page_{i+1:04d}.png"), float(page.rect.width) / 72.0, float(page.rect.height) / 72.0) for i, page in enumerate(doc)]
    groups = _group_small_format_pages(rendered) if auto_pair_small_pages else [
        {"kind": "single", "indices": [index], "width_in": rendered[index][1], "height_in": rendered[index][2], "gap_in": 0.0}
        for index in range(len(rendered))
    ]
    outputs = []
    for group in groups:
        pages = [index + 1 for index in group["indices"]]
        label = f"Pages {pages[0]}-{pages[-1]}" if len(pages) > 1 else f"Page {pages[0]}"
        outputs.append(
            {
                "kind": group["kind"],
                "pages": pages,
                "label": label,
                "size_in": [round(float(group["width_in"]), 2), round(float(group["height_in"]), 2)],
            }
        )
    return {
        "mode": "paired",
        "page_count": len(doc),
        "output_count": len(outputs),
        "download_type": "zip" if len(outputs) > 1 else "single",
        "summary": f"Export {len(outputs)} SVG file{'s' if len(outputs) != 1 else ''} using small-page pairing where possible.",
        "outputs": outputs,
    }


def pdf_fidelity_plan(
    input_path: Path,
    target: str,
    auto_pair_small_pages: bool = True,
) -> dict:
    """Return a lightweight plan for PDF->Word/PPTX fidelity exports."""
    import fitz  # type: ignore

    if target not in {"word", "pptx"}:
        raise ValueError("Fidelity plan target must be 'word' or 'pptx'.")

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to export.")

    rendered = [(Path(f"page_{i+1:04d}.png"), float(page.rect.width) / 72.0, float(page.rect.height) / 72.0) for i, page in enumerate(doc)]
    groups = _group_small_format_pages(rendered) if auto_pair_small_pages else [
        {"kind": "single", "indices": [index], "width_in": rendered[index][1], "height_in": rendered[index][2], "gap_in": 0.0}
        for index in range(len(rendered))
    ]

    outputs = []
    for group in groups:
        pages = [index + 1 for index in group["indices"]]
        label = f"Pages {pages[0]}-{pages[-1]}" if len(pages) > 1 else f"Page {pages[0]}"
        outputs.append(
            {
                "kind": group["kind"],
                "pages": pages,
                "label": label,
                "size_in": [round(float(group["width_in"]), 2), round(float(group["height_in"]), 2)],
            }
        )

    if target == "word":
        download_type = "single"
        summary = (
            f"Export one Word document with {len(outputs)} page layout block"
            f"{'' if len(outputs) == 1 else 's'}, using small-page pairing where possible."
        )
    else:
        download_type = "single"
        summary = (
            f"Export one PowerPoint deck with {len(outputs)} slide"
            f"{'' if len(outputs) == 1 else 's'}, using small-page pairing where possible."
        )

    return {
        "target": target,
        "page_count": len(doc),
        "output_count": 1,
        "layout_count": len(outputs),
        "download_type": download_type,
        "summary": summary,
        "outputs": outputs,
    }


def pdf_to_excel_plan(input_path: Path) -> dict:
    """Return a lightweight plan for PDF->Excel fidelity exports."""
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to export.")

    outputs = []
    table_pages = 0
    text_pages = 0

    for index, page in enumerate(doc, 1):
        page_width_in = round(float(page.rect.width) / 72.0, 2)
        page_height_in = round(float(page.rect.height) / 72.0, 2)
        table_count = 0
        table_objects = []
        try:
            tables = page.find_tables()
            table_count = len(tables.tables)
            table_objects = list(tables.tables)
        except Exception:
            table_count = 0
            table_objects = []

        page_dict = page.get_text("dict")
        inferred_table_rows = _extract_grid_like_table(page_dict, float(page.rect.width)) if table_count == 0 else []
        inferred_table_count = 1 if inferred_table_rows else 0
        narrative_blocks = 0
        for lines, bbox in _sorted_text_blocks(page_dict, float(page.rect.width)):
            if any(_bbox_overlaps(bbox, tuple(float(v) for v in table.bbox)) for table in table_objects):
                continue
            plain = _block_plain_text(lines)
            if plain:
                narrative_blocks += 1

        effective_table_count = table_count or inferred_table_count

        if effective_table_count:
            table_pages += 1
        if narrative_blocks:
            text_pages += 1

        if effective_table_count and narrative_blocks:
            mode = "snapshot + tables + text"
        elif effective_table_count:
            mode = "snapshot + tables"
        elif narrative_blocks:
            mode = "snapshot + text"
        else:
            mode = "snapshot only"

        outputs.append(
            {
                "label": f"Page {index}",
                "pages": [index],
                "sheet_name": f"Page {index}",
                "size_in": [page_width_in, page_height_in],
                "table_count": effective_table_count,
                "narrative_blocks": narrative_blocks,
                "mode": mode,
            }
        )

    return {
        "target": "excel",
        "page_count": len(doc),
        "output_count": 1,
        "sheet_count": len(outputs),
        "download_type": "single",
        "summary": (
            f"Export one workbook with {len(outputs)} worksheet"
            f"{'' if len(outputs) == 1 else 's'}, preserving page snapshots and adding detected tables/text below them."
        ),
        "table_pages": table_pages,
        "text_pages": text_pages,
        "outputs": outputs,
    }


def pdf_document_profile(input_path: Path) -> dict:
    """Return a cross-tool document profile for routing/export hints."""
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to inspect.")

    small_pages = 0
    drawing_pages = 0
    table_pages = 0
    text_pages = 0
    image_pages = 0
    total_images = 0
    total_drawings = 0
    page_sizes: list[tuple[float, float]] = []

    for page in doc:
        width_in = float(page.rect.width) / 72.0
        height_in = float(page.rect.height) / 72.0
        page_sizes.append((round(width_in, 2), round(height_in, 2)))
        if _is_small_format_page(width_in, height_in):
            small_pages += 1

        page_dict = page.get_text("dict")
        if _page_is_drawing_like(page_dict, float(page.rect.width), float(page.rect.height)):
            drawing_pages += 1

        try:
            tables = page.find_tables()
            if len(tables.tables) > 0:
                table_pages += 1
        except Exception:
            pass

        has_text = any(
            str(span.get("text", "")).strip()
            for block in page_dict.get("blocks", [])
            if block.get("type") == 0
            for line in block.get("lines", [])
            for span in line.get("spans", [])
        )
        if has_text:
            text_pages += 1

        images = page.get_images(full=True)
        total_images += len(images)
        if images:
            image_pages += 1

        try:
            total_drawings += len(page.get_drawings())
        except Exception:
            pass

    unique_sizes = sorted(set(page_sizes))
    primary_size = unique_sizes[0]
    if len(unique_sizes) == 1:
        size_label = f"{primary_size[0]} x {primary_size[1]} in"
    else:
        size_label = f"{primary_size[0]} x {primary_size[1]} in + {len(unique_sizes) - 1} other size set(s)"

    if drawing_pages >= max(1, len(doc) // 2):
        profile = "drawing sheet"
        recommendation = "Prefer SVG or visual-preserving fidelity output."
    elif small_pages == len(doc):
        profile = "small-format layout"
        recommendation = "Pair small front/back pages when appropriate."
    elif table_pages >= max(1, len(doc) // 2):
        profile = "table-heavy report"
        recommendation = "Excel fidelity should preserve snapshots and structured tables well."
    else:
        profile = "regular document"
        recommendation = "Word or PowerPoint fidelity should preserve layout while keeping text editable where possible."

    if profile == "drawing sheet":
        tool_fit = {
            "word": {
                "label": "use with caveats",
                "editability": "low",
                "reason": "Word will mainly preserve the page visually. Technical drawing content is not a strong native Word target.",
            },
            "pptx": {
                "label": "use with caveats",
                "editability": "low",
                "reason": "PowerPoint can preserve the sheet as a slide, but drawing objects are not rebuilt as native slide shapes.",
            },
            "excel": {
                "label": "weak fit",
                "editability": "low",
                "reason": "Excel will mostly become page snapshots with limited extracted content, not a real drawing workspace.",
            },
            "svg": {
                "label": "best fit",
                "editability": "high",
                "reason": "SVG is the strongest target for preserving vector lines, orientation, symbols, and text placement.",
            },
        }
        best_target = "svg"
    elif profile == "small-format layout":
        tool_fit = {
            "word": {
                "label": "use with caveats",
                "editability": "low-medium",
                "reason": "Word stays visual-first on small-format layouts to avoid fragile compatibility-mode overlays.",
            },
            "pptx": {
                "label": "good fit",
                "editability": "medium",
                "reason": "PowerPoint handles fixed layouts and front/back pairing better than Word while keeping placement easier to edit.",
            },
            "excel": {
                "label": "weak fit",
                "editability": "low",
                "reason": "Excel is not a natural target for card-style layouts unless you only need snapshots plus extracted text.",
            },
            "svg": {
                "label": "best fit",
                "editability": "high",
                "reason": "SVG keeps the compact design, placement, and vectors closest to the original layout.",
            },
        }
        best_target = "svg"
    elif profile == "table-heavy report":
        tool_fit = {
            "word": {
                "label": "good fit",
                "editability": "medium",
                "reason": "Word works well for readable report output, but dense tables may still be better in Excel.",
            },
            "pptx": {
                "label": "use with caveats",
                "editability": "low-medium",
                "reason": "PowerPoint is useful page-by-page, but large reports with tables are not usually best consumed as slides.",
            },
            "excel": {
                "label": "best fit",
                "editability": "high",
                "reason": "Excel is the strongest target when the PDF contains many tables, schedules, or structured rows.",
            },
            "svg": {
                "label": "use with caveats",
                "editability": "low",
                "reason": "SVG preserves appearance well, but it is usually less practical than Excel for structured business data.",
            },
        }
        best_target = "excel"
    else:
        tool_fit = {
            "word": {
                "label": "best fit",
                "editability": "medium",
                "reason": "Word is usually the best Office target for regular reports and narrative documents.",
            },
            "pptx": {
                "label": "good fit",
                "editability": "medium",
                "reason": "PowerPoint works well when each PDF page should become a slide, but it is less natural than Word for reading-heavy documents.",
            },
            "excel": {
                "label": "use with caveats",
                "editability": "medium",
                "reason": "Excel is helpful when the document contains useful tables, but freeform pages stay snapshot-led.",
            },
            "svg": {
                "label": "use with caveats",
                "editability": "low-medium",
                "reason": "SVG preserves appearance, but it is usually less practical than Word for everyday document editing.",
            },
        }
        best_target = "word"

    return {
        "profile": profile,
        "recommendation": recommendation,
        "best_target": best_target,
        "tool_fit": tool_fit,
        "page_count": len(doc),
        "size_label": size_label,
        "small_pages": small_pages,
        "drawing_pages": drawing_pages,
        "table_pages": table_pages,
        "text_pages": text_pages,
        "image_pages": image_pages,
        "total_images": total_images,
        "total_drawings": total_drawings,
    }


# ---------------------------------------------------------------------------
# PDF info helper
# ---------------------------------------------------------------------------

def pdf_info(input_path: Path) -> dict:
    from pypdf import PdfReader
    reader = PdfReader(str(input_path))
    meta = reader.metadata or {}
    return {
        "pages": len(reader.pages),
        "title": meta.get("/Title", ""),
        "author": meta.get("/Author", ""),
        "encrypted": reader.is_encrypted,
    }


# ---------------------------------------------------------------------------
# Redact PDF  (permanently black out phrases / words)
# ---------------------------------------------------------------------------

def redact_pdf(input_path: Path, output_path: Path, phrases: list[str]) -> int:
    """
    Search for each phrase in every page and apply permanent black redaction boxes.
    Returns the total number of redactions applied.
    """
    import fitz  # type: ignore

    clean_phrases = [p.strip() for p in phrases if p.strip()]
    if not clean_phrases:
        raise ValueError("Provide at least one word or phrase to redact.")

    doc = fitz.open(str(input_path))
    total = 0
    for page in doc:
        for phrase in clean_phrases:
            hits = page.search_for(phrase)
            for rect in hits:
                page.add_redact_annot(rect, fill=(0, 0, 0))
                total += 1
        if total:
            page.apply_redactions()
    doc.save(str(output_path), garbage=4, deflate=True)
    return total


# ---------------------------------------------------------------------------
# Compare PDF  (side-by-side HTML diff of extracted text)
# ---------------------------------------------------------------------------

def compare_pdfs(path_a: Path, path_b: Path, output_path: Path) -> None:
    """
    Extract text from both PDFs, produce a side-by-side HTML diff, and write to output_path.
    """
    import difflib
    from pypdf import PdfReader

    def _extract(p: Path) -> list[str]:
        reader = PdfReader(str(p))
        lines: list[str] = []
        for page in reader.pages:
            text = page.extract_text() or ""
            lines.extend(text.splitlines())
        return lines

    lines_a = _extract(path_a)
    lines_b = _extract(path_b)

    differ = difflib.HtmlDiff(wrapcolumn=80)
    table = differ.make_table(
        lines_a, lines_b,
        fromdesc=path_a.name,
        todesc=path_b.name,
        context=True,
        numlines=3,
    )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>PDF Comparison — {path_a.name} vs {path_b.name}</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
  *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: 'JetBrains Mono', 'Fira Code', 'Courier New', monospace;
    font-size: 13px;
    line-height: 1.6;
    background: #0e0e10;
    color: #f0ece4;
    padding: 0;
    min-height: 100vh;
  }}
  .report-header {{
    background: #18181c;
    border-bottom: 1px solid rgba(255,255,255,0.07);
    padding: 14px 24px;
    display: flex;
    align-items: center;
    gap: 16px;
  }}
  .report-logo {{
    display: flex;
    align-items: center;
    gap: 10px;
    text-decoration: none;
  }}
  .report-logo svg {{ flex-shrink: 0; }}
  .report-logo-copy {{
    font-family: 'Space Grotesk', system-ui, sans-serif;
    line-height: 1.2;
  }}
  .report-logo-copy strong {{ font-size: 14px; font-weight: 700; color: #f0ece4; letter-spacing: -.01em; }}
  .report-logo-copy strong span {{ color: #e8a838; }}
  .report-logo-copy small {{ font-size: 11px; color: #8a8690; display: block; }}
  .report-title-area {{
    margin-left: auto;
    text-align: right;
  }}
  .report-title-area h1 {{
    font-family: 'Space Grotesk', system-ui, sans-serif;
    font-size: 13px;
    font-weight: 600;
    color: #f0ece4;
  }}
  .report-title-area .filenames {{
    font-size: 11px;
    color: #8a8690;
    margin-top: 1px;
    word-break: break-all;
  }}
  .report-body {{ padding: 20px 24px; }}
  .legend {{
    display: flex;
    align-items: center;
    gap: 20px;
    margin-bottom: 16px;
    font-family: 'Space Grotesk', system-ui, sans-serif;
    font-size: 12px;
    color: #8a8690;
    flex-wrap: wrap;
  }}
  .legend-item {{ display: flex; align-items: center; gap: 6px; }}
  .legend-swatch {{
    width: 10px; height: 10px; border-radius: 2px; flex-shrink: 0;
  }}
  .legend-swatch--add {{ background: #34d399; }}
  .legend-swatch--del {{ background: #f87171; }}
  .legend-swatch--chg {{ background: #e8a838; }}
  .table-wrap {{
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 8px;
    overflow: hidden;
    overflow-x: auto;
  }}
  table.diff {{ width: 100%; border-collapse: collapse; }}
  table.diff td {{ padding: 2px 8px; vertical-align: top; white-space: pre-wrap; word-break: break-word; }}
  table.diff th {{
    background: #18181c;
    color: #8a8690;
    padding: 6px 8px;
    text-align: left;
    border-bottom: 1px solid rgba(255,255,255,0.07);
    font-family: 'Space Grotesk', system-ui, sans-serif;
    font-size: 12px;
    font-weight: 600;
    letter-spacing: .03em;
    text-transform: uppercase;
  }}
  .diff_add {{ background: rgba(52,211,153,0.10); }}
  .diff_chg {{ background: rgba(232,168,56,0.10); }}
  .diff_sub {{ background: rgba(248,113,113,0.10); }}
  .diff_next {{ background: rgba(30,42,58,0.7); color: #60a5fa; }}
  span.diff_add {{ background: rgba(52,211,153,0.28); color: #34d399; border-radius: 2px; padding: 0 1px; }}
  span.diff_chg {{ background: rgba(232,168,56,0.28); color: #e8a838; border-radius: 2px; padding: 0 1px; }}
  span.diff_sub {{ background: rgba(248,113,113,0.28); color: #f87171; border-radius: 2px; padding: 0 1px; }}
  td.diff_header {{ color: #55525c; padding: 2px 6px; font-size: 11px; }}
  colgroup col.diff_header {{ width: 40px; }}
  tr:hover td:not(.diff_next) {{ background-color: rgba(255,255,255,0.025); }}
  .report-footer {{
    margin-top: 24px;
    padding: 16px 24px;
    border-top: 1px solid rgba(255,255,255,0.07);
    font-family: 'Space Grotesk', system-ui, sans-serif;
    font-size: 11px;
    color: #55525c;
  }}
  @media (max-width: 600px) {{
    .report-header {{ flex-direction: column; align-items: flex-start; gap: 8px; }}
    .report-title-area {{ margin-left: 0; text-align: left; }}
    .report-body {{ padding: 16px; }}
  }}
</style>
</head>
<body>
<header class="report-header">
  <div class="report-logo">
    <svg viewBox="0 0 20 20" fill="none" xmlns="http://www.w3.org/2000/svg" width="28" height="28" aria-hidden="true">
      <circle cx="10" cy="10" r="8.25" stroke="rgba(232,168,56,.3)" stroke-width="1.2"/>
      <path d="M4.75 10h10.5M11.75 6.75L15 10l-3.25 3.25M8.25 13.25L5 10l3.25-3.25" stroke="#e8a838" stroke-width="1.9" stroke-linecap="round" stroke-linejoin="round"/>
    </svg>
    <div class="report-logo-copy">
      <strong>Convert<span>Flow</span></strong>
      <small>PDF Comparison Report</small>
    </div>
  </div>
  <div class="report-title-area">
    <h1>Side-by-side diff</h1>
    <div class="filenames">{path_a.name} &nbsp;vs&nbsp; {path_b.name}</div>
  </div>
</header>
<div class="report-body">
  <div class="legend">
    <span class="legend-item"><span class="legend-swatch legend-swatch--add"></span> Additions</span>
    <span class="legend-item"><span class="legend-swatch legend-swatch--del"></span> Deletions</span>
    <span class="legend-item"><span class="legend-swatch legend-swatch--chg"></span> Changes</span>
  </div>
  <div class="table-wrap">
{table}
  </div>
</div>
<footer class="report-footer">
  Generated by ConvertFlow &nbsp;&mdash;&nbsp; Core processing runs on your machine
</footer>
</body>
</html>"""

    output_path.write_text(html, encoding="utf-8")


# ---------------------------------------------------------------------------
# Organize PDF  (render thumbnails + reorder pages)
# ---------------------------------------------------------------------------

def render_pdf_thumbnails(input_path: Path, output_dir: Path, dpi: int = 72) -> list[Path]:
    """Render every PDF page as a small PNG thumbnail."""
    import fitz  # type: ignore
    doc = fitz.open(str(input_path))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    paths: list[Path] = []
    for i, page in enumerate(doc, 1):
        out = output_dir / f"thumb_{i:04d}.png"
        pix = page.get_pixmap(matrix=mat)
        pix.save(str(out))
        paths.append(out)
    return paths


def render_pdf_edit_previews(input_path: Path, output_dir: Path, dpi: int = 96) -> list[dict]:
    """Render edit previews and return page metadata for the browser editor."""
    import fitz  # type: ignore

    def _edit_text_blocks(page) -> list[dict]:
        page_width = float(page.rect.width) or 1.0
        page_height = float(page.rect.height) or 1.0
        page_dict = page.get_text("dict")
        blocks: list[dict] = []
        block_id = 1
        for block in page_dict.get("blocks", []):
            if block.get("type") != 0:
                continue
            bbox = block.get("bbox") or ()
            if len(bbox) != 4:
                continue
            x0, y0, x1, y1 = (float(v) for v in bbox)
            width = max(0.0, x1 - x0)
            height = max(0.0, y1 - y0)
            if width < 8 or height < 6:
                continue
            lines = block.get("lines", [])
            text_lines: list[str] = []
            span_texts: list[str] = []
            font_sizes: list[float] = []
            color_values: list[int] = []
            font_names: list[str] = []
            bold_hits = 0
            italic_hits = 0
            total_spans = 0
            for line in lines:
                line_parts: list[str] = []
                for span in line.get("spans", []):
                    text = str(span.get("text", "")).strip()
                    if not text:
                        continue
                    line_parts.append(text)
                    span_texts.append(text)
                    total_spans += 1
                    try:
                        font_sizes.append(float(span.get("size") or 0))
                    except Exception:
                        pass
                    try:
                        color_values.append(int(span.get("color") or 0))
                    except Exception:
                        pass
                    font_name = str(span.get("font") or "")
                    if font_name:
                        font_names.append(font_name)
                        upper_name = font_name.upper()
                        if "BOLD" in upper_name or upper_name.endswith("BD"):
                            bold_hits += 1
                        if "ITALIC" in upper_name or "OBLIQUE" in upper_name or upper_name.endswith("IT"):
                            italic_hits += 1
                if line_parts:
                    text_lines.append(" ".join(line_parts))
            block_text = "\n".join(text_lines).strip()
            if not block_text:
                continue
            if len(block_text) > 800:
                block_text = block_text[:800].rstrip() + "..."
            font_size = round(sum(font_sizes) / len(font_sizes), 2) if font_sizes else 12.0
            color_hex = "#111111"
            if color_values:
                rgb = color_values[0] & 0xFFFFFF
                color_hex = f"#{rgb:06x}"
            font_family = "helv"
            if font_names:
                name = font_names[0].lower()
                if "cour" in name:
                    font_family = "cour"
                elif "times" in name or "serif" in name or "roman" in name:
                    font_family = "tiro"
            blocks.append(
                {
                    "id": f"tb_{page.number + 1}_{block_id}",
                    "text": block_text,
                    "x": round(x0 / page_width, 6),
                    "y": round(y0 / page_height, 6),
                    "w": round(width / page_width, 6),
                    "h": round(height / page_height, 6),
                    "font_size": font_size,
                    "font_family": font_family,
                    "bold": bold_hits > 0 and bold_hits >= max(1, total_spans // 2),
                    "italic": italic_hits > 0 and italic_hits >= max(1, total_spans // 2),
                    "color": color_hex,
                }
            )
            block_id += 1
        return blocks

    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(input_path))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pages: list[dict] = []
    try:
        for i, page in enumerate(doc, 1):
            out = output_dir / f"edit_page_{i:04d}.png"
            pix = page.get_pixmap(matrix=mat, alpha=False)
            pix.save(str(out))
            pages.append(
                {
                    "index": i - 1,
                    "page_number": i,
                    "thumb_path": out,
                    "page_width_pt": float(page.rect.width),
                    "page_height_pt": float(page.rect.height),
                    "thumb_width_px": pix.width,
                    "thumb_height_px": pix.height,
                    "text_blocks": _edit_text_blocks(page),
                }
            )
    finally:
        doc.close()
    return pages


def _hex_to_rgb_fraction(value: str) -> tuple[float, float, float]:
    text = str(value or "").strip().lstrip("#")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if len(text) != 6 or any(ch not in "0123456789abcdefABCDEF" for ch in text):
        return (0.0, 0.0, 0.0)
    return tuple(int(text[i:i + 2], 16) / 255.0 for i in (0, 2, 4))


def edit_pdf(input_path: Path, output_path: Path, operations: Sequence[dict]) -> None:
    """Apply simple overlay edits to a PDF using normalized page coordinates."""
    import base64
    import fitz  # type: ignore

    doc = fitz.open(str(input_path))
    if len(doc) == 0:
        raise ValueError("The PDF has no pages to edit.")

    for op_index, operation in enumerate(operations, 1):
        page_index = int(operation.get("page", -1))
        if page_index < 0 or page_index >= len(doc):
            raise ValueError(f"Edit operation {op_index} references an invalid page index.")
        op_type = str(operation.get("type", "")).strip().lower()
        if op_type not in {"text", "replace_text", "whiteout", "image", "highlight", "underline", "strikethrough",
                           "rect_shape", "ellipse", "line", "stamp", "symbol", "link"}:
            raise ValueError(f"Edit operation {op_index} has an unsupported type: {op_type!r}")

        page = doc[page_index]
        page_width = float(page.rect.width)
        page_height = float(page.rect.height)

        def _norm(name: str, default: float = 0.0) -> float:
            try:
                value = float(operation.get(name, default))
            except Exception as exc:
                raise ValueError(f"Edit operation {op_index} has an invalid {name!r} value.") from exc
            if value < 0 or value > 1.2:
                raise ValueError(f"Edit operation {op_index} has out-of-range normalized {name!r}.")
            return value

        x = _norm("x") * page_width
        y = _norm("y") * page_height
        w = max(1.0, _norm("w", 0.12) * page_width)
        h = max(1.0, _norm("h", 0.04) * page_height)
        rect = fitz.Rect(x, y, min(page_width, x + w), min(page_height, y + h))

        def _opacity(default: float = 1.0) -> float:
            try:
                val = float(operation.get("opacity", default))
            except Exception:
                val = default
            return max(0.0, min(1.0, val))

        if op_type == "whiteout":
            page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1), width=0)
            continue

        if op_type == "highlight":
            fill_color = _hex_to_rgb_fraction(str(operation.get("color", "#ffeb3b")))
            opacity = _opacity(0.5)
            page.draw_rect(rect, color=None, fill=fill_color, fill_opacity=opacity, width=0)
            continue

        if op_type == "underline":
            line_color = _hex_to_rgb_fraction(str(operation.get("color", "#e53935")))
            opacity = _opacity(1.0)
            line_h = max(1.0, (rect.y1 - rect.y0) * 0.15)
            line_rect = fitz.Rect(rect.x0, rect.y1 - line_h, rect.x1, rect.y1)
            page.draw_rect(line_rect, color=None, fill=line_color, fill_opacity=opacity, width=0)
            continue

        if op_type == "strikethrough":
            line_color = _hex_to_rgb_fraction(str(operation.get("color", "#e53935")))
            opacity = _opacity(1.0)
            line_h = max(1.0, (rect.y1 - rect.y0) * 0.15)
            mid_y = (rect.y0 + rect.y1) / 2
            line_rect = fitz.Rect(rect.x0, mid_y - line_h / 2, rect.x1, mid_y + line_h / 2)
            page.draw_rect(line_rect, color=None, fill=line_color, fill_opacity=opacity, width=0)
            continue

        if op_type == "rect_shape":
            fill = _hex_to_rgb_fraction(str(operation.get("fill_color", "#4a90d9")))
            stroke = _hex_to_rgb_fraction(str(operation.get("stroke_color", "#1a5fa8")))
            bw = max(0.0, float(operation.get("border_width", 2)))
            opacity = _opacity(0.8)
            page.draw_rect(rect, color=stroke if bw > 0 else None, fill=fill,
                           fill_opacity=opacity, width=bw)
            continue

        if op_type == "ellipse":
            fill = _hex_to_rgb_fraction(str(operation.get("fill_color", "#4a90d9")))
            stroke = _hex_to_rgb_fraction(str(operation.get("stroke_color", "#1a5fa8")))
            bw = max(0.0, float(operation.get("border_width", 2)))
            opacity = _opacity(0.8)
            page.draw_oval(rect, color=stroke if bw > 0 else None, fill=fill,
                           fill_opacity=opacity, width=bw)
            continue

        if op_type == "line":
            stroke = _hex_to_rgb_fraction(str(operation.get("stroke_color", "#1a5fa8")))
            bw = max(1.0, float(operation.get("border_width", 3)))
            opacity = _opacity(1.0)
            arrowhead = str(operation.get("arrowhead", "end")).lower()
            p1 = fitz.Point(rect.x0, (rect.y0 + rect.y1) / 2)
            p2 = fitz.Point(rect.x1, (rect.y0 + rect.y1) / 2)
            annot = page.add_line_annot(p1, p2)
            annot.set_colors(stroke=stroke)
            annot.set_border(width=bw)
            if arrowhead in ("end", "both"):
                annot.set_line_ends(fitz.PDF_ANNOT_LE_NONE if arrowhead == "end" else fitz.PDF_ANNOT_LE_OPEN_ARROW,
                                    fitz.PDF_ANNOT_LE_OPEN_ARROW)
            annot.set_opacity(opacity)
            annot.update()
            continue

        if op_type == "stamp":
            stamp_text = str(operation.get("stamp_preset", "DRAFT")).strip().upper() or "DRAFT"
            color = _hex_to_rgb_fraction(str(operation.get("color", "#e53935")))
            opacity = _opacity(1.0)
            font_size = max(8.0, h * 0.55)
            page.draw_rect(rect, color=color, fill=None, width=max(1.5, h * 0.04))
            inserted = page.insert_textbox(
                rect, stamp_text,
                fontsize=font_size, fontname="helv",
                color=color, align=1,
            )
            if inserted < 0:
                page.insert_textbox(rect, stamp_text, fontsize=max(6.0, font_size * 0.7),
                                    fontname="helv", color=color, align=1)
            continue

        if op_type == "symbol":
            sym = str(operation.get("symbol", "✓")).strip() or "✓"
            color = _hex_to_rgb_fraction(str(operation.get("color", "#111111")))
            font_size = max(8.0, min(w, h) * 0.75)
            page.insert_textbox(rect, sym, fontsize=font_size, fontname="helv",
                                color=color, align=1)
            continue

        if op_type == "link":
            link_type  = str(operation.get("link_type", "url")).lower()
            link_dest  = str(operation.get("link_dest", "")).strip()
            link_label = str(operation.get("link_label", "")).strip()
            link_color = _hex_to_rgb_fraction(str(operation.get("link_color", "#1a5fa8")))
            link_style = str(operation.get("link_style", "underline")).lower()

            # Build the fitz link dict
            if link_type == "page":
                try:
                    target_page = max(0, int(link_dest) - 1)
                except ValueError:
                    target_page = 0
                target_page = min(target_page, len(doc) - 1)
                link_dict = {
                    "kind": fitz.LINK_GOTO,
                    "page": target_page,
                    "from": rect,
                    "to": fitz.Point(0, 0),
                }
            else:
                # url / email / phone — all become URI links
                uri = link_dest
                if not uri:
                    continue
                link_dict = {
                    "kind": fitz.LINK_URI,
                    "uri": uri,
                    "from": rect,
                }

            page.insert_link(link_dict)

            # Optionally draw a visible decoration (underline or box)
            if link_style == "underline":
                line_h = max(1.0, (rect.y1 - rect.y0) * 0.08)
                underline_rect = fitz.Rect(rect.x0, rect.y1 - line_h, rect.x1, rect.y1)
                page.draw_rect(underline_rect, color=None, fill=link_color, width=0)
            elif link_style == "box":
                page.draw_rect(rect, color=link_color, fill=None, width=max(0.8, h * 0.03))

            # Optionally draw label text
            if link_label:
                font_size = max(6.0, h * 0.55)
                page.insert_textbox(rect, link_label, fontsize=font_size,
                                    fontname="helv", color=link_color, align=0)
            continue

        if op_type == "image":
            image_data = str(operation.get("image_data", "")).strip()
            if not image_data:
                raise ValueError(f"Edit operation {op_index} is missing image data.")
            if "," in image_data and image_data.lower().startswith("data:image/"):
                image_data = image_data.split(",", 1)[1]
            try:
                image_bytes = base64.b64decode(image_data, validate=True)
            except Exception as exc:
                raise ValueError(f"Edit operation {op_index} has invalid image data.") from exc
            if not image_bytes:
                raise ValueError(f"Edit operation {op_index} image data is empty.")
            page.insert_image(rect, stream=image_bytes, keep_proportion=False, overlay=True)
            continue

        text = str(operation.get("text", "")).strip()
        if not text:
            raise ValueError(f"Edit operation {op_index} is missing text content.")
        font_size = operation.get("font_size", 14)
        try:
            font_size = float(font_size)
        except Exception as exc:
            raise ValueError(f"Edit operation {op_index} has an invalid font size.") from exc
        if font_size <= 0:
            raise ValueError(f"Edit operation {op_index} must use a positive font size.")

        # Resolve font name: pymupdf built-in base14 names encode bold/italic as suffixes.
        # helv=Helvetica, tiro=Times, cour=Courier. Bold adds 'o', italic adds 'i', both adds 'bi'.
        base_font = str(operation.get("font_family", "helv")).strip().lower()
        if base_font not in {"helv", "tiro", "cour"}:
            base_font = "helv"
        is_bold   = bool(operation.get("bold", False))
        is_italic = bool(operation.get("italic", False))
        font_suffix_map = {
            ("helv", False, False): "helv",
            ("helv", True,  False): "hebo",
            ("helv", False, True):  "heit",
            ("helv", True,  True):  "hebi",
            ("tiro", False, False): "tiro",
            ("tiro", True,  False): "tibi",
            ("tiro", False, True):  "tii",
            ("tiro", True,  True):  "tibi",
            ("cour", False, False): "cour",
            ("cour", True,  False): "cobo",
            ("cour", False, True):  "coit",
            ("cour", True,  True):  "cobi",
        }
        fontname = font_suffix_map.get((base_font, is_bold, is_italic), base_font)

        color = _hex_to_rgb_fraction(str(operation.get("color", "#111111")))
        align_name = str(operation.get("align", "left")).lower()
        align_map = {"left": 0, "center": 1, "right": 2, "justify": 3}
        align = align_map.get(align_name, 0)
        if op_type == "replace_text":
            # First cover the original block, then write the new text back into the same box.
            page.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1), width=0)
        inserted = page.insert_textbox(
            rect,
            text,
            fontsize=font_size,
            fontname=fontname,
            color=color,
            align=align,
        )
        if inserted < 0:
            fitted = False
            retry_size = font_size
            while retry_size > 6:
                retry_size = max(6.0, round(retry_size * 0.9, 2))
                inserted = page.insert_textbox(
                    rect,
                    text,
                    fontsize=retry_size,
                    fontname=fontname,
                    color=color,
                    align=align,
                )
                if inserted >= 0:
                    fitted = True
                    break
                if retry_size <= 6.0:
                    break
            if not fitted:
                raise ValueError(f"Edit operation {op_index} text box is too small for its content.")

    doc.save(str(output_path), garbage=4, deflate=True)


def organize_pdf(input_path: Path, output_path: Path, page_order: list[int]) -> None:
    """Reorder PDF pages. page_order is 0-indexed list of page indices."""
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    for i in page_order:
        if i < 0 or i >= total:
            raise ValueError(f"Page index {i} is out of range for a {total}-page PDF.")
    writer = PdfWriter()
    for i in page_order:
        writer.add_page(reader.pages[i])
    with output_path.open("wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# Sign PDF  (embed a signature image at a chosen position)
# ---------------------------------------------------------------------------

def sign_pdf(
    input_path: Path,
    output_path: Path,
    signature_png: bytes,
    position: str = "bottom-right",
    page_number: int = -1,
    sig_width_mm: float = 60,
) -> None:
    """
    Embed a signature PNG into the PDF at the given position.
    page_number: 0-indexed page; -1 = last page.
    """
    import fitz  # type: ignore
    import io

    doc = fitz.open(str(input_path))
    total = len(doc)
    if page_number == -1:
        page_number = total - 1
    if page_number < 0 or page_number >= total:
        raise ValueError(f"Page number {page_number} is out of range for a {total}-page PDF.")

    page = doc[page_number]
    pw = page.rect.width
    ph = page.rect.height

    MM_TO_PT = 72 / 25.4
    sig_w = sig_width_mm * MM_TO_PT

    # Detect aspect ratio from the PNG bytes
    try:
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(signature_png))
        iw, ih = img.size
        aspect = ih / iw if iw else 0.3
    except Exception:
        aspect = 0.3  # fallback ratio

    sig_h = sig_w * aspect
    margin = 20.0  # pt

    pos_map = {
        "top-left":     fitz.Rect(margin, margin, margin + sig_w, margin + sig_h),
        "top-right":    fitz.Rect(pw - margin - sig_w, margin, pw - margin, margin + sig_h),
        "center":       fitz.Rect((pw - sig_w) / 2, (ph - sig_h) / 2, (pw + sig_w) / 2, (ph + sig_h) / 2),
        "bottom-left":  fitz.Rect(margin, ph - margin - sig_h, margin + sig_w, ph - margin),
        "bottom-right": fitz.Rect(pw - margin - sig_w, ph - margin - sig_h, pw - margin, ph - margin),
    }
    rect = pos_map.get(position, pos_map["bottom-right"])
    page.insert_image(rect, stream=signature_png)
    doc.save(str(output_path), garbage=4, deflate=True)
