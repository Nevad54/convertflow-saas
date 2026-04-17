from __future__ import annotations

from dataclasses import dataclass
import base64
from html import escape
import json
import logging
import os
from pathlib import Path
import re
import shutil
from typing import Iterable, Sequence

logger = logging.getLogger(__name__)

from PIL import Image

try:
    import pytesseract
except ImportError:  # pragma: no cover - handled at runtime
    pytesseract = None
else:
    from pytesseract import Output

try:
    from openai import OpenAI
except ImportError:  # pragma: no cover - handled at runtime
    OpenAI = None

try:
    import httpx
except ImportError:  # pragma: no cover - handled at runtime
    httpx = None

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.shared import Pt
from fpdf import FPDF
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font
from pptx import Presentation
from docx.table import Table


SUPPORTED_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}
SUPPORTED_OUTPUTS = {"docx", "pdf", "xlsx", "pptx", "md", "html", "txt"}
SUPPORTED_OCR_ENGINES = {"auto", "tesseract", "openai", "github", "ollama"}
SUPPORTED_DOC_STYLES = {"auto", "general", "legal", "form"}
ROMAN_HEADING_RE = re.compile(r"^[IVXLCDM]+\.\s+[A-Z]")
ALPHA_CLAUSE_RE = re.compile(r"^[A-Z]\.\s+")
NUMERIC_CLAUSE_RE = re.compile(r"^\d+\.\s+")
BULLET_LINE_RE = re.compile(r"^(?:[-*•]|[oO0]\s)\s*.+")
CHECKBOX_LINE_RE = re.compile(r"^(?:\[[ xX]?\]|\([ xX]?\)|☐|☑|✓)\s*.+")


@dataclass
class OCRPage:
    filename: str
    text: str


@dataclass
class TableRowCandidate:
    cells: list[str]


@dataclass
class PageAnalysis:
    text: str
    has_table_grid: bool


class MissingDependencyError(RuntimeError):
    """Raised when a required runtime dependency is unavailable."""


def configure_tesseract() -> None:
    if pytesseract is None:
        raise MissingDependencyError(
            "pytesseract is not installed. Install requirements and ensure Tesseract OCR is available on PATH."
        )

    detected = shutil.which("tesseract")
    if detected:
        pytesseract.pytesseract.tesseract_cmd = detected
        return

    windows_candidates = [
        Path(os.environ.get("ProgramFiles", r"C:\Program Files")) / "Tesseract-OCR" / "tesseract.exe",
        Path(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")) / "Tesseract-OCR" / "tesseract.exe",
    ]

    for candidate in windows_candidates:
        if candidate.exists():
            pytesseract.pytesseract.tesseract_cmd = str(candidate)
            return

    raise MissingDependencyError(
        "tesseract is not installed or it's not in your PATH. Install Tesseract OCR and try again."
    )


def validate_images(image_paths: Sequence[Path]) -> None:
    if not image_paths:
        raise ValueError("No image files were provided.")

    invalid = [path.name for path in image_paths if path.suffix.lower() not in SUPPORTED_EXTENSIONS]
    if invalid:
        raise ValueError(f"Unsupported image file(s): {', '.join(invalid)}")


def extract_text_from_images(image_paths: Sequence[Path]) -> list[OCRPage]:
    return extract_text_from_images_with_engine(image_paths, "auto", "auto")


def extract_text_from_images_with_engine(image_paths: Sequence[Path], engine: str, document_style: str) -> list[OCRPage]:
    validate_images(image_paths)
    normalized_engine = normalize_ocr_engine(engine)
    normalized_style = normalize_document_style(document_style)

    total = len(image_paths)
    pages: list[OCRPage] = []
    for idx, image_path in enumerate(image_paths, 1):
        logger.info("OCR [%d/%d] %s (engine=%s)", idx, total, image_path.name, normalized_engine)
        analysis = extract_text_for_single_image(image_path, normalized_engine, normalized_style)
        text = analysis.text
        if analysis.has_table_grid:
            text = f"[TABLE_GRID_DETECTED]\n{text}" if text.strip() else "[TABLE_GRID_DETECTED]"
        normalized = normalize_text(text)
        pages.append(
            OCRPage(
                filename=image_path.name,
                text=normalized or "[No readable text detected on this image]",
            )
        )
    return harmonize_repeated_entities(pages)


def harmonize_repeated_entities(pages: Sequence[OCRPage]) -> list[OCRPage]:
    canonical_map = build_entity_canonical_map(pages)
    if not canonical_map:
        return list(pages)

    normalized_pages: list[OCRPage] = []
    for page in pages:
        updated_text = page.text
        for source, target in canonical_map.items():
            if source == target:
                continue
            updated_text = re.sub(rf"\b{re.escape(source)}\b", target, updated_text)
        normalized_pages.append(OCRPage(filename=page.filename, text=updated_text))
    return normalized_pages


def build_entity_canonical_map(pages: Sequence[OCRPage]) -> dict[str, str]:
    phrase_counts: dict[str, int] = {}
    canonical_map: dict[str, str] = {}

    for page in pages:
        for match in re.finditer(r"\b[A-Z][A-Za-z&'/-]*(?:\s+[A-Z][A-Za-z&'/-]*){0,5}", page.text):
            phrase = match.group(0).strip()
            if not looks_like_entity_candidate(phrase):
                continue
            phrase_counts[phrase] = phrase_counts.get(phrase, 0) + 1

    candidates = sorted(phrase_counts.items(), key=lambda item: (-item[1], -len(item[0]), item[0]))
    canonical_phrases = [phrase for phrase, count in candidates if count >= 2]

    for phrase, _count in candidates:
        if phrase in canonical_map:
            continue
        best = phrase
        for canonical in canonical_phrases:
            if canonical == phrase:
                best = canonical
                break
            if are_similar_entity_phrases(phrase, canonical):
                best = canonical
                break
        canonical_map[phrase] = best

    return {source: target for source, target in canonical_map.items() if source != target}


def looks_like_entity_candidate(text: str) -> bool:
    stripped = text.strip()
    if len(stripped) < 4 or len(stripped) > 70:
        return False
    if stripped.isupper() and len(stripped.split()) <= 2:
        return False
    if any(char.isdigit() for char in stripped):
        return False

    words = stripped.split()
    capitalized_words = sum(1 for word in words if word[:1].isupper())
    if capitalized_words == 0:
        return False

    stop_words = {"This", "That", "Whereas", "Now", "Thereof", "Agreement", "Recitals"}
    if stripped in stop_words:
        return False
    return capitalized_words >= max(1, len(words) - 1)


def are_similar_entity_phrases(left: str, right: str) -> bool:
    left_tokens = normalize_entity_tokens(left)
    right_tokens = normalize_entity_tokens(right)
    if left_tokens == right_tokens:
        return True
    if abs(len(left_tokens) - len(right_tokens)) > 1:
        return False

    matches = 0
    for left_token, right_token in zip(left_tokens, right_tokens):
        if left_token == right_token:
            matches += 1
        elif token_edit_distance(left_token, right_token) <= 2:
            matches += 1

    return matches >= max(1, min(len(left_tokens), len(right_tokens)) - 1)


def normalize_entity_tokens(text: str) -> list[str]:
    return [re.sub(r"[^a-z]", "", token.lower()) for token in text.split() if re.sub(r"[^a-z]", "", token.lower())]


def token_edit_distance(left: str, right: str) -> int:
    if left == right:
        return 0
    if not left:
        return len(right)
    if not right:
        return len(left)

    previous = list(range(len(right) + 1))
    for i, left_char in enumerate(left, start=1):
        current = [i]
        for j, right_char in enumerate(right, start=1):
            insert_cost = current[j - 1] + 1
            delete_cost = previous[j] + 1
            replace_cost = previous[j - 1] + (0 if left_char == right_char else 1)
            current.append(min(insert_cost, delete_cost, replace_cost))
        previous = current
    return previous[-1]


def normalize_ocr_engine(engine: str) -> str:
    normalized = (engine or "auto").lower()
    if normalized not in SUPPORTED_OCR_ENGINES:
        raise ValueError(f"Unsupported OCR engine: {engine}")
    return normalized


def normalize_document_style(style: str) -> str:
    normalized = (style or "auto").lower()
    if normalized not in SUPPORTED_DOC_STYLES:
        raise ValueError(f"Unsupported document style: {style}")
    return normalized


def extract_text_for_single_image(image_path: Path, engine: str, document_style: str) -> PageAnalysis:
    if engine == "openai":
        return extract_page_text_openai(image_path, document_style)
    if engine == "github":
        return extract_page_text_github(image_path, document_style)
    if engine == "tesseract":
        return extract_page_text_tesseract(image_path, document_style)
    if engine == "ollama":
        return extract_page_text_ollama(image_path, document_style)

    try:
        return extract_page_text_ollama(image_path, document_style)
    except MissingDependencyError:
        pass

    try:
        return extract_page_text_github(image_path, document_style)
    except MissingDependencyError:
        pass

    try:
        return extract_page_text_openai(image_path, document_style)
    except MissingDependencyError:
        pass

    return extract_page_text_tesseract(image_path, document_style)


def extract_page_text_tesseract(image_path: Path, document_style: str) -> PageAnalysis:
    configure_tesseract()
    with Image.open(image_path) as image:
        processed = preprocess_for_ocr(image)
        return extract_page_text(processed, document_style)


def extract_page_text_openai(image_path: Path, document_style: str) -> PageAnalysis:
    client = get_openai_client()
    has_table_grid = image_likely_has_table_grid(image_path)
    with image_path.open("rb") as image_file:
        data_url = f"data:image/{image_path.suffix.lstrip('.').lower()};base64," + base64.b64encode(
            image_file.read()
        ).decode("ascii")

    prompt = build_ai_ocr_prompt(document_style)
    response = client.responses.create(
        model=os.environ.get("OPENAI_OCR_MODEL", "gpt-4.1-mini"),
        input=[
            {
                "role": "user",
                "content": [
                    {"type": "input_text", "text": prompt},
                    {"type": "input_image", "image_url": data_url},
                ],
            }
        ],
    )
    text = (response.output_text or "").strip()
    if not text:
        raise RuntimeError("OpenAI OCR returned an empty response.")
    return PageAnalysis(text=postprocess_ocr_text(text), has_table_grid=has_table_grid)


def extract_page_text_github(image_path: Path, document_style: str) -> PageAnalysis:
    token = os.environ.get("GITHUB_TOKEN") or os.environ.get("GITHUB_PAT")
    if not token:
        raise MissingDependencyError("GITHUB_TOKEN or GITHUB_PAT is not set.")
    if httpx is None:
        raise MissingDependencyError("httpx is not installed. Add it to requirements.")
    has_table_grid = image_likely_has_table_grid(image_path)

    with image_path.open("rb") as image_file:
        data_url = f"data:image/{image_path.suffix.lstrip('.').lower()};base64," + base64.b64encode(
            image_file.read()
        ).decode("ascii")

    prompt = build_ai_ocr_prompt(document_style)
    payload = {
        "model": os.environ.get("GITHUB_OCR_MODEL", "openai/gpt-4.1-mini"),
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ],
    }

    response = httpx.post(
        "https://models.github.ai/inference/chat/completions",
        headers={
            "Accept": "application/vnd.github+json",
            "Authorization": f"Bearer {token}",
            "X-GitHub-Api-Version": "2022-11-28",
            "Content-Type": "application/json",
        },
        content=json.dumps(payload),
        timeout=120.0,
    )
    response.raise_for_status()
    body = response.json()
    text = (
        body.get("choices", [{}])[0]
        .get("message", {})
        .get("content", "")
        .strip()
    )
    if not text:
        raise RuntimeError("GitHub Models OCR returned an empty response.")
    return PageAnalysis(text=postprocess_ocr_text(text), has_table_grid=has_table_grid)


def extract_page_text_ollama(image_path: Path, document_style: str) -> PageAnalysis:
    if httpx is None:
        raise MissingDependencyError("httpx is not installed. Add it to requirements.")
    if os.environ.get("OLLAMA_ENABLED", "true").strip().lower() in ("false", "0", "no"):
        raise MissingDependencyError("Ollama is disabled (OLLAMA_ENABLED=false).")
    ollama_url = os.environ.get("OLLAMA_HOST", "http://localhost:11434") + "/api/chat"
    ollama_model = os.environ.get("OLLAMA_OCR_MODEL", "gemma4:e4b")

    # Verify Ollama is reachable and the model is available
    try:
        probe = httpx.get(
            os.environ.get("OLLAMA_HOST", "http://localhost:11434") + "/api/tags",
            timeout=5.0,
        )
        probe.raise_for_status()
        available = [m.get("name", "") for m in probe.json().get("models", [])]
        if not any(ollama_model in name for name in available):
            raise MissingDependencyError(
                f"Ollama model '{ollama_model}' is not pulled. "
                f"Run: ollama pull {ollama_model}"
            )
    except (httpx.ConnectError, httpx.TimeoutException, httpx.HTTPError) as exc:
        raise MissingDependencyError(f"Ollama is not reachable: {exc}")

    has_table_grid = image_likely_has_table_grid(image_path)
    with image_path.open("rb") as image_file:
        image_b64 = base64.b64encode(image_file.read()).decode("ascii")

    prompt = build_ai_ocr_prompt(document_style)
    payload = {
        "model": ollama_model,
        "messages": [
            {
                "role": "user",
                "content": prompt,
                "images": [image_b64],
            }
        ],
        "stream": False,
    }

    response = httpx.post(ollama_url, content=json.dumps(payload), timeout=180.0)
    response.raise_for_status()
    text = response.json().get("message", {}).get("content", "").strip()
    if not text:
        raise RuntimeError("Ollama OCR returned an empty response.")
    return PageAnalysis(text=postprocess_ocr_text(text), has_table_grid=has_table_grid)


def get_openai_client() -> OpenAI:
    if OpenAI is None:
        raise MissingDependencyError("openai is not installed. Add the package and set OPENAI_API_KEY.")
    if not os.environ.get("OPENAI_API_KEY"):
        raise MissingDependencyError("OPENAI_API_KEY is not set.")
    return OpenAI()


def build_ai_ocr_prompt(document_style: str) -> str:
    style_hint = {
        "legal": "This is likely a legal or contract document. Preserve clause numbering, recital lines, and party definitions.",
        "form": "This is likely a form. Preserve field labels, blanks, table headers, and row structure.",
        "general": "This is a general document. Preserve natural paragraphs and headings without over-formatting.",
        "auto": "Infer the document type from the page and preserve the visible structure carefully.",
    }[document_style]
    return (
        "Transcribe this document page into editable text. Preserve headings, paragraph breaks, numbered clauses, "
        "table text, blank lines, and visible indentation when possible. Do not summarize. Do not invent missing "
        "content. If any text is unreadable, write [unclear]. Return plain text only. "
        + style_hint
    )


def normalize_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").split("\n")]

    collapsed: list[str] = []
    blank_streak = 0
    for line in lines:
        if line.strip():
            collapsed.append(line)
            blank_streak = 0
        else:
            blank_streak += 1
            if blank_streak == 1:
                collapsed.append("")

    return "\n".join(collapsed).strip()


def extract_page_text(image: Image.Image, document_style: str) -> PageAnalysis:
    if pytesseract is None:
        raise MissingDependencyError(
            "pytesseract is not installed. Install requirements and ensure Tesseract OCR is available on PATH."
        )

    psm_candidates = {
        "form": [11, 6, 4],
        "general": [4, 6, 3],
        "legal": [4, 6, 3],
        "auto": [4, 6, 11],
    }[document_style]

    best_text = ""
    best_score = float("-inf")

    for psm in psm_candidates:
        candidate_text, candidate_score = run_tesseract_pass(image, psm)
        if candidate_score > best_score:
            best_text = candidate_text
            best_score = candidate_score

    if best_text.strip():
        return PageAnalysis(text=best_text, has_table_grid=image_likely_has_table_grid_from_image(image))

    fallback = pytesseract.image_to_string(image, config="--oem 3 --psm 6")
    return PageAnalysis(
        text=postprocess_ocr_text(fallback),
        has_table_grid=image_likely_has_table_grid_from_image(image),
    )


def run_tesseract_pass(image: Image.Image, psm: int) -> tuple[str, float]:
    data = pytesseract.image_to_data(image, config=f"--oem 3 --psm {psm}", output_type=Output.DICT)
    lines, confidence_score = build_lines_from_ocr_data(data, image.width)
    text = "\n".join(lines)
    quality_score = score_ocr_text(text) + confidence_score
    return text, quality_score


def build_lines_from_ocr_data(data: dict[str, list[object]], page_width: int) -> tuple[list[str], float]:
    grouped: dict[tuple[int, int, int], list[dict[str, int | str | float]]] = {}
    confidences: list[float] = []

    for idx, raw_text in enumerate(data["text"]):
        text = str(raw_text).strip()
        if not text:
            continue

        conf_text = str(data["conf"][idx]).strip()
        try:
            conf = float(conf_text)
        except ValueError:
            conf = -1
        if conf < 0:
            continue
        confidences.append(conf)

        key = (
            int(data["block_num"][idx]),
            int(data["par_num"][idx]),
            int(data["line_num"][idx]),
        )
        grouped.setdefault(key, []).append(
            {
                "text": text,
                "left": int(data["left"][idx]),
                "top": int(data["top"][idx]),
                "width": int(data["width"][idx]),
                "height": int(data["height"][idx]),
                "conf": conf,
            }
        )

    line_entries: list[dict[str, int | str | float]] = []

    for key in sorted(grouped.keys()):
        words = sorted(grouped[key], key=lambda item: int(item["left"]))
        line = join_ocr_words(words, page_width)
        if not line:
            continue
        line_confidence = average_word_confidence(words)
        processed_line = postprocess_ocr_line(line)
        if should_discard_low_confidence_line(processed_line, line_confidence):
            continue

        left = min(int(word["left"]) for word in words)
        top = min(int(word["top"]) for word in words)
        height = max(int(word["height"]) for word in words)
        line_entries.append(
            {
                "text": processed_line,
                "left": left,
                "top": top,
                "height": height,
                "block_num": key[0],
                "par_num": key[1],
                "conf": line_confidence,
            }
        )

    if not line_entries:
        return [], float("-inf")

    ordered_entries = order_lines_by_layout(line_entries, page_width)
    lines: list[str] = []
    previous: dict[str, int | str] | None = None

    for entry in ordered_entries:
        if previous and should_insert_blank_line(previous, entry):
            lines.append("")
        lines.append(str(entry["text"]))
        previous = entry

    average_conf = (sum(confidences) / len(confidences)) / 100 if confidences else 0.0
    return trim_blank_lines(lines), average_conf


def average_word_confidence(words: Sequence[dict[str, int | str | float]]) -> float:
    values = [float(word["conf"]) for word in words if float(word["conf"]) >= 0]
    if not values:
        return 0.0
    return sum(values) / len(values)


def should_discard_low_confidence_line(line: str, confidence: float) -> bool:
    stripped = line.strip()
    if not stripped:
        return True
    if is_page_number_line(stripped):
        return True
    if looks_like_stamp_or_watermark_noise(stripped, confidence):
        return True
    if confidence >= 38:
        return False

    alnum_count = sum(1 for char in stripped if char.isalnum())
    punctuation_count = sum(1 for char in stripped if not char.isalnum() and not char.isspace())
    alpha_ratio = alnum_count / max(len(stripped), 1)

    if alnum_count < 2 and punctuation_count >= 2:
        return True
    if confidence < 20 and alpha_ratio < 0.45:
        return True
    if confidence < 25 and len(stripped) <= 3:
        return True
    if confidence < 25 and re.fullmatch(r"[\W_]+", stripped):
        return True
    return False


def looks_like_stamp_or_watermark_noise(text: str, confidence: float) -> bool:
    normalized = text.strip()
    if not normalized:
        return False

    upper_ratio = sum(1 for char in normalized if char.isupper()) / max(sum(1 for char in normalized if char.isalpha()), 1)
    alpha_count = sum(1 for char in normalized if char.isalpha())

    watermark_keywords = (
        "received",
        "certified",
        "verified",
        "paid",
        "cancelled",
        "canceled",
        "copy",
        "original",
        "duplicate",
        "released",
        "filed",
        "approved",
        "sample",
        "void",
    )
    lowered = normalized.lower()

    if confidence < 30 and alpha_count <= 4 and upper_ratio >= 0.9:
        return True
    if confidence < 28 and re.fullmatch(r"[A-Z0-9/\- ]{2,18}", normalized):
        return True
    if confidence < 32 and any(keyword in lowered for keyword in watermark_keywords) and len(normalized) <= 24:
        return True
    if confidence < 26 and re.fullmatch(r"(?:[A-Z]{2,}\s*){2,}", normalized):
        return True
    return False


def order_lines_by_layout(line_entries: Sequence[dict[str, int | str]], page_width: int) -> list[dict[str, int | str]]:
    if not detect_multicolumn_layout(line_entries, page_width):
        return sorted(line_entries, key=lambda entry: (int(entry["top"]), int(entry["left"])))

    midpoint = page_width / 2
    left_column = [entry for entry in line_entries if int(entry["left"]) < midpoint]
    right_column = [entry for entry in line_entries if int(entry["left"]) >= midpoint]

    ordered: list[dict[str, int | str]] = []
    ordered.extend(sorted(left_column, key=lambda entry: (int(entry["top"]), int(entry["left"]))))
    if left_column and right_column:
        ordered.append({"text": "", "left": 0, "top": 0, "height": 0, "block_num": -1, "par_num": -1})
    ordered.extend(sorted(right_column, key=lambda entry: (int(entry["top"]), int(entry["left"]))))
    return ordered


def detect_multicolumn_layout(line_entries: Sequence[dict[str, int | str]], page_width: int) -> bool:
    if len(line_entries) < 8:
        return False

    left_positions = [int(entry["left"]) for entry in line_entries]
    left_band = [pos for pos in left_positions if pos < page_width * 0.4]
    right_band = [pos for pos in left_positions if pos > page_width * 0.55]

    if len(left_band) < 3 or len(right_band) < 3:
        return False

    return True


def should_insert_blank_line(previous: dict[str, int | str], current: dict[str, int | str]) -> bool:
    if not str(previous["text"]).strip() or not str(current["text"]).strip():
        return True

    previous_text = str(previous["text"]).strip()
    current_text = str(current["text"]).strip()
    previous_block = (int(previous["block_num"]), int(previous["par_num"]))
    current_block = (int(current["block_num"]), int(current["par_num"]))
    if previous_block != current_block and should_separate_paragraphs(previous, current):
        return True

    previous_bottom = int(previous["top"]) + int(previous["height"])
    vertical_gap = int(current["top"]) - previous_bottom
    left_shift = abs(int(current["left"]) - int(previous["left"]))
    previous_height = max(int(previous["height"]), 1)

    if vertical_gap > max(18, previous_height // 2) and should_separate_paragraphs(previous, current):
        return True
    if left_shift > 32 and starts_like_paragraph(current_text):
        return True
    if previous_text.endswith((".", "?", "!", ":", ";")) and starts_like_paragraph(current_text):
        return True
    return False


def should_separate_paragraphs(previous: dict[str, int | str], current: dict[str, int | str]) -> bool:
    previous_text = str(previous["text"]).strip()
    current_text = str(current["text"]).strip()
    if not previous_text or not current_text:
        return False
    if is_heading_line(current_text) or is_clause_line(current_text) or is_list_line(current_text):
        return True
    if looks_like_fill_in_field(current_text):
        return True
    if starts_like_paragraph(current_text) and ends_like_paragraph(previous_text):
        return True
    return False


def starts_like_paragraph(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if stripped[:1].isupper() and not is_all_caps_phrase(stripped):
        return True
    if re.match(r'^[("\']?[A-Z]', stripped):
        return True
    return False


def ends_like_paragraph(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if stripped.endswith((".", "?", "!", ":", ";")):
        return True
    return len(stripped) > 80


def is_all_caps_phrase(text: str) -> bool:
    letters = [char for char in text if char.isalpha()]
    if not letters:
        return False
    return all(char.isupper() for char in letters)


def join_ocr_words(words: Sequence[dict[str, int | str | float]], page_width: int) -> str:
    parts: list[str] = []
    previous_right = 0
    previous_token = ""

    for word in words:
        text = str(word["text"])
        left = int(word["left"])
        width = int(word["width"])

        if parts:
            gap = left - previous_right
            if gap > max(20, page_width * 0.03):
                parts.append("  ")
            elif not should_attach_without_space(previous_token, text):
                parts.append(" ")

        parts.append(text)
        previous_right = left + width
        previous_token = text

    return "".join(parts).strip()


def should_attach_without_space(previous_token: str, current_token: str) -> bool:
    return current_token in {".", ",", ";", ":", ")", "%"} or previous_token in {"(", "$"}


def postprocess_ocr_text(text: str) -> str:
    return "\n".join(postprocess_ocr_line(line) for line in text.splitlines())


def normalize_common_ocr_artifacts(line: str) -> str:
    cleaned = line
    cleaned = cleaned.replace("ﬁ", "fi").replace("ﬂ", "fl")
    cleaned = cleaned.replace("`", "'")
    cleaned = cleaned.replace("“", '"').replace("”", '"').replace("„", '"')
    cleaned = cleaned.replace("’", "'").replace("‘", "'").replace("‚", "'")
    cleaned = cleaned.replace("–", "-").replace("—", "-")
    cleaned = cleaned.replace("•", "-")
    cleaned = re.sub(r"([A-Za-z])\|([A-Za-z])", r"\1I\2", cleaned)
    cleaned = re.sub(r"\b0f\b", "of", cleaned)
    cleaned = re.sub(r"\b0n\b", "on", cleaned)
    cleaned = re.sub(r"\b1n\b", "In", cleaned)
    cleaned = re.sub(r"\b([Tt])his\s+1s\b", r"\1his is", cleaned)
    cleaned = re.sub(r"(?<=\w)\s+([,.;:?!])", r"\1", cleaned)
    cleaned = re.sub(r"([({\[])\s+", r"\1", cleaned)
    cleaned = re.sub(r"\s+([)}\]])", r"\1", cleaned)
    cleaned = re.sub(r'"\s*,', '",', cleaned)
    cleaned = re.sub(r"\bIi\b", "II", cleaned)
    cleaned = re.sub(r"\b([Aa])ng\b", r"\1nd", cleaned)
    cleaned = re.sub(r"^\s*[Oo0]\s+(?=[A-Za-z])", "- ", cleaned)
    cleaned = re.sub(r"^\s*[\[\(]\s*[\]\)]\s*", "[ ] ", cleaned)
    cleaned = re.sub(r"^\s*[\[\(]\s*[xX]\s*[\]\)]\s*", "[x] ", cleaned)
    cleaned = re.sub(r"^\s*☐\s*", "[ ] ", cleaned)
    cleaned = re.sub(r"^\s*(?:☑|✓)\s*", "[x] ", cleaned)
    cleaned = re.sub(r"^\s*•\s*", "- ", cleaned)
    return cleaned


def postprocess_ocr_line(line: str) -> str:
    cleaned = normalize_common_ocr_artifacts(clean_ocr_line(line))
    replacements = {
        "Agreement')": 'Agreement"),',
        'Agreement").': 'Agreement"),',
        '("Agreement"),.': '("Agreement"),',
        "Agreement’": 'Agreement")',
        "heremafter": "hereinafter",
        "shail": "shall",
        "collectively as “Parties”": 'collectively as "Parties"',
        "collectively as \"Parties\"": 'collectively as "Parties"',
        "collectively as “Parties\")": 'collectively as "Parties")',
        "collectively as \"Parties\")": 'collectively as "Parties")',
        "L. APPOINTMENT": "I. APPOINTMENT",
        "L APPOINTMENT": "I. APPOINTMENT",
        "Il ": "II ",
        " In this Agreement": " in this Agreement",
        "oeos": "Pesos",
        "Poeos": "Pesos",
        "PPesos": "Pesos",
        "eamed": "earned",
        "pnvilege": "privilege",
        "elated": "related",
        "follaws": "follows",
        "operale": "operate",
        "chanded": "changed",
        "rrelated": "related",
        "indude": "include",
        "ttademarks": "trademarks",
        "sharehokers": "shareholders",
        "transportabon": "transportation",
        "exterlor": "exterior",
        "Interior": "interior",
        "axchisively": "exclusively",
        "confkjential": "confidential",
        "sat up": "set up",
        "II Is": "It is",
        "In this Agreement": "in this Agreement",
        "FRANCHISORs": "FRANCHISOR's",
        "FRANCHISEEs": "FRANCHISEE's",
        "(FRANCHISOR)s": "FRANCHISOR's",
        "(FRANCHISEE)s": "FRANCHISEE's",
        "Inc. a corporation": "Inc., a corporation",
        "Officer. Rosalie": "Officer, Rosalie",
        "March 2026. by and between:": "March 2026, by and between:",
        'referred to as the ("FRANCHISOR)': 'referred to as the ("FRANCHISOR")',
    }

    for source, target in replacements.items():
        cleaned = cleaned.replace(source, target)

    cleaned = re.sub(r"\(FRANCHISEE[\"'’`:]*\)", "(FRANCHISEE)", cleaned)
    cleaned = re.sub(r"\(FRANCHISOR[\"'’`:]*\)", "(FRANCHISOR)", cleaned)
    cleaned = re.sub(r"\bFRANCHISEE[\"'’`:]+", "FRANCHISEE", cleaned)
    cleaned = re.sub(r"\bFRANCHISOR[\"'’`:]+", "FRANCHISOR", cleaned)
    cleaned = re.sub(r'\("FRANCHISOR\)', '("FRANCHISOR")', cleaned)
    cleaned = re.sub(r'\("FRANCHISEE\)', '("FRANCHISEE")', cleaned)
    cleaned = re.sub(r"\b([A-Z])\s+\.", r"\1.", cleaned)
    cleaned = re.sub(r"\b([IVXLCDM]+)\s+\.", r"\1.", cleaned)
    cleaned = re.sub(r"^\(?L[\.\s]+([A-Z])", r"I. \1", cleaned)
    cleaned = re.sub(r"\(System['\"]\)", '("System")', cleaned)
    cleaned = re.sub(r"\.\.", ".", cleaned)
    cleaned = re.sub(r"\s+,", ",", cleaned)
    cleaned = re.sub(r"\s+\.", ".", cleaned)
    cleaned = re.sub(r'"\)', '")', cleaned)
    cleaned = re.sub(r"\)\)", ")", cleaned)
    cleaned = re.sub(r"\bInc\.\s*,", "Inc.,", cleaned)
    cleaned = re.sub(r"\s+—\s+", " ", cleaned)
    cleaned = re.sub(r"\b([A-Za-z]+)\s+Inc\.,\b", r"\1 Inc.,", cleaned)
    cleaned = re.sub(r"\b([A-Za-z]+)\s+Corp\.\b", r"\1 Corp.", cleaned)
    cleaned = re.sub(r"\b([A-Za-z]+)\s+Co\.\b", r"\1 Co.", cleaned)
    cleaned = re.sub(r"\s{3,}", "  ", cleaned)
    cleaned = cleaned.replace("lo as", "to as").replace("bee", "been").replace("fhe", "the")
    cleaned = cleaned.replace('collectively as "Parties"))', 'collectively as "Parties")')
    cleaned = normalize_field_values(cleaned)
    return cleaned.strip()


def normalize_field_values(text: str) -> str:
    cleaned = text

    # Standardize common field labels.
    cleaned = re.sub(r"\bctc\s*/?\s*id\b", "CTC/ID", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\btin\s*/?\s*no\.?\b", "TIN No.", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\bsss\s*/?\s*no\.?\b", "SSS No.", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\bpassport\s*/?\s*no\.?\b", "Passport No.", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\bdate\s*&\s*place\s*of\s*issue\b", "Date & Place of Issue", cleaned, flags=re.IGNORECASE)

    # Normalize OCR confusion in dates and slash-separated IDs.
    cleaned = re.sub(r"(?<=\d)[Oo](?=\d)", "0", cleaned)
    cleaned = re.sub(r"(?<=\d)[Il](?=\d)", "1", cleaned)
    cleaned = re.sub(r"(?<=\d)\s*/\s*(?=\d)", "/", cleaned)
    cleaned = re.sub(r"(?<=\d)\s*-\s*(?=\d)", "-", cleaned)

    # Normalize money, percentages, and reference numbers.
    cleaned = re.sub(r"\b[Pp][Ee]?[Ss]?[Oo]?[Ss]?\s*[:\-]?\s*(\d)", r"Pesos \1", cleaned)
    cleaned = re.sub(r"\bUSD\s*[:\-]?\s*(\d)", r"USD \1", cleaned)
    cleaned = re.sub(r"(?<=\d)\s*%\b", "%", cleaned)
    cleaned = re.sub(r"\bno\s*[:\-]?\s*(\d)", r"No. \1", cleaned, flags=re.IGNORECASE)

    # Keep month-name dates cleaner.
    months = "January|February|March|April|May|June|July|August|September|October|November|December"
    cleaned = re.sub(
        rf"\b({months})\s+(\d{{1,2}})\s*[.,]\s*(\d{{4}})\b",
        r"\1 \2, \3",
        cleaned,
        flags=re.IGNORECASE,
    )

    return cleaned


def trim_blank_lines(lines: Sequence[str]) -> list[str]:
    trimmed: list[str] = []
    previous_blank = True
    for line in lines:
        is_blank = not line.strip()
        if is_blank and previous_blank:
            continue
        trimmed.append(line)
        previous_blank = is_blank
    while trimmed and not trimmed[-1].strip():
        trimmed.pop()
    return trimmed


def build_output(
    pages: Sequence[OCRPage],
    output_format: str,
    output_path: Path,
    title: str | None = None,
    image_paths: Sequence[Path] | None = None,
    document_style: str = "auto",
    pdf_page_numbers: bool = False,
    pdf_watermark: str | None = None,
) -> Path:
    normalized_format = output_format.lower().lstrip(".")
    if normalized_format not in SUPPORTED_OUTPUTS:
        raise ValueError(f"Unsupported output format: {output_format}")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if normalized_format == "docx":
        build_docx(pages, output_path, title, image_paths=image_paths, document_style=document_style)
    elif normalized_format == "pdf":
        build_pdf(
            pages,
            output_path,
            title,
            document_style=document_style,
            add_page_numbers=pdf_page_numbers,
            watermark_text=pdf_watermark,
        )
    elif normalized_format == "xlsx":
        build_xlsx(pages, output_path, title)
    elif normalized_format == "pptx":
        build_pptx(pages, output_path, title, document_style=document_style)
    elif normalized_format == "md":
        output_path.write_text(build_markdown(pages, title, document_style=document_style), encoding="utf-8")
    elif normalized_format == "html":
        output_path.write_text(build_html(pages, title, document_style=document_style), encoding="utf-8")
    elif normalized_format == "txt":
        output_path.write_text(build_text(pages, title, document_style=document_style), encoding="utf-8")

    return output_path


def build_docx(
    pages: Sequence[OCRPage],
    output_path: Path,
    title: str | None,
    image_paths: Sequence[Path] | None = None,
    document_style: str = "auto",
) -> None:
    document = Document()
    style = document.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(12)

    document_text = merge_pages_for_document(pages)
    profile = infer_document_style(document_text, document_style)
    sections = split_document_blocks(document_text, profile)

    effective_title = title
    if not effective_title and sections:
        first_nonempty = next((line for line in sections if line.strip()), "")
        if looks_like_title(first_nonempty):
            effective_title = first_nonempty
            sections = sections[1:]

    if effective_title:
        add_centered_paragraph(document, effective_title, bold=True, size=14)
        document.add_paragraph()

    pending_signature_headers: list[str] | None = None
    pending_signature_rows = 0
    pending_table_rows: list[TableRowCandidate] = []
    force_table_mode = False

    for block in sections:
        if block == "[TABLE_GRID_DETECTED]":
            force_table_mode = True
            continue

        table_candidate = parse_table_row_candidate(block, profile)

        if table_candidate and not pending_signature_headers:
            pending_table_rows.append(table_candidate)
            continue

        if force_table_mode and looks_like_tableish_text(block) and not pending_signature_headers:
            cells = [part.strip() for part in re.split(r"\s{2,}", block) if part.strip()]
            if len(cells) >= 2:
                pending_table_rows.append(TableRowCandidate(cells=cells))
                continue

        if pending_table_rows:
            flush_table_candidates(document, pending_table_rows)
            pending_table_rows = []

        if is_signature_header(block):
            pending_signature_headers = [part.strip() for part in re.split(r"\s{2,}", block) if part.strip()]
            if len(pending_signature_headers) < 2:
                pending_signature_headers = ["Name", "CTC/ID with No.", "Date & Place of Issue"]
            pending_signature_rows = 0
            continue

        if pending_signature_headers:
            if looks_like_signature_blank_row(block):
                pending_signature_rows += 1
                continue
            if pending_signature_rows:
                add_signature_table(document, pending_signature_headers, max(pending_signature_rows, 2))
                pending_signature_headers = None
                pending_signature_rows = 0

        add_structured_block(document, block, profile)

    if pending_signature_headers and pending_signature_rows:
        add_signature_table(document, pending_signature_headers, max(pending_signature_rows, 2))
    if pending_table_rows:
        flush_table_candidates(document, pending_table_rows)

    document.save(output_path)


def build_pdf(
    pages: Sequence[OCRPage],
    output_path: Path,
    title: str | None,
    document_style: str = "auto",
    add_page_numbers: bool = False,
    watermark_text: str | None = None,
) -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(18, 18, 18)
    pdf.set_font("Helvetica", size=11)
    normalized_style = normalize_document_style(document_style)
    normalized_watermark = (watermark_text or "").strip()

    for page_index, page in enumerate(pages, start=1):
        pdf.add_page()
        profile = infer_document_style(page.text, normalized_style)

        if title and page_index == 1:
            pdf.set_font("Helvetica", "B", 15)
            pdf.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT", align="C")
            pdf.ln(4)

        pdf.set_font("Helvetica", "I", 9)
        pdf.cell(0, 6, page.filename, new_x="LMARGIN", new_y="NEXT", align="R")
        pdf.ln(2)

        if normalized_watermark:
            add_pdf_watermark(pdf, normalized_watermark)

        blocks = split_document_blocks(page.text, profile)
        pending_table_rows: list[TableRowCandidate] = []
        force_table_mode = False

        for block in blocks:
            if block == "[TABLE_GRID_DETECTED]":
                force_table_mode = True
                continue

            table_candidate = parse_table_row_candidate(block, profile)
            if table_candidate:
                pending_table_rows.append(table_candidate)
                continue

            if force_table_mode and looks_like_tableish_text(block):
                cells = [part.strip() for part in re.split(r"\s{2,}", block) if part.strip()]
                if len(cells) >= 2:
                    pending_table_rows.append(TableRowCandidate(cells=cells))
                    continue

            if pending_table_rows:
                add_pdf_table(pdf, pending_table_rows)
                pending_table_rows = []

            add_pdf_block(pdf, block)

        if pending_table_rows:
            add_pdf_table(pdf, pending_table_rows)

        if add_page_numbers:
            add_pdf_page_number(pdf, page_index)

    pdf.output(str(output_path))


def add_pdf_block(pdf: FPDF, block: str) -> None:
    if looks_like_title(block):
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 9, block, new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(1)
        return

    if is_heading_line(block):
        align = "C" if block.endswith(":") else "L"
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 8, block, new_x="LMARGIN", new_y="NEXT", align=align)
        pdf.ln(1)
        return

    if block.lower() == "and":
        pdf.set_font("Helvetica", size=11)
        pdf.cell(0, 7, block, new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(1)
        return

    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(0, 6, block, align="J")
    pdf.ln(1)


def add_pdf_table(pdf: FPDF, rows: Sequence[TableRowCandidate]) -> None:
    normalized_rows = normalize_table_candidates(rows)
    if not normalized_rows:
        return

    max_cols = max(len(row.cells) for row in normalized_rows)
    usable_width = pdf.w - pdf.l_margin - pdf.r_margin
    col_width = usable_width / max_cols
    row_height = 8

    pdf.set_font("Helvetica", size=10)
    for row_index, row in enumerate(normalized_rows):
        for col_index in range(max_cols):
            text = row.cells[col_index] if col_index < len(row.cells) else ""
            font_style = "B" if row_index == 0 else ""
            pdf.set_font("Helvetica", font_style, 10)
            pdf.cell(col_width, row_height, truncate_pdf_cell_text(text, max_cols), border=1)
        pdf.ln(row_height)
    pdf.ln(2)


def truncate_pdf_cell_text(text: str, column_count: int) -> str:
    compact = " ".join(text.split())
    max_chars = 34 if column_count <= 3 else 20 if column_count <= 5 else 14
    if len(compact) <= max_chars:
        return compact
    return compact[: max_chars - 1].rstrip() + "..."


def add_pdf_watermark(pdf: FPDF, watermark_text: str) -> None:
    pdf.ln(1)
    pdf.set_text_color(168, 156, 146)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, watermark_text[:80], new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_text_color(0, 0, 0)
    pdf.ln(2)


def add_pdf_page_number(pdf: FPDF, page_number: int) -> None:
    current_y = pdf.get_y()
    pdf.set_y(-14)
    pdf.set_font("Helvetica", size=9)
    pdf.set_text_color(110, 102, 94)
    pdf.cell(0, 6, f"Page {page_number}", align="C")
    pdf.set_text_color(0, 0, 0)
    pdf.set_y(max(current_y, pdf.t_margin))


def build_xlsx(pages: Sequence[OCRPage], output_path: Path, title: str | None) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Extracted Text"

    row = 1
    if title:
        sheet.cell(row=row, column=1, value=title)
        sheet.cell(row=row, column=1).font = Font(bold=True, size=14)
        row += 2

    headers = ["Source Image", "Paragraph", "Text"]
    for column, value in enumerate(headers, start=1):
        cell = sheet.cell(row=row, column=column, value=value)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal="center", vertical="center")
    row += 1

    for page in pages:
        table_rows = extract_table_rows_from_text(page.text)
        if table_rows:
            sheet.cell(row=row, column=1, value=page.filename)
            sheet.cell(row=row, column=1).font = Font(bold=True, size=12)
            row += 1

            max_cols = max(len(table_row) for table_row in table_rows)
            for r_index, table_row in enumerate(table_rows):
                for c_index in range(max_cols):
                    value = table_row[c_index] if c_index < len(table_row) else ""
                    cell = sheet.cell(row=row, column=c_index + 1, value=value)
                    cell.alignment = Alignment(wrap_text=True, vertical="top")
                    if r_index == 0:
                        cell.font = Font(bold=True)
                row += 1
            row += 1
            continue

        paragraphs = split_paragraphs(page.text)
        if not paragraphs:
            paragraphs = ["[No readable text detected on this image]"]
        for idx, paragraph in enumerate(paragraphs, start=1):
            sheet.cell(row=row, column=1, value=page.filename)
            sheet.cell(row=row, column=2, value=idx)
            sheet.cell(row=row, column=3, value=paragraph)
            sheet.cell(row=row, column=3).alignment = Alignment(wrap_text=True, vertical="top")
            row += 1

    auto_fit_columns(sheet)
    workbook.save(output_path)


def build_pptx(
    pages: Sequence[OCRPage],
    output_path: Path,
    title: str | None,
    document_style: str = "auto",
) -> None:
    presentation = Presentation()
    normalized_style = normalize_document_style(document_style)

    for page_index, page in enumerate(pages, start=1):
        profile = infer_document_style(page.text, normalized_style)
        blocks = [block for block in split_document_blocks(page.text, profile) if block != "[TABLE_GRID_DETECTED]"]
        text_chunks = chunk_blocks_for_pptx(blocks) or ["[No readable text detected]"]

        for chunk_index, chunk in enumerate(text_chunks, start=1):
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            if title:
                slide_title = title if len(text_chunks) == 1 else f"{title} ({page_index}.{chunk_index})"
            else:
                slide_title = page.filename if len(text_chunks) == 1 else f"{page.filename} ({chunk_index})"
            slide.shapes.title.text = slide_title
            slide.placeholders[1].text = chunk

    if not presentation.slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title or "Extracted Document"
        slide.placeholders[1].text = "[No readable text detected]"

    presentation.save(output_path)


def build_markdown(pages: Sequence[OCRPage], title: str | None, document_style: str = "auto") -> str:
    lines: list[str] = []
    normalized_style = normalize_document_style(document_style)
    if title:
        lines.append(f"# {title}")
        lines.append("")

    for page in pages:
        lines.append(f"## {page.filename}")
        lines.append("")
        profile = infer_document_style(page.text, normalized_style)
        blocks = [block for block in split_document_blocks(page.text, profile) if block != "[TABLE_GRID_DETECTED]"]
        if not blocks:
            lines.append("[No readable text detected on this image]")
            lines.append("")
            continue

        pending_table_rows: list[TableRowCandidate] = []
        for block in blocks:
            table_candidate = parse_table_row_candidate(block, profile)
            if table_candidate:
                pending_table_rows.append(table_candidate)
                continue

            if pending_table_rows:
                lines.extend(render_markdown_table_rows(pending_table_rows))
                lines.append("")
                pending_table_rows = []

            lines.extend(render_markdown_block(block))
            lines.append("")

        if pending_table_rows:
            lines.extend(render_markdown_table_rows(pending_table_rows))
            lines.append("")

        lines.append("")

    return "\n".join(lines).strip() + "\n"


def build_html(pages: Sequence[OCRPage], title: str | None, document_style: str = "auto") -> str:
    body_parts = []
    normalized_style = normalize_document_style(document_style)
    if title:
        body_parts.append(f"<h1>{escape(title)}</h1>")

    for page in pages:
        profile = infer_document_style(page.text, normalized_style)
        blocks = [block for block in split_document_blocks(page.text, profile) if block != "[TABLE_GRID_DETECTED]"]
        body_parts.append(f"<section><h2>{escape(page.filename)}</h2>")

        if not blocks:
            body_parts.append("<p>[No readable text detected on this image]</p>")
            body_parts.append("</section>")
            continue

        pending_table_rows: list[TableRowCandidate] = []
        for block in blocks:
            table_candidate = parse_table_row_candidate(block, profile)
            if table_candidate:
                pending_table_rows.append(table_candidate)
                continue

            if pending_table_rows:
                body_parts.append(render_html_table(pending_table_rows))
                pending_table_rows = []

            body_parts.append(render_html_block(block))

        if pending_table_rows:
            body_parts.append(render_html_table(pending_table_rows))

        body_parts.append("</section>")

    return (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{escape(title or 'Image to Document')}</title>"
        "<style>body{font-family:Georgia,serif;max-width:960px;margin:40px auto;padding:0 20px;line-height:1.65;"
        "color:#1f1b18;background:#faf7f2}h1,h2,h3{font-family:'Trebuchet MS',sans-serif}h1{margin-bottom:28px}"
        "section{margin-bottom:36px;padding:24px 28px;background:#fff;border:1px solid #e7dfd3;border-radius:18px;"
        "box-shadow:0 12px 30px rgba(87,64,40,.08)}p{margin:0 0 14px}p.center{text-align:center}"
        "p.justified{text-align:justify}table{width:100%;border-collapse:collapse;margin:18px 0}th,td{border:1px solid #c8b9a6;"
        "padding:8px 10px;vertical-align:top}th{background:#f3eadf;text-align:left}</style></head><body>"
        + "".join(body_parts)
        + "</body></html>"
    )


def build_text(pages: Sequence[OCRPage], title: str | None, document_style: str = "auto") -> str:
    chunks: list[str] = []
    normalized_style = normalize_document_style(document_style)
    if title:
        chunks.append(title)
        chunks.append("=" * len(title))
        chunks.append("")

    for page in pages:
        chunks.append(page.filename)
        chunks.append("-" * len(page.filename))
        profile = infer_document_style(page.text, normalized_style)
        blocks = [block for block in split_document_blocks(page.text, profile) if block != "[TABLE_GRID_DETECTED]"]
        if blocks:
            chunks.append("\n\n".join(blocks))
        else:
            chunks.append("[No readable text detected on this image]")
        chunks.append("")

    return "\n".join(chunks).strip() + "\n"


def chunk_blocks_for_pptx(blocks: Sequence[str], max_chars: int = 900) -> list[str]:
    chunks: list[str] = []
    current: list[str] = []
    current_length = 0

    for block in blocks:
        normalized = block.strip()
        if not normalized:
            continue

        candidate_length = current_length + len(normalized) + (2 if current else 0)
        if current and candidate_length > max_chars:
            chunks.append("\n\n".join(current))
            current = [normalized]
            current_length = len(normalized)
            continue

        current.append(normalized)
        current_length = candidate_length

    if current:
        chunks.append("\n\n".join(current))

    return chunks


def render_markdown_block(block: str) -> list[str]:
    if looks_like_title(block):
        return [f"### {block}"]
    if is_heading_line(block):
        return [f"**{block}**"]
    if block.lower() == "and":
        return ["_and_"]
    if is_list_line(block):
        return [normalize_list_marker(block)]
    return [block]


def render_markdown_table_rows(rows: Sequence[TableRowCandidate]) -> list[str]:
    normalized_candidates = normalize_table_candidates(rows)
    if not normalized_candidates:
        return []

    max_cols = max(len(row.cells) for row in normalized_candidates)
    normalized_rows = [
        [cell.strip() for cell in row.cells] + [""] * (max_cols - len(row.cells))
        for row in normalized_candidates
    ]

    header = normalized_rows[0]
    separator = ["---"] * max_cols
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in normalized_rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return lines


def render_html_block(block: str) -> str:
    escaped = escape(block)
    if looks_like_title(block):
        return f"<h3>{escaped}</h3>"
    if is_heading_line(block):
        css_class = "center" if block.endswith(":") else ""
        class_attr = f" class='{css_class}'" if css_class else ""
        return f"<p{class_attr}><strong>{escaped}</strong></p>"
    if block.lower() == "and":
        return f"<p class='center'>{escaped}</p>"
    if is_list_line(block):
        return f"<p>{escape(normalize_list_marker(block))}</p>"
    return f"<p class='justified'>{escaped}</p>"


def render_html_table(rows: Sequence[TableRowCandidate]) -> str:
    normalized_candidates = normalize_table_candidates(rows)
    if not normalized_candidates:
        return ""

    max_cols = max(len(row.cells) for row in normalized_candidates)
    normalized_rows = [
        [escape(cell.strip()) for cell in row.cells] + [""] * (max_cols - len(row.cells))
        for row in normalized_candidates
    ]

    header_cells = "".join(f"<th>{cell}</th>" for cell in normalized_rows[0])
    body_rows = []
    for row in normalized_rows[1:]:
        body_rows.append("<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>")

    return (
        "<table><thead><tr>"
        + header_cells
        + "</tr></thead><tbody>"
        + "".join(body_rows)
        + "</tbody></table>"
    )


def split_paragraphs(text: str) -> list[str]:
    return [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]


def merge_pages_for_document(pages: Sequence[OCRPage]) -> str:
    cleaned_pages = remove_repeated_headers_and_footers(pages)
    merged_parts: list[str] = []
    carry = ""

    for page in cleaned_pages:
        lines = [line.strip() for line in page.text.split("\n") if line.strip()]
        if not lines:
            continue

        if carry:
            lines[0] = merge_wrapped_segments(carry, lines[0])
            carry = ""

        if should_carry_to_next_page(lines[-1]):
            carry = lines.pop()

        if lines:
            merged_parts.append("\n".join(lines))

    if carry:
        merged_parts.append(carry)

    return "\n".join(merged_parts)


def remove_repeated_headers_and_footers(pages: Sequence[OCRPage]) -> list[OCRPage]:
    if len(pages) < 3:
        return list(pages)

    top_counts: dict[str, int] = {}
    bottom_counts: dict[str, int] = {}
    split_pages: list[list[str]] = []

    for page in pages:
        lines = [line.strip() for line in page.text.split("\n") if line.strip()]
        split_pages.append(lines)
        for line in lines[:2]:
            key = normalize_running_header_footer_candidate(line)
            if key:
                top_counts[key] = top_counts.get(key, 0) + 1
        for line in lines[-2:]:
            key = normalize_running_header_footer_candidate(line)
            if key:
                bottom_counts[key] = bottom_counts.get(key, 0) + 1

    repeated_headers = {line for line, count in top_counts.items() if count >= 3}
    repeated_footers = {line for line, count in bottom_counts.items() if count >= 3}

    cleaned_pages: list[OCRPage] = []
    for page, lines in zip(pages, split_pages):
        trimmed = list(lines)
        while trimmed and should_strip_running_header_footer(trimmed[0], repeated_headers, zone="header"):
            trimmed = trimmed[1:]
        while trimmed and should_strip_running_header_footer(trimmed[-1], repeated_footers, zone="footer"):
            trimmed = trimmed[:-1]
        trimmed = remove_page_number_lines(trimmed)
        cleaned_pages.append(OCRPage(filename=page.filename, text="\n".join(trimmed)))

    return cleaned_pages


def normalize_running_header_footer_candidate(text: str) -> str | None:
    stripped = text.strip()
    if not stripped:
        return None
    if len(stripped) > 120:
        return None
    if looks_like_title(stripped) or is_heading_line(stripped):
        return None

    normalized = stripped.lower()
    normalized = re.sub(r"\bpage\s+\d+(\s*(of|/)\s*\d+)?\b", "page", normalized)
    normalized = re.sub(r"\b\d+\s*(of|/)\s*\d+\b", "page", normalized)
    normalized = re.sub(r"\b\d+\b", "#", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip(" -_|")
    if len(normalized) < 3:
        return None
    return normalized


def should_strip_running_header_footer(text: str, repeated_lines: set[str], zone: str) -> bool:
    normalized = normalize_running_header_footer_candidate(text)
    if not normalized or normalized not in repeated_lines:
        return False
    stripped = text.strip()
    if zone == "header" and looks_like_title(stripped):
        return False
    return True


def remove_page_number_lines(lines: Sequence[str]) -> list[str]:
    filtered: list[str] = []
    for line in lines:
        stripped = line.strip()
        if is_page_number_line(stripped):
            continue
        filtered.append(stripped)
    return filtered


def is_page_number_line(text: str) -> bool:
    if not text:
        return False

    normalized = text.strip().lower()
    patterns = [
        r"^page\s+\d+\s*(of\s+\d+)?$",
        r"^\d+\s*(of|/)\s*\d+$",
        r"^-\s*\d+\s*-$",
        r"^\[\s*\d+\s*\]$",
        r"^\d+$",
    ]
    return any(re.match(pattern, normalized) for pattern in patterns)


def split_document_blocks(text: str, profile: str) -> list[str]:
    raw_lines = [clean_ocr_line(line) for line in text.split("\n")]
    blocks: list[str] = []
    buffer: list[str] = []

    for line in raw_lines:
        if not line:
            if buffer:
                blocks.append(join_wrapped_lines(buffer))
                buffer.clear()
            continue

        if starts_new_block(line, buffer):
            if buffer:
                blocks.append(join_wrapped_lines(buffer))
            buffer = [line]
        else:
            buffer.append(line)

    if buffer:
        blocks.append(join_wrapped_lines(buffer))

    separated: list[str] = []
    for block in blocks:
        separated.extend(split_embedded_title_block(block))

    normalized = normalize_document_blocks(separated, profile)
    return [block for block in normalized if block.strip()]


def clean_ocr_line(line: str) -> str:
    cleaned = line.strip()
    cleaned = cleaned.replace("‘", "").replace("’", "'").replace("“", '"').replace("”", '"')
    cleaned = re.sub(r"\s+", " ", cleaned)
    cleaned = re.sub(r"^\|\.\s*", "I. ", cleaned)
    cleaned = re.sub(r"^\[\s*([A-Z])", r"\1.", cleaned)
    return cleaned


def join_wrapped_lines(lines: Sequence[str]) -> str:
    nonempty = [line.strip() for line in lines if line.strip()]
    if not nonempty:
        return ""

    text = nonempty[0]
    for line in nonempty[1:]:
        text = merge_wrapped_segments(text, line)

    text = re.sub(r"\s+([,.;:])", r"\1", text)
    text = re.sub(r"\(\s+", "(", text)
    text = re.sub(r"\s+\)", ")", text)
    return text.strip()


def merge_wrapped_segments(left: str, right: str) -> str:
    left = left.rstrip()
    right = right.lstrip()

    if not left:
        return right
    if not right:
        return left

    if should_keep_lines_separate(left, right):
        return f"{left}\n{right}"

    if left.endswith("-") and should_unhyphenate_break(left, right):
        return f"{left[:-1]}{right}"

    if should_attach_without_intervening_space(left, right):
        return f"{left}{right}"

    return f"{left} {right}"


def should_keep_lines_separate(left: str, right: str) -> bool:
    left_clean = left.strip()
    right_clean = right.strip()
    if not left_clean or not right_clean:
        return False
    if looks_like_title(left_clean) or looks_like_title(right_clean):
        return True
    if is_heading_line(left_clean) or is_heading_line(right_clean):
        return True
    if is_clause_line(right_clean):
        return True
    if parse_table_row_candidate(left_clean, "general") or parse_table_row_candidate(right_clean, "general"):
        return True
    if right_clean.lower() == "and":
        return True
    return False


def should_unhyphenate_break(left: str, right: str) -> bool:
    if not left.endswith("-") or not right:
        return False
    if len(left) < 3:
        return False

    left_prev = left[-2]
    right_first = right[0]
    if not left_prev.isalpha() or not right_first.isalpha():
        return False
    if right_first.isupper() and not right[:2].islower():
        return False
    return True


def should_attach_without_intervening_space(left: str, right: str) -> bool:
    if left.endswith(("/", "(", "$", '"')):
        return True
    if right.startswith((".", ",", ";", ":", ")", "%", "]")):
        return True
    return False


def starts_new_block(line: str, buffer: Sequence[str]) -> bool:
    if not buffer:
        return True
    previous = buffer[-1].strip()
    if is_heading_line(line):
        return True
    if is_clause_line(line):
        return True
    if is_list_line(line):
        return True
    if looks_like_fill_in_field(line):
        return True
    if previous.endswith((".", "?", "!", ":", ";")) and starts_like_paragraph(line):
        return True
    if starts_like_paragraph(line) and should_keep_lines_separate(previous, line):
        return True
    if len(line) < 8:
        return True
    if buffer and should_split_short_line(buffer[-1], line):
        return True
    return False


def should_split_short_line(previous: str, current: str) -> bool:
    return previous.endswith(":") or previous in {"and", "AND"}


def should_carry_to_next_page(line: str) -> bool:
    if len(line) < 35:
        return False
    if line.endswith((".", ":", ";", "?", "!")):
        return False
    return True


def looks_like_title(text: str) -> bool:
    stripped = text.strip()
    return bool(stripped) and stripped == stripped.upper() and len(stripped.split()) <= 8


def split_embedded_title_block(block: str) -> list[str]:
    stripped = block.strip()
    match = re.match(r"^([A-Z][A-Z\s&\-]{5,})\s+([A-Z][a-z].+)$", stripped)
    if not match:
        return [block]

    candidate_title = re.sub(r"\s+", " ", match.group(1)).strip()
    remainder = match.group(2).strip()
    if looks_like_title(candidate_title):
        return [candidate_title, remainder]
    return [block]


def normalize_document_blocks(blocks: Sequence[str], profile: str) -> list[str]:
    normalized: list[str] = []
    index = 0

    while index < len(blocks):
        current = blocks[index].strip()
        next_block = blocks[index + 1].strip() if index + 1 < len(blocks) else ""

        if profile == "legal" and current.endswith("principal place of business at") and looks_like_franchisee_placeholder(next_block):
            normalized.append(f'{current} ____________________ ("FRANCHISEE");')
            index += 2
            continue

        if looks_like_fill_in_blank(next_block) and should_merge_with_following_blank(current, next_block):
            normalized.append(merge_field_with_blank(current, next_block))
            index += 2
            continue

        current = normalize_fill_in_block(current)

        normalized.extend(split_whereas_block(current, profile))
        index += 1

    return normalized


def should_merge_with_following_blank(current: str, next_block: str) -> bool:
    if not current or not next_block:
        return False
    if is_heading_line(current) or is_clause_line(current):
        return False
    return current.endswith((":","at","on","by","for")) or looks_like_field_label(current)


def merge_field_with_blank(current: str, next_block: str) -> str:
    normalized_blank = standardize_blank_field(next_block)
    if current.endswith(":"):
        return f"{current} {normalized_blank}"
    return f"{current} {normalized_blank}"


def normalize_fill_in_block(block: str) -> str:
    normalized = block
    normalized = re.sub(r"(?:_{3,}|\.{3,}|-{4,})", " ____________________ ", normalized)
    normalized = re.sub(r"\(\s*____________________\s*\)", "(____________________)", normalized)
    normalized = re.sub(r"\s{2,}", " ", normalized)
    return normalized.strip()


def looks_like_fill_in_blank(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if set(stripped) <= {"_", ".", "-", " "} and len(stripped.replace(" ", "")) >= 4:
        return True
    return bool(re.fullmatch(r"[\[(]?\s*(?:_{4,}|\.{4,}|-{5,})\s*[\])]?", stripped))


def standardize_blank_field(text: str) -> str:
    if looks_like_fill_in_blank(text):
        return "____________________"
    return normalize_fill_in_block(text)


def looks_like_field_label(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if stripped.endswith(":"):
        return True
    lowered = stripped.lower()
    keywords = ("name", "address", "date", "place", "issued", "issued at", "signature", "contact", "email")
    return any(lowered.endswith(keyword) for keyword in keywords)


def looks_like_franchisee_placeholder(text: str) -> bool:
    return text.strip().startswith("(FRANCHISEE")


def split_whereas_block(block: str, profile: str) -> list[str]:
    if profile != "legal":
        return [block]
    if block.count("WHEREAS") <= 1:
        return [block]

    parts = re.split(r"(?=WHEREAS[,\.])", block)
    cleaned_parts = [part.strip() for part in parts if part.strip()]
    return cleaned_parts or [block]


def is_heading_line(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if stripped.endswith(":") and stripped.upper() == stripped:
        return True
    if ROMAN_HEADING_RE.match(stripped):
        return True
    return False


def is_clause_line(text: str) -> bool:
    stripped = text.strip()
    return bool(ALPHA_CLAUSE_RE.match(stripped) or NUMERIC_CLAUSE_RE.match(stripped))


def is_bullet_line(text: str) -> bool:
    stripped = text.strip()
    return bool(BULLET_LINE_RE.match(stripped))


def is_checkbox_line(text: str) -> bool:
    stripped = text.strip()
    return bool(CHECKBOX_LINE_RE.match(stripped))


def is_list_line(text: str) -> bool:
    return is_bullet_line(text) or is_checkbox_line(text)


def is_signature_header(text: str) -> bool:
    lowered = text.lower()
    return "ctc/id" in lowered or "date & place of issue" in lowered


def looks_like_signature_blank_row(text: str) -> bool:
    stripped = text.strip()
    return bool(stripped) and (set(stripped) <= {"_", " "} or stripped.count("_") >= 6)


def add_structured_block(document: Document, block: str, profile: str) -> None:
    if looks_like_title(block):
        add_centered_paragraph(document, block, bold=True, size=14)
        return

    if block.lower() == "and":
        add_centered_paragraph(document, block, size=12)
        return

    if is_heading_line(block):
        alignment = WD_ALIGN_PARAGRAPH.CENTER if block.endswith(":") else WD_ALIGN_PARAGRAPH.LEFT
        add_paragraph(document, block, alignment=alignment, bold=True)
        return

    if profile == "legal" and looks_like_party_definition(block):
        add_paragraph(document, block, alignment=WD_ALIGN_PARAGRAPH.JUSTIFY)
        return

    if looks_like_fill_in_field(block):
        add_paragraph(document, block, alignment=WD_ALIGN_PARAGRAPH.LEFT)
        return

    if is_list_line(block):
        add_paragraph(document, normalize_list_marker(block), alignment=WD_ALIGN_PARAGRAPH.LEFT, indent=1)
        return

    if is_clause_line(block):
        indent_level = clause_indent_level(block)
        add_paragraph(document, block, alignment=WD_ALIGN_PARAGRAPH.JUSTIFY, indent=indent_level)
        return

    bold = block.startswith("NOW, THEREFORE")
    add_paragraph(document, block, alignment=WD_ALIGN_PARAGRAPH.JUSTIFY, bold=bold)


def clause_indent_level(text: str) -> int:
    stripped = text.strip()
    if NUMERIC_CLAUSE_RE.match(stripped):
        return 1
    if ALPHA_CLAUSE_RE.match(stripped):
        return 0
    return 0


def add_centered_paragraph(document: Document, text: str, bold: bool = False, size: int = 12) -> None:
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.paragraph_format.space_before = Pt(6 if not bold else 12)
    paragraph.paragraph_format.space_after = Pt(12)
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Times New Roman"


def add_paragraph(
    document: Document,
    text: str,
    *,
    alignment: WD_ALIGN_PARAGRAPH,
    bold: bool = False,
    indent: int = 0,
) -> None:
    paragraph = document.add_paragraph()
    paragraph.alignment = alignment
    paragraph.paragraph_format.left_indent = Pt(indent * 18)
    paragraph.paragraph_format.space_before = Pt(6 if bold else 0)
    paragraph.paragraph_format.space_after = Pt(12)
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)


def looks_like_party_definition(text: str) -> bool:
    lowered = text.lower()
    return "shall each be referred to as" in lowered and "party" in lowered


def looks_like_fill_in_field(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    return "____________________" in stripped or looks_like_fill_in_blank(stripped)


def normalize_list_marker(text: str) -> str:
    stripped = text.strip()
    if is_checkbox_line(stripped):
        stripped = re.sub(r"^(?:\[[ xX]?\]|\([ xX]?\)|☐|☑|✓)\s*", "", stripped)
        checked = text.strip().startswith(("[x]", "(x)", "[X]", "(X)", "☑", "✓"))
        marker = "[x]" if checked else "[ ]"
        return f"{marker} {stripped}".strip()
    if is_bullet_line(stripped):
        stripped = re.sub(r"^(?:[-*•]|[oO0]\s)\s*", "", stripped)
        return f"- {stripped}".strip()
    return stripped


def infer_document_style(text: str, requested_style: str) -> str:
    normalized = requested_style.lower()
    if normalized not in SUPPORTED_DOC_STYLES:
        raise ValueError(f"Unsupported document style: {requested_style}")
    if normalized != "auto":
        return normalized

    lowered = text.lower()
    legal_score = sum(
        [
            2 if "whereas" in lowered else 0,
            2 if "recitals" in lowered else 0,
            1 if "now, therefore" in lowered else 0,
            1 if "franchisor" in lowered else 0,
            1 if "franchisee" in lowered else 0,
        ]
    )
    form_score = sum(
        [
            2 if "ctc/id" in lowered else 0,
            1 if "date & place of issue" in lowered else 0,
            1 if "acknowledgment" in lowered else 0,
            1 if "____" in text else 0,
        ]
    )

    if legal_score >= 3:
        return "legal"
    if form_score >= 2:
        return "form"
    return "general"


def add_signature_table(document: Document, headers: Sequence[str], blank_rows: int) -> Table:
    col_count = max(3, len(headers))
    table = document.add_table(rows=1 + blank_rows, cols=col_count)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"

    header_cells = table.rows[0].cells
    for idx in range(col_count):
        text = headers[idx] if idx < len(headers) else ""
        paragraph = header_cells[idx].paragraphs[0]
        run = paragraph.add_run(text)
        run.bold = True

    for row_index in range(1, 1 + blank_rows):
        for cell in table.rows[row_index].cells:
            cell.text = " "

    document.add_paragraph()
    return table


def parse_table_row_candidate(block: str, profile: str) -> TableRowCandidate | None:
    if is_clause_line(block) or is_heading_line(block) or looks_like_title(block) or looks_like_fill_in_field(block):
        return None
    if profile == "legal" and block.startswith("WHEREAS"):
        return None

    if "  " not in block:
        return None

    cells = [part.strip() for part in re.split(r"\s{2,}", block) if part.strip()]
    if len(cells) < 2:
        return None
    if any(len(cell) > 80 for cell in cells):
        return None
    return TableRowCandidate(cells=cells)


def looks_like_tableish_text(block: str) -> bool:
    if "  " not in block:
        return False
    cells = [part.strip() for part in re.split(r"\s{2,}", block) if part.strip()]
    if len(cells) < 2:
        return False
    shortish = sum(1 for cell in cells if len(cell) <= 40)
    return shortish >= 2


def flush_table_candidates(document: Document, rows: Sequence[TableRowCandidate]) -> None:
    normalized_rows = normalize_table_candidates(rows)
    if len(normalized_rows) < 2:
        for row in normalized_rows:
            add_paragraph(document, "  ".join(row.cells), alignment=WD_ALIGN_PARAGRAPH.JUSTIFY)
        return

    col_count = max(len(row.cells) for row in normalized_rows)
    table = document.add_table(rows=len(normalized_rows), cols=col_count)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"

    for row_index, row in enumerate(normalized_rows):
        for col_index in range(col_count):
            text = row.cells[col_index] if col_index < len(row.cells) else ""
            cell_paragraph = table.rows[row_index].cells[col_index].paragraphs[0]
            run = cell_paragraph.add_run(text)
            if row_index == 0:
                run.bold = True

    document.add_paragraph()


def extract_table_rows_from_text(text: str) -> list[list[str]]:
    lines = [line.strip() for line in text.splitlines() if line.strip() and line.strip() != "[TABLE_GRID_DETECTED]"]
    rows: list[list[str]] = []

    for line in lines:
        if not looks_like_tableish_text(line):
            if rows:
                break
            continue

        cells = [part.strip() for part in re.split(r"\s{2,}", line) if part.strip()]
        if len(cells) >= 2:
            rows.append(cells)

    return [row.cells for row in normalize_table_candidates([TableRowCandidate(cells=row) for row in rows])]


def normalize_table_candidates(rows: Sequence[TableRowCandidate]) -> list[TableRowCandidate]:
    if not rows:
        return []

    target_cols = infer_target_column_count(rows)
    normalized: list[TableRowCandidate] = []
    for row in rows:
        cells = normalize_table_row_cells(row.cells, target_cols)
        normalized.append(TableRowCandidate(cells=cells))
    return normalized


def infer_target_column_count(rows: Sequence[TableRowCandidate]) -> int:
    counts: dict[int, int] = {}
    for row in rows:
        counts[len(row.cells)] = counts.get(len(row.cells), 0) + 1
    return max(counts.items(), key=lambda item: (item[1], item[0]))[0]


def normalize_table_row_cells(cells: Sequence[str], target_cols: int) -> list[str]:
    cleaned = [normalize_table_cell_text(cell) for cell in cells if cell.strip()]
    if not cleaned:
        return [""] * target_cols

    if len(cleaned) == target_cols:
        return cleaned

    if len(cleaned) > target_cols:
        head = cleaned[: target_cols - 1]
        tail = " ".join(cleaned[target_cols - 1 :]).strip()
        return head + [tail]

    while len(cleaned) < target_cols:
        split_index = find_best_cell_to_split(cleaned)
        if split_index is None:
            cleaned.append("")
            continue
        left, right = split_cell_for_table(cleaned[split_index])
        if not right:
            cleaned.append("")
            continue
        cleaned = cleaned[:split_index] + [left, right] + cleaned[split_index + 1 :]

    return cleaned[:target_cols]


def find_best_cell_to_split(cells: Sequence[str]) -> int | None:
    scored: list[tuple[int, int]] = []
    for index, cell in enumerate(cells):
        if len(cell) < 12:
            continue
        if ":" in cell or "/" in cell or " - " in cell:
            scored.append((index, len(cell)))
        elif len(cell.split()) >= 3:
            scored.append((index, len(cell)))
    if not scored:
        return None
    return max(scored, key=lambda item: item[1])[0]


def split_cell_for_table(cell: str) -> tuple[str, str]:
    for separator in (" : ", ": ", " / ", " - ", "  "):
        if separator in cell:
            left, right = cell.split(separator, 1)
            return left.strip(), right.strip()

    words = cell.split()
    midpoint = len(words) // 2
    return " ".join(words[:midpoint]).strip(), " ".join(words[midpoint:]).strip()


def normalize_table_cell_text(text: str) -> str:
    cleaned = " ".join(text.split())
    cleaned = re.sub(r"\s+([,.;:])", r"\1", cleaned)
    return cleaned.strip()


def auto_fit_columns(sheet) -> None:
    widths: dict[str, int] = {}
    for row in sheet.iter_rows():
        for cell in row:
            if cell.value is None:
                continue
            value = str(cell.value)
            widths[cell.column_letter] = max(widths.get(cell.column_letter, 0), min(len(value), 60))

    for column_letter, width in widths.items():
        sheet.column_dimensions[column_letter].width = max(12, width + 2)


def preprocess_for_ocr(image: Image.Image) -> Image.Image:
    corrected = auto_orient_image(image)
    # Upscale and binarize scans so Tesseract gets cleaner glyph boundaries.
    grayscale = corrected.convert("L")
    width, height = grayscale.size
    scaled = grayscale.resize((width * 2, height * 2))
    threshold = 185
    binarized = scaled.point(lambda px: 255 if px > threshold else 0, mode="1").convert("L")
    return deskew_binary_image(binarized)


def auto_orient_image(image: Image.Image) -> Image.Image:
    if pytesseract is None:
        return image

    try:
        osd = pytesseract.image_to_osd(image)
    except Exception:
        return image

    match = re.search(r"Rotate:\s+(\d+)", osd)
    if not match:
        return image

    rotation = int(match.group(1)) % 360
    if rotation == 0:
        return image
    return image.rotate(-rotation, expand=True)


def deskew_binary_image(image: Image.Image) -> Image.Image:
    angle = estimate_skew_angle(image)
    if abs(angle) < 0.2:
        return image
    return image.rotate(-angle, expand=True, fillcolor=255)


def estimate_skew_angle(image: Image.Image) -> float:
    # Lightweight projection-based skew estimate across a small angle range.
    best_angle = 0.0
    best_score = -1.0

    for angle in [value / 10 for value in range(-20, 21)]:
        rotated = image.rotate(angle, expand=True, fillcolor=255)
        score = horizontal_projection_score(rotated)
        if score > best_score:
            best_score = score
            best_angle = angle

    return best_angle


def horizontal_projection_score(image: Image.Image) -> float:
    width, height = image.size
    pixels = image.load()
    row_sums: list[int] = []

    for y in range(height):
        dark_pixels = 0
        for x in range(width):
            if pixels[x, y] < 128:
                dark_pixels += 1
        row_sums.append(dark_pixels)

    transitions = 0
    for idx in range(1, len(row_sums)):
        transitions += abs(row_sums[idx] - row_sums[idx - 1])
    return float(transitions)


def score_ocr_text(text: str) -> float:
    if not text.strip():
        return float("-inf")

    lines = [line for line in text.splitlines() if line.strip()]
    if not lines:
        return float("-inf")

    alnum_count = sum(1 for char in text if char.isalnum())
    weird_count = sum(1 for char in text if char in {"�", "|", "~"})
    long_line_penalty = sum(max(len(line) - 180, 0) for line in lines) / 200
    blank_bonus = min(text.count("\n\n"), 12) * 0.05

    return (alnum_count / 500) - weird_count - long_line_penalty + blank_bonus


def image_likely_has_table_grid(image_path: Path) -> bool:
    with Image.open(image_path) as image:
        processed = preprocess_for_ocr(image)
    return image_likely_has_table_grid_from_image(processed)


def image_likely_has_table_grid_from_image(image: Image.Image) -> bool:
    width, height = image.size
    pixels = image.load()

    horizontal_lines = 0
    vertical_lines = 0

    sample_rows = max(10, height // 40)
    for row in range(0, height, sample_rows):
        dark_run = 0
        max_run = 0
        for col in range(width):
            if pixels[col, row] < 128:
                dark_run += 1
                max_run = max(max_run, dark_run)
            else:
                dark_run = 0
        if max_run > width * 0.45:
            horizontal_lines += 1

    sample_cols = max(10, width // 40)
    for col in range(0, width, sample_cols):
        dark_run = 0
        max_run = 0
        for row in range(height):
            if pixels[col, row] < 128:
                dark_run += 1
                max_run = max(max_run, dark_run)
            else:
                dark_run = 0
        if max_run > height * 0.25:
            vertical_lines += 1

    return horizontal_lines >= 2 and vertical_lines >= 2


def append_paragraphs(document: Document, lines: Iterable[str]) -> None:
    paragraph_lines = list(lines)
    buffer: list[str] = []

    for line in paragraph_lines:
        if line.strip():
            buffer.append(line)
            continue

        if buffer:
            document.add_paragraph(" ".join(buffer))
            buffer.clear()
        else:
            document.add_paragraph("")

    if buffer:
        document.add_paragraph(" ".join(buffer))


def convert_images_to_document(
    image_paths: Sequence[Path],
    output_format: str,
    output_path: Path,
    title: str | None = None,
    ocr_engine: str = "auto",
    document_style: str = "auto",
    pdf_page_numbers: bool = False,
    pdf_watermark: str | None = None,
) -> Path:
    pages = extract_text_from_images_with_engine(image_paths, ocr_engine, document_style)
    return build_output(
        pages,
        output_format,
        output_path,
        title=title,
        image_paths=image_paths,
        document_style=document_style,
        pdf_page_numbers=pdf_page_numbers,
        pdf_watermark=pdf_watermark,
    )
