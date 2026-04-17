"""
ConvertFlow — Smoke Test Suite (Section A of QA Checklist)

Run from the project root:
    .venv/Scripts/python test_all_tools.py

Creates all fixtures in memory or in a temp dir, calls each tool function
directly, and checks that the output exists and is non-zero bytes.

Skips tests that require external services (Ollama, LibreOffice, etc.) if
those services are unavailable.
"""
from __future__ import annotations

import io
import base64
import sys as _sys
if hasattr(_sys.stdout, "reconfigure"):
    _sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    _sys.stderr.reconfigure(encoding="utf-8", errors="replace")
import os
import sys
import time
import shutil
import tempfile
import traceback
from pathlib import Path

# ── Make sure the project root is on sys.path ─────────────────────────────────
ROOT = Path(__file__).parent
sys.path.insert(0, str(ROOT))

from execution import pdf_tools  # noqa: E402

# ── Colours ───────────────────────────────────────────────────────────────────
GREEN  = "\033[92m"
RED    = "\033[91m"
YELLOW = "\033[93m"
RESET  = "\033[0m"
BOLD   = "\033[1m"

PASS = f"{GREEN}PASS{RESET}"
FAIL = f"{RED}FAIL{RESET}"
SKIP = f"{YELLOW}SKIP{RESET}"

results: list[tuple[str, str, float, str]] = []   # (name, status, seconds, msg)


def run(name: str, fn, *args, skip_reason: str = ""):
    """Run fn(*args), record result."""
    if skip_reason:
        results.append((name, "SKIP", 0.0, skip_reason))
        print(f"  {SKIP}  {name}  — {skip_reason}")
        return

    t0 = time.perf_counter()
    try:
        fn(*args)
        elapsed = time.perf_counter() - t0
        results.append((name, "PASS", elapsed, ""))
        print(f"  {PASS}  {name}  ({elapsed:.2f}s)")
    except Exception as exc:
        elapsed = time.perf_counter() - t0
        msg = str(exc)
        results.append((name, "FAIL", elapsed, msg))
        print(f"  {FAIL}  {name}  ({elapsed:.2f}s)")
        print(f"       {RED}{msg[:200]}{RESET}")
        if "--verbose" in sys.argv or "-v" in sys.argv:
            traceback.print_exc()


def assert_nonempty(path: Path):
    """Raise if output file is missing or zero bytes."""
    if not path.exists():
        raise FileNotFoundError(f"Output not created: {path}")
    size = path.stat().st_size
    if size == 0:
        raise ValueError(f"Output is zero bytes: {path}")


# ── Fixtures ──────────────────────────────────────────────────────────────────

def make_pdf_3pages(path: Path, text_prefix: str = "Page") -> Path:
    """Create a simple 3-page PDF with text using PyMuPDF."""
    import fitz
    doc = fitz.open()
    for i in range(1, 4):
        page = doc.new_page(width=595, height=842)   # A4 in points
        page.insert_text((72, 100), f"{text_prefix} {i}", fontsize=24)
        page.insert_text((72, 140), f"Sample content for page {i}.", fontsize=12)
        if i == 1:
            page.insert_text((72, 180), "Invoice #12345", fontsize=12)
            page.insert_text((72, 200), "Hello World", fontsize=12)
    doc.save(str(path))
    return path


def make_pdf_mixed_layout(path: Path) -> Path:
    """Create a one-page PDF containing both narrative text and a simple table."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)

    page.insert_text((72, 90), "Quarterly Summary", fontsize=22)
    page.insert_text((72, 125), "This report mixes narrative text with a tabular section.", fontsize=11)
    page.insert_text((72, 145), "The workbook output should keep both kinds of content.", fontsize=11)

    left = 72
    top = 220
    col_widths = [150, 110, 110]
    row_height = 28
    rows = [
        ["Item", "Qty", "Price"],
        ["Widget A", "2", "19.99"],
        ["Widget B", "5", "7.50"],
    ]

    x = left
    for width in col_widths:
        page.draw_line((x, top), (x, top + row_height * len(rows)), color=(0, 0, 0), width=1)
        x += width
    page.draw_line((x, top), (x, top + row_height * len(rows)), color=(0, 0, 0), width=1)
    for row_index in range(len(rows) + 1):
        y = top + row_index * row_height
        page.draw_line((left, y), (left + sum(col_widths), y), color=(0, 0, 0), width=1)

    for row_index, row in enumerate(rows):
        y = top + 18 + row_index * row_height
        x = left + 8
        for col_index, value in enumerate(row):
            page.insert_text((x, y), value, fontsize=11)
            x += col_widths[col_index]

    page.insert_text((72, 360), "Notes", fontsize=14)
    page.insert_text((72, 385), "Widget A outperformed expectations.", fontsize=11)
    page.insert_text((72, 405), "Widget B had the highest unit volume.", fontsize=11)

    doc.save(str(path))
    return path


def make_pdf_two_columns(path: Path) -> Path:
    """Create a one-page PDF with left and right columns to test reading order."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)

    page.insert_text((72, 80), "Two Column Example", fontsize=20)

    left_lines = [
        "Left A",
        "Left B",
        "Left C",
    ]
    right_lines = [
        "Right A",
        "Right B",
        "Right C",
    ]

    y = 130
    for line in left_lines:
        page.insert_text((72, y), line, fontsize=11)
        y += 22

    y = 130
    for line in right_lines:
        page.insert_text((330, y), line, fontsize=11)
        y += 22

    doc.save(str(path))
    return path


def make_pdf_small_card(path: Path, pages: int = 2) -> Path:
    """Create a small-format PDF to exercise card/layout heuristics."""
    import fitz

    doc = fitz.open()
    for index in range(pages):
        page = doc.new_page(width=252, height=144)
        label = "FRONT SIDE" if index == 0 else f"PAGE {index + 1}"
        page.insert_text((24, 54), label, fontsize=16)
        page.insert_text((24, 86), "Business card layout sample", fontsize=10)
    doc.save(str(path))
    return path


def make_pdf_typed_table(path: Path) -> Path:
    """Create a table PDF with date, amount, and percent-like values."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 90), "Typed Table Example", fontsize=20)

    left = 72
    top = 160
    col_widths = [130, 120, 120, 90]
    row_height = 28
    rows = [
        ["Date", "Amount", "Margin", "Count"],
        ["02-Nov-22", "12345.67", "12.5%", "42"],
        ["2026-04-10", "(1,250.00)", "3.0%", "7"],
    ]

    x = left
    for width in col_widths:
        page.draw_line((x, top), (x, top + row_height * len(rows)), color=(0, 0, 0), width=1)
        x += width
    page.draw_line((x, top), (x, top + row_height * len(rows)), color=(0, 0, 0), width=1)
    for row_index in range(len(rows) + 1):
        y = top + row_index * row_height
        page.draw_line((left, y), (left + sum(col_widths), y), color=(0, 0, 0), width=1)

    for row_index, row in enumerate(rows):
        y = top + 18 + row_index * row_height
        x = left + 8
        for col_index, value in enumerate(row):
            page.insert_text((x, y), value, fontsize=11)
            x += col_widths[col_index]

    doc.save(str(path))
    return path


def make_pdf_many_pages(path: Path, count: int = 24) -> Path:
    """Create a longer PDF to exercise adaptive fidelity render settings."""
    import fitz

    doc = fitz.open()
    for i in range(1, count + 1):
        page = doc.new_page(width=595, height=842)
        page.insert_text((72, 90), f"Long Report Page {i}", fontsize=18)
        for line_index in range(18):
            page.insert_text(
                (72, 130 + line_index * 24),
                f"Paragraph {line_index + 1} on page {i} with enough text to simulate a report export.",
                fontsize=11,
            )
    doc.save(str(path))
    return path


def make_pdf_text_only_table(path: Path) -> Path:
    """Create a text-only table layout without drawn borders."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 90), "Project Register", fontsize=18)
    rows = [
        ("Item", "Project Title", "Date", "Owner"),
        ("1", "Alpha Build", "2025", "Client A"),
        ("2", "Beta Plant", "2026", "Client B"),
        ("3", "Gamma Works", "2027", "Client C"),
    ]
    y = 150
    x_positions = [72, 130, 330, 430]
    for row in rows:
        for x, value in zip(x_positions, row):
            page.insert_text((x, y), value, fontsize=12)
        y += 28
    doc.save(str(path))
    return path


def make_pdf_wrapped_text_table(path: Path) -> Path:
    """Create a text-only table where logical rows wrap across multiple PDF lines."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 90), "Wrapped Project Register", fontsize=18)
    x_positions = [72, 130, 330, 430, 470]
    rows = [
        [("Item", 0), ("Project Title", 1), ("Date", 2), ("Owner", 3), ("Address", 4)],
        [("31", 0), ("Hoist Alarm and", 1), ("02-Nov-22", 2), ("Mitsubishi", 3), ("Autopark,", 4)],
        [("Signal Light", 1), ("Motors Phils.", 3), ("Greenfield PEZA,", 4)],
        [("Installation", 1), ("Corp.", 3), ("Sta Rosa City,", 4)],
        [("32", 0), ("Power Piping Banding", 1), ("11-Sep-22", 2), ("Honda", 3), ("Laguna", 4)],
    ]
    y = 150
    for row in rows:
        for value, column_index in row:
            page.insert_text((x_positions[column_index], y), value, fontsize=12)
        y += 28
    doc.save(str(path))
    return path


def make_png(path: Path, width: int = 200, height: int = 60) -> Path:
    """Create a small PNG image using PIL."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, width - 10, height - 10], outline=(0, 0, 200), width=2)
    draw.line([(20, height // 2), (width - 20, height // 2)], fill=(0, 0, 200), width=3)
    img.save(str(path))
    return path


def make_signature_png() -> bytes:
    """Return PNG bytes suitable for sign_pdf."""
    from PIL import Image, ImageDraw
    img = Image.new("RGBA", (600, 180), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    pts = [(50 + i * 10, 90 + int(40 * __import__("math").sin(i * 0.5))) for i in range(51)]
    draw.line(pts, fill=(0, 0, 150, 255), width=3)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def make_docx(path: Path) -> Path:
    """Create a minimal .docx file."""
    from docx import Document
    doc = Document()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a sample paragraph for testing.")
    doc.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    """Create a minimal .xlsx file."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"])
    ws.append(["Bob", 82, "B"])
    wb.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    """Create a minimal .pptx file."""
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "Sample slide content"
    prs.save(str(path))
    return path


def make_html(path: Path) -> Path:
    path.write_text(
        "<html><body><h1>Hello PDF</h1><p>Sample HTML page for testing.</p></body></html>",
        encoding="utf-8",
    )
    return path


def _ollama_available() -> bool:
    import json
    import urllib.error
    import urllib.request

    host = os.environ.get("OLLAMA_HOST", "http://localhost:11434").rstrip("/")
    model = os.environ.get("OLLAMA_OCR_MODEL", "gemma4:e4b")
    wait_seconds = max(0.0, float(os.environ.get("OLLAMA_WAIT_FOR_MODEL_SECONDS", "0") or "0"))
    deadline = time.time() + wait_seconds

    def _model_present() -> bool:
        with urllib.request.urlopen(host + "/api/tags", timeout=2) as response:
            payload = json.loads(response.read().decode("utf-8"))
        available = [entry.get("name", "") for entry in payload.get("models", [])]
        return any(model in name for name in available)

    try:
        while True:
            if _model_present():
                return True
            if time.time() >= deadline:
                return False
            time.sleep(2)
    except urllib.error.URLError:
        return False
    except Exception:
        return False


def _libreoffice_available() -> bool:
    return shutil.which("soffice") is not None or shutil.which("libreoffice") is not None


def _office_available() -> bool:
    """Return True if Microsoft Office COM automation is available (pywin32 + registry)."""
    try:
        import winreg
        for app in ("POWERPNT.EXE", "WINWORD.EXE", "EXCEL.EXE"):
            try:
                winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                    r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\\" + app)
                return True
            except FileNotFoundError:
                continue
    except Exception:
        pass
    return _libreoffice_available()


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    tmp = Path(tempfile.mkdtemp(prefix="cf_test_"))
    print(f"\n{BOLD}ConvertFlow — Smoke Test{RESET}")
    print(f"Temp dir: {tmp}\n")

    ollama_ok   = _ollama_available()
    libre_ok    = _office_available()

    # =========================================================================
    # SECTION A — Smoke Tests (all tools return output)
    # =========================================================================
    print(f"{BOLD}Section A — Smoke Tests{RESET}\n")

    # ── Build core fixtures ───────────────────────────────────────────────────
    pdf3   = make_pdf_3pages(tmp / "sample3.pdf")
    pdf_a  = make_pdf_3pages(tmp / "doc_a.pdf", text_prefix="Alpha")
    pdf_b  = make_pdf_3pages(tmp / "doc_b.pdf", text_prefix="Beta")
    img1   = make_png(tmp / "img1.png")
    img2   = make_png(tmp / "img2.png", 300, 200)
    sig_bytes = make_signature_png()

    # ── 1. Merge PDF ──────────────────────────────────────────────────────────
    def t_merge():
        out = tmp / "merged.pdf"
        pdf_tools.merge_pdfs([pdf_a, pdf_b], out)
        assert_nonempty(out)
    run("merge_pdf", t_merge)

    # ── 2. Split PDF ──────────────────────────────────────────────────────────
    def t_split():
        out_dir = tmp / "split"
        out_dir.mkdir(exist_ok=True)
        pages = pdf_tools.split_pdf(pdf3, out_dir)
        if len(pages) != 3:
            raise ValueError(f"Expected 3 files, got {len(pages)}")
        for p in pages:
            assert_nonempty(p)
    run("split_pdf", t_split)

    # ── 3. Extract Pages ──────────────────────────────────────────────────────
    def t_extract():
        out = tmp / "extracted.pdf"
        pdf_tools.extract_pages(pdf3, "1,3", out)
        assert_nonempty(out)
    run("extract_pages", t_extract)

    # ── 4. Remove Pages ───────────────────────────────────────────────────────
    def t_remove():
        out = tmp / "removed.pdf"
        pdf_tools.remove_pages(pdf3, "2", out)
        assert_nonempty(out)
    run("remove_pages", t_remove)

    # ── 5. Rotate PDF ─────────────────────────────────────────────────────────
    def t_rotate():
        out = tmp / "rotated.pdf"
        pdf_tools.rotate_pdf(pdf3, 90, out)
        assert_nonempty(out)
    run("rotate_pdf", t_rotate)

    # ── 6. Compress PDF ───────────────────────────────────────────────────────
    def t_compress():
        out = tmp / "compressed.pdf"
        pdf_tools.compress_pdf(pdf3, out)
        assert_nonempty(out)
    run("compress_pdf", t_compress)

    # ── 7. Repair PDF ─────────────────────────────────────────────────────────
    def t_repair():
        out = tmp / "repaired.pdf"
        pdf_tools.repair_pdf(pdf3, out)
        assert_nonempty(out)
    run("repair_pdf", t_repair)

    # ── 8. Images to PDF ──────────────────────────────────────────────────────
    def t_images_to_pdf():
        out = tmp / "from_images.pdf"
        pdf_tools.images_to_pdf([img1, img2], out)
        assert_nonempty(out)
    run("images_to_pdf", t_images_to_pdf)

    # ── 9. Word to PDF ────────────────────────────────────────────────────────
    def t_word_to_pdf():
        docx_path = make_docx(tmp / "test.docx")
        out = tmp / "from_word.pdf"
        pdf_tools.word_to_pdf(docx_path, out)
        assert_nonempty(out)
    run("word_to_pdf", t_word_to_pdf,
        skip_reason="" if libre_ok else "LibreOffice not found")

    # ── 10. PPTX to PDF ───────────────────────────────────────────────────────
    def t_pptx_to_pdf():
        pptx_path = make_pptx(tmp / "test.pptx")
        out = tmp / "from_pptx.pdf"
        pdf_tools.pptx_to_pdf(pptx_path, out)
        assert_nonempty(out)
    run("pptx_to_pdf", t_pptx_to_pdf,
        skip_reason="" if libre_ok else "LibreOffice not found")

    # ── 11. Excel to PDF ──────────────────────────────────────────────────────
    def t_excel_to_pdf():
        xlsx_path = make_xlsx(tmp / "test.xlsx")
        out = tmp / "from_excel.pdf"
        pdf_tools.excel_to_pdf(xlsx_path, out)
        assert_nonempty(out)
    run("excel_to_pdf", t_excel_to_pdf,
        skip_reason="" if libre_ok else "LibreOffice not found")

    # ── 12. HTML to PDF ───────────────────────────────────────────────────────
    def t_html_to_pdf():
        html_path = make_html(tmp / "test.html")
        out = tmp / "from_html.pdf"
        pdf_tools.html_to_pdf(html_path, out)
        assert_nonempty(out)
    run("html_to_pdf", t_html_to_pdf)

    # ── 13. PDF to Word (basic) ───────────────────────────────────────────────
    def t_pdf_to_word_basic():
        out = tmp / "to_word_basic.docx"
        pdf_tools.pdf_to_word(pdf3, out, engine="basic")
        assert_nonempty(out)
    run("pdf_to_word (basic)", t_pdf_to_word_basic)

    def t_pdf_to_word_fidelity():
        import zipfile
        out = tmp / "to_word_fidelity.docx"
        pdf_tools.pdf_to_word(pdf3, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            media = [name for name in zf.namelist() if name.startswith("word/media/")]
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
        if len(media) != 3:
            raise ValueError(f"Expected 3 embedded page images, got {len(media)}")
        if "txbxContent" not in doc_xml and "w:pict" not in doc_xml:
            raise ValueError("Expected Word fidelity output to contain textbox overlay markup")
    run("pdf_to_word (fidelity)", t_pdf_to_word_fidelity)

    def t_fidelity_render_settings():
        import fitz

        long_pdf = make_pdf_many_pages(tmp / "many_pages.pdf", count=24)
        with fitz.open(str(long_pdf)) as long_doc:
            long_settings = pdf_tools._fidelity_render_settings(long_doc)
        if str(long_settings["image_format"]).lower() not in {"jpg", "jpeg"}:
            raise ValueError("Expected long documents to use JPEG-backed fidelity rendering")
        if int(long_settings["dpi"]) > 110:
            raise ValueError("Expected long documents to reduce fidelity render DPI")

        card_pdf = tmp / "small_card.pdf"
        card_doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE"):
            page = card_doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        card_doc.save(str(card_pdf))
        with fitz.open(str(card_pdf)) as small_doc:
            small_settings = pdf_tools._fidelity_render_settings(small_doc)
        if str(small_settings["image_format"]).lower() != "png":
            raise ValueError("Expected small-format PDFs to keep PNG fidelity rendering")
        if int(small_settings["dpi"]) != 144:
            raise ValueError("Expected small-format PDFs to keep sharp 144 DPI rendering")
    run("fidelity render settings", t_fidelity_render_settings)

    def t_pdf_document_profile_guidance():
        regular = pdf_tools.pdf_document_profile(pdf3)
        if regular.get("best_target") != "word":
            raise ValueError(f"Expected regular documents to prefer Word, got {regular.get('best_target')!r}")
        if regular.get("tool_fit", {}).get("word", {}).get("label") != "best fit":
            raise ValueError("Expected regular document profile to mark Word as the best fit")

        small_card = make_pdf_small_card(tmp / "profile_small_card.pdf")
        small = pdf_tools.pdf_document_profile(small_card)
        if small.get("profile") != "small-format layout":
            raise ValueError(f"Expected small card PDF to classify as small-format layout, got {small.get('profile')!r}")
        if small.get("best_target") != "svg":
            raise ValueError(f"Expected small-format layout to prefer SVG, got {small.get('best_target')!r}")
        if small.get("tool_fit", {}).get("pptx", {}).get("label") != "good fit":
            raise ValueError("Expected small-format layout guidance to describe PowerPoint as a good fit")
    run("pdf document profile guidance", t_pdf_document_profile_guidance)

    def t_pdf_to_word_fidelity_two_columns():
        import zipfile
        cols_pdf = make_pdf_two_columns(tmp / "two_columns_word.pdf")
        out = tmp / "to_word_fidelity_columns.docx"
        pdf_tools.pdf_to_word(cols_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
        required = ["Left A", "Left B", "Left C", "Right A", "Right B", "Right C"]
        missing = [value for value in required if f"<w:t>{value}</w:t>" not in doc_xml]
        if missing:
            raise ValueError(f"Two-column Word fidelity output is missing: {missing}")
        left_a = doc_xml.index("<w:t>Left A</w:t>")
        left_c = doc_xml.index("<w:t>Left C</w:t>")
        right_a = doc_xml.index("<w:t>Right A</w:t>")
        if not (left_a < left_c < right_a):
            raise ValueError(
                f"Unexpected two-column Word reading order: Left A idx={left_a}, Left C idx={left_c}, Right A idx={right_a}"
            )
    run("pdf_to_word (fidelity, two columns)", t_pdf_to_word_fidelity_two_columns)

    def t_pdf_to_word_fidelity_editable_appendix():
        import zipfile

        long_pdf = make_pdf_many_pages(tmp / "word_appendix_long.pdf", count=10)
        out = tmp / "to_word_fidelity_appendix.docx"
        pdf_tools.pdf_to_word(long_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
        if "Editable Text Extract" not in doc_xml:
            raise ValueError("Expected long Word fidelity export to include an editable text appendix")
        if "Long Report Page 1" not in doc_xml or "Paragraph 1 on page 1" not in doc_xml:
            raise ValueError("Expected editable Word appendix to include extracted page text")
    run("pdf_to_word (fidelity, editable appendix)", t_pdf_to_word_fidelity_editable_appendix)

    def t_pdf_to_word_fidelity_mixed_layout_appendix():
        import zipfile

        mixed_pdf = make_pdf_mixed_layout(tmp / "mixed_layout_word.doc.pdf")
        out = tmp / "to_word_fidelity_mixed.docx"
        pdf_tools.pdf_to_word(mixed_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
        if "Editable Text Extract" not in doc_xml:
            raise ValueError("Expected mixed-layout Word fidelity output to include an editable appendix")
        if "Quarterly Summary" not in doc_xml:
            raise ValueError("Expected mixed-layout Word fidelity appendix to preserve page text")
    run("pdf_to_word (fidelity, mixed layout appendix)", t_pdf_to_word_fidelity_mixed_layout_appendix)

    # ── 13b. PDF to Word (ollama) ─────────────────────────────────────────────
    def t_pdf_to_word_ollama():
        out = tmp / "to_word_ollama.docx"
        pdf_tools.pdf_to_word(pdf3, out, engine="ollama")
        assert_nonempty(out)
    run("pdf_to_word (ollama)", t_pdf_to_word_ollama,
        skip_reason="" if ollama_ok else "Ollama offline")

    # ── 14. PDF to Text (basic) ───────────────────────────────────────────────
    def t_pdf_to_text_basic():
        out = tmp / "to_text_basic.txt"
        pdf_tools.pdf_to_text(pdf3, out, engine="basic")
        assert_nonempty(out)
    run("pdf_to_text (basic)", t_pdf_to_text_basic)

    # ── 14b. PDF to Text (ollama) ─────────────────────────────────────────────
    def t_pdf_to_text_ollama():
        out = tmp / "to_text_ollama.txt"
        pdf_tools.pdf_to_text(pdf3, out, engine="ollama")
        assert_nonempty(out)
    run("pdf_to_text (ollama)", t_pdf_to_text_ollama,
        skip_reason="" if ollama_ok else "Ollama offline")

    # ── 15. PDF to Excel (basic) ──────────────────────────────────────────────
    def t_pdf_to_excel_basic():
        out = tmp / "to_excel_basic.xlsx"
        pdf_tools.pdf_to_excel(pdf3, out, engine="basic")
        assert_nonempty(out)
    run("pdf_to_excel (basic)", t_pdf_to_excel_basic)

    def t_pdf_to_excel_fidelity():
        import openpyxl
        out = tmp / "to_excel_fidelity.xlsx"
        pdf_tools.pdf_to_excel(pdf3, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        if wb.sheetnames != ["Page 1", "Page 2", "Page 3"]:
            raise ValueError(f"Unexpected sheet names: {wb.sheetnames}")
        ws = wb["Page 1"]
        values = [ws.cell(r, 1).value for r in range(1, min(ws.max_row, 140) + 1) if ws.cell(r, 1).value]
        if "PDF page 1 snapshot" not in values:
            raise ValueError("Expected worksheet snapshot header in fidelity workbook")
        if "Extracted content" not in values:
            raise ValueError("Expected extracted content section in fidelity workbook")
        if "Invoice #12345" not in values:
            raise ValueError("Expected extracted page text in fidelity workbook")
    run("pdf_to_excel (fidelity)", t_pdf_to_excel_fidelity)

    # ── 15b. PDF to Excel (ollama) ────────────────────────────────────────────
    def t_pdf_to_excel_fidelity_mixed_layout():
        import openpyxl
        mixed_pdf = make_pdf_mixed_layout(tmp / "mixed_layout.pdf")
        out = tmp / "to_excel_fidelity_mixed.xlsx"
        pdf_tools.pdf_to_excel(mixed_pdf, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        ws = wb["Page 1"]
        values = [
            ws.cell(r, c).value
            for r in range(1, min(ws.max_row, 220) + 1)
            for c in range(1, min(ws.max_column, 8) + 1)
            if ws.cell(r, c).value
        ]
        expected = [
            "Extracted content",
            "Table 1",
            "Item",
            "Widget A",
            "Widget B",
            "Notes",
            "Widget A outperformed expectations.",
        ]
        missing = [value for value in expected if value not in values]
        if missing:
            raise ValueError(f"Mixed-layout Excel fidelity output is missing: {missing}")
    run("pdf_to_excel (fidelity, mixed layout)", t_pdf_to_excel_fidelity_mixed_layout)

    def t_pdf_to_excel_fidelity_text_only_table():
        import openpyxl

        text_table_pdf = make_pdf_wrapped_text_table(tmp / "text_only_table.pdf")
        out = tmp / "to_excel_fidelity_text_table.xlsx"
        pdf_tools.pdf_to_excel(text_table_pdf, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        ws = wb["Page 1"]
        values = [
            ws.cell(r, c).value
            for r in range(1, min(ws.max_row, 120) + 1)
            for c in range(1, min(ws.max_column, 6) + 1)
            if ws.cell(r, c).value
        ]
        required = [
            "Inferred table",
            "Item",
            "Project Title",
            "Hoist Alarm and Signal Light Installation",
            "Mitsubishi Motors Phils. Corp.",
            "Autopark, Greenfield PEZA, Sta Rosa City,",
        ]
        missing = [value for value in required if value not in values]
        if missing:
            raise ValueError(f"Text-only table Excel fidelity output is missing: {missing}")
    run("pdf_to_excel (fidelity, text-only table)", t_pdf_to_excel_fidelity_text_only_table)

    def t_pdf_to_excel_fidelity_typed_values():
        import openpyxl
        from datetime import datetime

        typed_pdf = make_pdf_typed_table(tmp / "typed_table.pdf")
        out = tmp / "to_excel_fidelity_typed.xlsx"
        pdf_tools.pdf_to_excel(typed_pdf, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        ws = wb["Page 1"]
        found_date = None
        found_amount = None
        found_margin = None
        found_count = None
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 180), max_col=min(ws.max_column, 6)):
            for cell in row:
                if isinstance(cell.value, datetime) and found_date is None:
                    found_date = cell
                if cell.value == 12345.67 and found_amount is None:
                    found_amount = cell
                if cell.value == 0.125 and found_margin is None:
                    found_margin = cell
                if cell.value == 42 and found_count is None:
                    found_count = cell
        if found_date is None:
            raise ValueError("Expected typed Excel output to include a real date cell")
        if found_amount is None or found_amount.number_format != "#,##0.00":
            raise ValueError("Expected typed Excel output to include a formatted numeric amount cell")
        if found_margin is None or found_margin.number_format != "0.00%":
            raise ValueError("Expected typed Excel output to include a formatted percent cell")
        if found_count is None or found_count.number_format != "#,##0":
            raise ValueError("Expected typed Excel output to include a formatted integer count cell")
    run("pdf_to_excel (fidelity, typed values)", t_pdf_to_excel_fidelity_typed_values)

    def t_blob_table_amount_hints():
        rows = [
            ["Item Project Title Contract Amount Date Completed\n39 Installation of Cooling System 6,675,295.00 01-Jul-19\n40 Construction of Oil Storage 2,899,907.58 01-Apr-19"]
        ]
        hints = pdf_tools._extract_blob_table_amount_hints(rows)
        if hints.get("39") != "6,675,295.00":
            raise ValueError(f"Expected amount hint for item 39, got {hints!r}")
        if hints.get("40") != "2,899,907.58":
            raise ValueError(f"Expected amount hint for item 40, got {hints!r}")
    run("pdf_to_excel (blob amount hints)", t_blob_table_amount_hints)

    def t_excel_recovered_cell_marker():
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        cell = ws["A1"]
        cell.value = 6675295
        pdf_tools._mark_excel_recovered_cell(cell, "Recovered from malformed blob table fallback for item 39.")
        if cell.comment is None or "item 39" not in cell.comment.text:
            raise ValueError("Expected recovered Excel cell to include a comment marker")
        if cell.fill.fgColor.rgb != "00FFF4CC":
            raise ValueError(f"Expected recovered Excel cell highlight fill, got {cell.fill.fgColor.rgb!r}")
    run("pdf_to_excel (recovered cell marker)", t_excel_recovered_cell_marker)

    def t_pdf_to_excel_fidelity_table_objects():
        import openpyxl

        mixed_pdf = make_pdf_mixed_layout(tmp / "mixed_layout_table_object.pdf")
        out = tmp / "to_excel_fidelity_table_object.xlsx"
        pdf_tools.pdf_to_excel(mixed_pdf, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        ws = wb["Page 1"]
        if not ws.tables:
            raise ValueError("Expected fidelity Excel output to create at least one real Excel table")
        if ws.freeze_panes is None:
            raise ValueError("Expected fidelity Excel output to freeze panes above the first extracted table")
    run("pdf_to_excel (fidelity, table objects)", t_pdf_to_excel_fidelity_table_objects)

    def t_pdf_to_excel_fidelity_two_columns():
        import openpyxl
        cols_pdf = make_pdf_two_columns(tmp / "two_columns.pdf")
        out = tmp / "to_excel_fidelity_columns.xlsx"
        pdf_tools.pdf_to_excel(cols_pdf, out, engine="fidelity")
        assert_nonempty(out)
        wb = openpyxl.load_workbook(out)
        ws = wb["Page 1"]
        values = [ws.cell(r, 1).value for r in range(1, min(ws.max_row, 220) + 1) if ws.cell(r, 1).value]
        required = ["Left A", "Left B", "Left C", "Right A", "Right B", "Right C"]
        missing = [value for value in required if value not in values]
        if missing:
            raise ValueError(f"Two-column Excel fidelity output is missing: {missing}")
        left_a = values.index("Left A")
        left_c = values.index("Left C")
        right_a = values.index("Right A")
        if not (left_a < left_c < right_a):
            raise ValueError(
                f"Unexpected two-column reading order: Left A idx={left_a}, Left C idx={left_c}, Right A idx={right_a}"
            )
    run("pdf_to_excel (fidelity, two columns)", t_pdf_to_excel_fidelity_two_columns)

    def t_pdf_to_excel_ollama():
        out = tmp / "to_excel_ollama.xlsx"
        pdf_tools.pdf_to_excel(pdf3, out, engine="ollama")
        assert_nonempty(out)
    run("pdf_to_excel (ollama)", t_pdf_to_excel_ollama,
        skip_reason="" if ollama_ok else "Ollama offline")

    # ── 16. PDF to Images ─────────────────────────────────────────────────────
    def t_pdf_to_images():
        out_dir = tmp / "pdf_images"
        out_dir.mkdir(exist_ok=True)
        imgs = pdf_tools.pdf_to_images(pdf3, out_dir, fmt="jpg")
        if len(imgs) != 3:
            raise ValueError(f"Expected 3 images, got {len(imgs)}")
        for p in imgs:
            assert_nonempty(p)
    run("pdf_to_images", t_pdf_to_images)

    # ── 17. PDF to PDF/A ──────────────────────────────────────────────────────
    def t_pdf_to_pdfa():
        out = tmp / "to_pdfa.pdf"
        pdf_tools.pdf_to_pdfa(pdf3, out)
        assert_nonempty(out)
    run("pdf_to_pdfa", t_pdf_to_pdfa)

    # ── 18. Add Page Numbers ──────────────────────────────────────────────────
    def t_add_page_numbers():
        out = tmp / "numbered.pdf"
        pdf_tools.add_page_numbers(pdf3, out, position="bottom-center")
        assert_nonempty(out)
    run("add_page_numbers", t_add_page_numbers)

    # ── 19. Watermark PDF ─────────────────────────────────────────────────────
    def t_watermark():
        out = tmp / "watermarked.pdf"
        pdf_tools.add_watermark(pdf3, "CONFIDENTIAL", out, position="center")
        assert_nonempty(out)
    run("add_watermark", t_watermark)

    # ── 20. Crop PDF ──────────────────────────────────────────────────────────
    def t_crop():
        out = tmp / "cropped.pdf"
        pdf_tools.crop_pdf(pdf3, out, left_mm=10, top_mm=10, right_mm=10, bottom_mm=10)
        assert_nonempty(out)
    run("crop_pdf", t_crop)

    # ── 21. Protect PDF ───────────────────────────────────────────────────────
    def t_protect():
        out = tmp / "protected.pdf"
        pdf_tools.protect_pdf(pdf3, "secret123", out)
        assert_nonempty(out)
    run("protect_pdf", t_protect)

    # ── 22. Unlock PDF ────────────────────────────────────────────────────────
    def t_unlock():
        # First protect it, then unlock
        protected = tmp / "to_unlock.pdf"
        pdf_tools.protect_pdf(pdf3, "secret123", protected)
        out = tmp / "unlocked.pdf"
        pdf_tools.unlock_pdf(protected, "secret123", out)
        assert_nonempty(out)
    run("unlock_pdf", t_unlock)

    # ── 23. AI Summarizer ─────────────────────────────────────────────────────
    def t_summarize():
        out = tmp / "summary.txt"
        pdf_tools.summarize_pdf(pdf3, out, engine="ollama", output_format="txt")
        assert_nonempty(out)
    run("summarize_pdf (ollama)", t_summarize,
        skip_reason="" if ollama_ok else "Ollama offline")

    # ── 24. Translate PDF ─────────────────────────────────────────────────────
    def t_translate():
        out = tmp / "translated.txt"
        pdf_tools.translate_pdf(pdf3, out, engine="ollama", language="Spanish", output_format="txt")
        assert_nonempty(out)
    run("translate_pdf (ollama)", t_translate,
        skip_reason="" if ollama_ok else "Ollama offline")

    # ── 25. PDF to PPTX ───────────────────────────────────────────────────────
    def t_pdf_to_pptx():
        out = tmp / "to_pptx.pptx"
        pdf_tools.pdf_to_pptx(pdf3, out, engine="basic")
        assert_nonempty(out)
    run("pdf_to_pptx (basic)", t_pdf_to_pptx)

    def t_pdf_to_pptx_fidelity():
        import zipfile
        out = tmp / "to_pptx_fidelity.pptx"
        pdf_tools.pdf_to_pptx(pdf3, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            slide_xml = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            media = [name for name in zf.namelist() if name.startswith("ppt/media/")]
            first_slide = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="replace")
        if len(slide_xml) != 3:
            raise ValueError(f"Expected 3 slide xml files, got {len(slide_xml)}")
        if len(media) != 3:
            raise ValueError(f"Expected 3 embedded slide images, got {len(media)}")
        if "<a:t>Page 1</a:t>" not in first_slide:
            raise ValueError("Expected editable text overlay content in fidelity slide")
    run("pdf_to_pptx (fidelity)", t_pdf_to_pptx_fidelity)

    def t_pdf_to_pptx_fidelity_two_columns():
        import zipfile
        cols_pdf = make_pdf_two_columns(tmp / "two_columns_ppt.pdf")
        out = tmp / "to_pptx_fidelity_columns.pptx"
        pdf_tools.pdf_to_pptx(cols_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="replace")
        required = ["Left A", "Left B", "Left C", "Right A", "Right B", "Right C"]
        missing = [value for value in required if f"<a:t>{value}</a:t>" not in slide_xml]
        if missing:
            raise ValueError(f"Two-column PPTX fidelity output is missing: {missing}")
        left_a = slide_xml.index("<a:t>Left A</a:t>")
        left_c = slide_xml.index("<a:t>Left C</a:t>")
        right_a = slide_xml.index("<a:t>Right A</a:t>")
        if not (left_a < left_c < right_a):
            raise ValueError(
                f"Unexpected two-column PPTX reading order: Left A idx={left_a}, Left C idx={left_c}, Right A idx={right_a}"
            )
    run("pdf_to_pptx (fidelity, two columns)", t_pdf_to_pptx_fidelity_two_columns)

    def t_pdf_to_pptx_fidelity_notes():
        import zipfile

        long_pdf = make_pdf_many_pages(tmp / "ppt_notes_long.pdf", count=10)
        out = tmp / "to_pptx_fidelity_notes.pptx"
        pdf_tools.pdf_to_pptx(long_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            note_names = [name for name in zf.namelist() if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")]
            if not note_names:
                raise ValueError("Expected long PPTX fidelity export to create notes slides")
            note_xml = zf.read(note_names[0]).decode("utf-8", errors="replace")
        if "Page 1" not in note_xml or "Long Report Page 1" not in note_xml:
            raise ValueError("Expected slide notes to include extracted editable text")
    run("pdf_to_pptx (fidelity, slide notes)", t_pdf_to_pptx_fidelity_notes)

    def t_pdf_to_pptx_fidelity_mixed_layout_notes():
        import zipfile

        mixed_pdf = make_pdf_mixed_layout(tmp / "mixed_layout_ppt.pdf")
        out = tmp / "to_pptx_fidelity_mixed.pptx"
        pdf_tools.pdf_to_pptx(mixed_pdf, out, engine="fidelity")
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            note_names = [name for name in zf.namelist() if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")]
            if not note_names:
                raise ValueError("Expected mixed-layout PowerPoint fidelity output to create notes slides")
            note_xml = zf.read(note_names[0]).decode("utf-8", errors="replace")
        if "Quarterly Summary" not in note_xml:
            raise ValueError("Expected mixed-layout PowerPoint notes to preserve page text")
    run("pdf_to_pptx (fidelity, mixed layout notes)", t_pdf_to_pptx_fidelity_mixed_layout_notes)

    # ── 27. Redact PDF ────────────────────────────────────────────────────────
    def t_redact():
        out = tmp / "redacted.pdf"
        # pdf3 page 1 has "Hello World" and "Invoice #12345"
        count = pdf_tools.redact_pdf(pdf3, out, phrases=["Hello World", "Invoice"])
        assert_nonempty(out)
        if count == 0:
            raise ValueError("Expected at least 1 redaction, got 0")
        print(f"       -> {count} redaction(s) applied")
    run("redact_pdf", t_redact)

    # ── 28. Compare PDF ───────────────────────────────────────────────────────
    def t_compare():
        out = tmp / "comparison.html"
        pdf_tools.compare_pdfs(pdf_a, pdf_b, out)
        assert_nonempty(out)
        content = out.read_text(encoding="utf-8")
        if "<table" not in content:
            raise ValueError("HTML diff table not found in output")
    run("compare_pdfs", t_compare)

    # ── 29a. Render Thumbnails ────────────────────────────────────────────────
    def t_thumbnails():
        out_dir = tmp / "thumbs"
        out_dir.mkdir(exist_ok=True)
        thumbs = pdf_tools.render_pdf_thumbnails(pdf3, out_dir, dpi=72)
        if len(thumbs) != 3:
            raise ValueError(f"Expected 3 thumbnails, got {len(thumbs)}")
        for p in thumbs:
            assert_nonempty(p)
    run("render_pdf_thumbnails", t_thumbnails)

    # ── 29b. Organize PDF ─────────────────────────────────────────────────────
    def t_organize():
        out = tmp / "organized.pdf"
        pdf_tools.organize_pdf(pdf3, out, page_order=[2, 0, 1])   # reverse then front
        assert_nonempty(out)
    run("organize_pdf", t_organize)

    # ── 30. Sign PDF ──────────────────────────────────────────────────────────
    def t_edit_pdf():
        import fitz

        out = tmp / "edited.pdf"
        pdf_tools.edit_pdf(
            pdf3,
            out,
            operations=[
                {"type": "whiteout", "page": 0, "x": 0.08, "y": 0.09, "w": 0.26, "h": 0.05},
                {
                    "type": "text",
                    "page": 0,
                    "x": 0.08,
                    "y": 0.09,
                    "w": 0.26,
                    "h": 0.07,
                    "text": "Replaced heading",
                    "font_size": 14,
                    "color": "#112233",
                    "align": "left",
                },
                {
                    "type": "image",
                    "page": 0,
                    "x": 0.62,
                    "y": 0.08,
                    "w": 0.14,
                    "h": 0.07,
                    "image_data": base64.b64encode(sig_bytes).decode("ascii"),
                },
            ],
        )
        assert_nonempty(out)
        with fitz.open(str(out)) as doc:
            text = doc[0].get_text("text")
            drawings = doc[0].get_drawings()
            images = doc[0].get_images(full=True)
        if "Replaced heading" not in text:
            raise ValueError("Expected edited PDF to contain inserted overlay text")
        if not drawings:
            raise ValueError("Expected edited PDF to contain at least one drawn overlay rectangle")
        if not images:
            raise ValueError("Expected edited PDF to contain at least one inserted image overlay")
    run("edit_pdf", t_edit_pdf)

    def t_edit_pdf_replace_text_blocks():
        import fitz

        sample = tmp / "live_text_replace_source.pdf"
        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        page.insert_text((72, 120), "Editable title block", fontsize=24, fontname="helv")
        doc.save(sample)
        doc.close()

        previews = pdf_tools.render_pdf_edit_previews(sample, tmp / "edit_previews_replace", dpi=96)
        if not previews or not previews[0].get("text_blocks"):
            raise ValueError("Expected edit preview metadata to include detected text blocks")

        block = previews[0]["text_blocks"][0]
        out = tmp / "edited_replace_text.pdf"
        pdf_tools.edit_pdf(
            sample,
            out,
            operations=[
                {
                    "type": "replace_text",
                    "page": 0,
                    "x": block["x"],
                    "y": block["y"],
                    "w": block["w"],
                    "h": block["h"],
                    "text": "Updated title block",
                    "font_size": block.get("font_size", 24),
                    "font_family": block.get("font_family", "helv"),
                    "bold": block.get("bold", False),
                    "italic": block.get("italic", False),
                    "color": block.get("color", "#111111"),
                    "align": "left",
                }
            ],
        )
        assert_nonempty(out)
        with fitz.open(str(out)) as edited:
            text = edited[0].get_text("text")
        if "Updated title block" not in text:
            raise ValueError("Expected replace_text edit to write updated text into the PDF")
    run("edit_pdf_replace_text", t_edit_pdf_replace_text_blocks)

    def t_sign():
        out = tmp / "signed.pdf"
        pdf_tools.sign_pdf(pdf3, out, sig_bytes, position="bottom-right", page_number=-1)
        assert_nonempty(out)
    run("sign_pdf", t_sign)

    # =========================================================================
    # SECTION B — Engine Fallback
    # =========================================================================
    print(f"\n{BOLD}Section B — Engine Fallback{RESET}\n")

    # B1. engine="basic" never calls _ollama_check
    def t_basic_no_ollama():
        from unittest.mock import patch
        sentinel = AssertionError("_ollama_check called when engine=basic")
        with patch("execution.pdf_tools._ollama_check", side_effect=sentinel):
            out = tmp / "b1_basic.txt"
            pdf_tools.pdf_to_text(pdf3, out, engine="basic")
            assert_nonempty(out)
    run("engine=basic  — never touches Ollama", t_basic_no_ollama)

    # B2. engine="auto" + Ollama running → produces output
    def t_auto_ollama_up():
        out = tmp / "b2_auto_up.txt"
        pdf_tools.pdf_to_text(pdf3, out, engine="auto")
        assert_nonempty(out)
    run("engine=auto   — Ollama up → uses Ollama", t_auto_ollama_up,
        skip_reason="" if ollama_ok else "Ollama offline")

    # B3. engine="auto" + Ollama down → falls back to basic silently
    def t_auto_ollama_down():
        from unittest.mock import patch
        with patch("execution.pdf_tools._ollama_check",
                   side_effect=RuntimeError("Ollama is not running. Start it with: ollama serve")):
            out = tmp / "b3_auto_fallback.txt"
            pdf_tools.pdf_to_text(pdf3, out, engine="auto")
            assert_nonempty(out)   # must still produce output via basic
    run("engine=auto   — Ollama down → silent fallback", t_auto_ollama_down)

    # B4. engine="ollama" + Ollama down → raises RuntimeError with clear message
    def t_ollama_explicit_down():
        from unittest.mock import patch
        with patch("execution.pdf_tools._ollama_check",
                   side_effect=RuntimeError("Ollama is not running. Start it with: ollama serve")):
            try:
                pdf_tools.pdf_to_text(pdf3, tmp / "b4_never.txt", engine="ollama")
                raise AssertionError("Expected RuntimeError — none was raised")
            except RuntimeError as exc:
                msg = str(exc).lower()
                if "ollama" not in msg:
                    raise AssertionError(f"Error message does not mention Ollama: {exc}")
    run("engine=ollama — Ollama down → clear RuntimeError", t_ollama_explicit_down)

    # =========================================================================
    # SECTION C — Edge Cases
    # =========================================================================
    # B5. engine="ollama" + runtime allocation failure → pdf_to_text falls back to extracted text
    def t_ollama_runtime_fallback_text():
        from unittest.mock import patch
        out = tmp / "b5_runtime_text.txt"
        with patch("execution.pdf_tools._render_pdf_pages", return_value=[tmp / "fake.png"]), \
             patch(
                 "execution.pdf_tools._ollama_call_page",
                 side_effect=RuntimeError(
                     "Ollama request failed (500): {\"error\":\"memory layout cannot be allocated\"}"
                 ),
             ):
            pdf_tools.pdf_to_text(pdf3, out, engine="ollama")
        assert_nonempty(out)
        content = out.read_text(encoding="utf-8")
        if "Invoice #12345" not in content:
            raise AssertionError("Expected basic text extraction fallback content in pdf_to_text output")
    run("engine=ollama — runtime fail → pdf_to_text fallback", t_ollama_runtime_fallback_text)

    # B6. summarize_pdf recovers to non-Ollama text AI when Ollama runtime fails
    def t_summarize_runtime_fallback_ai():
        from unittest.mock import patch
        out = tmp / "b6_summary.txt"
        with patch("execution.pdf_tools._render_pdf_pages", return_value=[tmp / "fake.png"]), \
             patch(
                 "execution.pdf_tools._ollama_call_page",
                 side_effect=RuntimeError(
                     "Ollama request failed (500): {\"error\":\"memory layout cannot be allocated\"}"
                 ),
             ), \
             patch(
                 "execution.pdf_tools._ollama_text",
                 side_effect=RuntimeError(
                     "Ollama request failed (500): {\"error\":\"memory layout cannot be allocated\"}"
                 ),
             ), \
             patch("execution.pdf_tools._fallback_text_ai", return_value="Fallback summary from alternate AI."):
            pdf_tools.summarize_pdf(pdf3, out, engine="ollama", output_format="txt")
        assert_nonempty(out)
        if "Fallback summary from alternate AI." not in out.read_text(encoding="utf-8"):
            raise AssertionError("Expected summarize_pdf to use non-Ollama fallback text AI")
    run("engine=ollama — runtime fail → summarize fallback", t_summarize_runtime_fallback_ai)

    # B7. translate_pdf recovers to non-Ollama text AI when Ollama runtime fails
    def t_translate_runtime_fallback_ai():
        from unittest.mock import patch
        out = tmp / "b7_translate.txt"
        translated_inputs: list[str] = []

        def fake_fallback(_system: str, user_content: str) -> str:
            translated_inputs.append(user_content)
            return "Texto traducido."

        with patch("execution.pdf_tools._render_pdf_pages", return_value=[tmp / "fake.png"]), \
             patch(
                 "execution.pdf_tools._ollama_call_page",
                 side_effect=RuntimeError(
                     "Ollama request failed (500): {\"error\":\"memory layout cannot be allocated\"}"
                 ),
             ), \
             patch(
                 "execution.pdf_tools._ollama_text",
                 side_effect=RuntimeError(
                     "Ollama request failed (500): {\"error\":\"memory layout cannot be allocated\"}"
                 ),
             ), \
             patch("execution.pdf_tools._fallback_text_ai", side_effect=fake_fallback):
            pdf_tools.translate_pdf(pdf3, out, engine="ollama", language="Spanish", output_format="txt")
        assert_nonempty(out)
        content = out.read_text(encoding="utf-8")
        if "Texto traducido." not in content or not translated_inputs:
            raise AssertionError("Expected translate_pdf to use non-Ollama fallback text AI")
    run("engine=ollama — runtime fail → translate fallback", t_translate_runtime_fallback_ai)

    print(f"\n{BOLD}Section C — Edge Cases{RESET}\n")

    # C1. Empty PDF (0 pages) → split returns 0 files (no crash, caller detects)
    def t_empty_pdf_no_crash():
        from pypdf import PdfWriter
        empty = tmp / "empty.pdf"
        writer = PdfWriter()
        with empty.open("wb") as f:
            writer.write(f)
        out_dir = tmp / "split_empty"
        out_dir.mkdir(exist_ok=True)
        pages = pdf_tools.split_pdf(empty, out_dir)
        # Acceptable: returns [] rather than crashing
        if pages:
            for p in pages:
                assert_nonempty(p)
    run("empty PDF (0 pages) — no crash", t_empty_pdf_no_crash)

    # C2. Encrypted PDF with wrong password → raises, message references password/decrypt
    def t_wrong_password():
        protected = tmp / "pw_test.pdf"
        pdf_tools.protect_pdf(pdf3, "correct123", protected)
        out = tmp / "pw_unlocked.pdf"
        try:
            pdf_tools.unlock_pdf(protected, "wrongpassword", out)
            raise AssertionError("Expected an error for wrong password — none raised")
        except AssertionError:
            raise
        except Exception as exc:
            # Any exception is fine; check message is somewhat informative
            msg = str(exc).lower()
            keywords = ("password", "incorrect", "decrypt", "wrong", "invalid", "not decrypted")
            if not any(k in msg for k in keywords):
                raise AssertionError(f"Error message not informative for wrong password: {exc!r}")
    run("wrong password      — clear error", t_wrong_password)

    # C3. Single-page PDF through split → exactly 1 file returned
    def t_split_single_page():
        single = tmp / "single_page.pdf"
        pdf_tools.extract_pages(pdf3, "1", single)
        out_dir = tmp / "split_single"
        out_dir.mkdir(exist_ok=True)
        pages = pdf_tools.split_pdf(single, out_dir)
        if len(pages) != 1:
            raise ValueError(f"Expected 1 file, got {len(pages)}")
        assert_nonempty(pages[0])
    run("single-page split   — 1 file returned", t_split_single_page)

    # C4. Zero-byte uploaded file → error before any PDF processing
    def t_zero_byte_file():
        zero = tmp / "zero.pdf"
        zero.write_bytes(b"")
        try:
            pdf_tools.merge_pdfs([zero], tmp / "zero_merged.pdf")
            raise AssertionError("Expected an error for zero-byte file — none raised")
        except AssertionError:
            raise
        except Exception:
            pass   # any exception satisfies the test
    run("zero-byte file      — error raised", t_zero_byte_file)

    # C5. Out-of-range page spec "1-999" on 3-page PDF → ValueError with range info
    def t_out_of_range_spec():
        try:
            pdf_tools.extract_pages(pdf3, "1-999", tmp / "range_out.pdf")
            raise AssertionError("Expected ValueError for out-of-range spec — none raised")
        except AssertionError:
            raise
        except ValueError as exc:
            msg = str(exc)
            if "999" not in msg and "exceed" not in msg.lower() and "range" not in msg.lower():
                raise AssertionError(f"Error message not informative: {exc!r}")
    run("page spec '1-999'   — clear ValueError", t_out_of_range_spec)

    # ── Summary ───────────────────────────────────────────────────────────────
    def t_pdf_to_svg():
        out = tmp / "to_svg.svg"
        pdf_tools.pdf_page_to_svg(pdf3, out, page_number=0, text_as_path=False)
        assert_nonempty(out)
        svg_text = out.read_text(encoding="utf-8", errors="replace")
        if "<svg" not in svg_text:
            raise ValueError("Expected SVG output to contain an <svg root element")
        if "<text" not in svg_text:
            raise ValueError("Expected SVG output to preserve text nodes")
    run("pdf_to_svg", t_pdf_to_svg)

    def t_pdf_to_svg_all_pages():
        out_dir = tmp / "svg_all"
        outputs = pdf_tools.pdf_to_svg(pdf3, out_dir, export_mode="all")
        if len(outputs) != 3:
            raise ValueError(f"Expected 3 SVG files from all-pages export, got {len(outputs)}")
        for out in outputs:
            assert_nonempty(out)
    run("pdf_to_svg (all pages)", t_pdf_to_svg_all_pages)

    def t_pdf_to_svg_small_page_pairing():
        import fitz

        src = tmp / "svg_card_three.pdf"
        doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE", "EXTRA PAGE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        doc.save(str(src))

        out_dir = tmp / "svg_paired"
        outputs = pdf_tools.pdf_to_svg(src, out_dir, export_mode="paired", auto_pair_small_pages=True)
        if len(outputs) != 2:
            raise ValueError(f"Expected 2 SVG outputs from paired export, got {len(outputs)}")
        paired_svg = outputs[0].read_text(encoding="utf-8", errors="replace")
        extra_svg = outputs[1].read_text(encoding="utf-8", errors="replace")
        if 'viewBox="0 0 522.00 144.00"' not in paired_svg:
            raise ValueError("Expected paired SVG export to combine two 3.5 x 2 inch pages plus the pairing gap")
        if "FRONT SIDE" not in paired_svg or "BACK SIDE" not in paired_svg:
            raise ValueError("Expected paired SVG export to include front and back content")
        if "EXTRA PAGE" not in extra_svg:
            raise ValueError("Expected extra page to remain as its own SVG output")
    run("pdf_to_svg (paired small pages)", t_pdf_to_svg_small_page_pairing)

    def t_pdf_to_word_business_card_size():
        import fitz
        import zipfile

        src = tmp / "business_card_size.pdf"
        doc = fitz.open()
        page = doc.new_page(width=252, height=144)
        page.insert_text((18, 34), "BUSINESS CARD", fontsize=16)
        page.insert_text((18, 58), "Jane Doe", fontsize=11)
        doc.save(str(src))

        out = tmp / "business_card_size.docx"
        pdf_tools.pdf_to_word(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")

        if 'w:pgSz w:w="5040" w:h="2880"' not in doc_xml and 'w:h="2880" w:w="5040"' not in doc_xml:
            raise ValueError("Expected Word fidelity export to preserve 3.5 x 2 inch page size")
    run("pdf_to_word (fidelity, business card size)", t_pdf_to_word_business_card_size)

    def t_pdf_to_word_business_card_no_appendix():
        import fitz
        import zipfile

        src = tmp / "business_card_appendix.pdf"
        doc = fitz.open()
        for label in ("FRONT SIDE", "BACK SIDE", "EXTRA SIDE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), label, fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_appendix.docx"
        pdf_tools.pdf_to_word(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")

        if "Editable Text Extract" in doc_xml:
            raise ValueError("Expected small-format Word fidelity export to stay appendix-free")
    run("pdf_to_word (fidelity, business card no appendix)", t_pdf_to_word_business_card_no_appendix)

    def t_pdf_to_pptx_business_card_size():
        import fitz
        import zipfile

        src = tmp / "business_card_size_ppt.pdf"
        doc = fitz.open()
        page = doc.new_page(width=252, height=144)
        page.insert_text((18, 34), "BUSINESS CARD", fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_size.pptx"
        pdf_tools.pdf_to_pptx(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            presentation_xml = zf.read("ppt/presentation.xml").decode("utf-8", errors="replace")

        if 'cx="3200400"' not in presentation_xml or 'cy="1828800"' not in presentation_xml:
            raise ValueError("Expected PowerPoint fidelity export to preserve 3.5 x 2 inch slide size")
    run("pdf_to_pptx (fidelity, business card size)", t_pdf_to_pptx_business_card_size)

    def t_pdf_to_word_business_card_pairing():
        import fitz
        import zipfile

        src = tmp / "business_card_pair_word.pdf"
        doc = fitz.open()
        front = doc.new_page(width=252, height=144)
        front.insert_text((18, 34), "FRONT SIDE", fontsize=16)
        back = doc.new_page(width=252, height=144)
        back.insert_text((18, 34), "BACK SIDE", fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_pair.docx"
        pdf_tools.pdf_to_word(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")

        if 'w:pgSz w:w="10440" w:h="2880"' not in doc_xml and 'w:h="2880" w:w="10440"' not in doc_xml:
            raise ValueError("Expected paired Word export to combine front and back card widths")
        if "<w:t>FRONT SIDE</w:t>" not in doc_xml or "<w:t>BACK SIDE</w:t>" not in doc_xml:
            raise ValueError("Expected paired Word export to keep editable overlays for front and back")
    run("pdf_to_word (fidelity, business card pairing)", t_pdf_to_word_business_card_pairing)

    def t_pdf_to_pptx_business_card_pairing():
        import fitz
        import zipfile

        src = tmp / "business_card_pair_ppt.pdf"
        doc = fitz.open()
        front = doc.new_page(width=252, height=144)
        front.insert_text((18, 34), "FRONT SIDE", fontsize=16)
        back = doc.new_page(width=252, height=144)
        back.insert_text((18, 34), "BACK SIDE", fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_pair.pptx"
        pdf_tools.pdf_to_pptx(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            presentation_xml = zf.read("ppt/presentation.xml").decode("utf-8", errors="replace")
            slide_names = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="replace")

        if len(slide_names) != 1:
            raise ValueError(f"Expected paired PowerPoint export to use one slide, got {len(slide_names)}")
        if 'cx="6629400"' not in presentation_xml or 'cy="1828800"' not in presentation_xml:
            raise ValueError("Expected paired PowerPoint export to combine front and back card widths")
        if "<a:t>FRONT SIDE</a:t>" not in slide_xml or "<a:t>BACK SIDE</a:t>" not in slide_xml:
            raise ValueError("Expected paired PowerPoint export to keep editable overlays for front and back")
    run("pdf_to_pptx (fidelity, business card pairing)", t_pdf_to_pptx_business_card_pairing)

    def t_pdf_to_word_business_card_three_pages():
        import fitz
        import zipfile

        src = tmp / "business_card_three_word.pdf"
        doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE", "EXTRA PAGE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_three.docx"
        pdf_tools.pdf_to_word(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")

        paired_ok = ('w:pgSz w:w="10440" w:h="2880"' in doc_xml) or ('w:h="2880" w:w="10440"' in doc_xml)
        single_ok = ('w:pgSz w:w="5040" w:h="2880"' in doc_xml) or ('w:h="2880" w:w="5040"' in doc_xml)
        if not paired_ok or not single_ok:
            raise ValueError("Expected Word fidelity export to pair first two card pages and keep the third separate")
        for value in ("FRONT SIDE", "BACK SIDE", "EXTRA PAGE"):
            if f"<w:t>{value}</w:t>" not in doc_xml:
                raise ValueError(f"Expected Word fidelity export to preserve editable overlay text for {value}")
    run("pdf_to_word (fidelity, business card 3 pages)", t_pdf_to_word_business_card_three_pages)

    def t_pdf_to_pptx_business_card_three_pages():
        import fitz
        import zipfile

        src = tmp / "business_card_three_ppt.pdf"
        doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE", "EXTRA PAGE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_three.pptx"
        pdf_tools.pdf_to_pptx(src, out, engine="fidelity")
        assert_nonempty(out)

        with zipfile.ZipFile(out) as zf:
            presentation_xml = zf.read("ppt/presentation.xml").decode("utf-8", errors="replace")
            slide_names = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            slide1 = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="replace")
            slide2 = zf.read("ppt/slides/slide2.xml").decode("utf-8", errors="replace")

        if len(slide_names) != 2:
            raise ValueError(f"Expected 2 slides for 3-page business card export, got {len(slide_names)}")
        if 'cx="6629400"' not in presentation_xml or 'cy="1828800"' not in presentation_xml:
            raise ValueError("Expected PowerPoint fidelity export to keep the paired card deck size")
        if "<a:t>FRONT SIDE</a:t>" not in slide1 or "<a:t>BACK SIDE</a:t>" not in slide1:
            raise ValueError("Expected the first slide to contain the paired front and back card overlays")
        if "<a:t>EXTRA PAGE</a:t>" not in slide2:
            raise ValueError("Expected the second slide to preserve the extra page overlay")
    run("pdf_to_pptx (fidelity, business card 3 pages)", t_pdf_to_pptx_business_card_three_pages)

    def t_pdf_to_word_business_card_pairing_disabled():
        import fitz
        import zipfile

        src = tmp / "business_card_pair_word_off.pdf"
        doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_pair_off.docx"
        pdf_tools.pdf_to_word(src, out, engine="fidelity", auto_pair_small_pages=False)
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            doc_xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
        if 'w:pgSz w:w="10440" w:h="2880"' in doc_xml or 'w:h="2880" w:w="10440"' in doc_xml:
            raise ValueError("Expected Word export with pairing disabled to keep pages separate")
    run("pdf_to_word (fidelity, pairing disabled)", t_pdf_to_word_business_card_pairing_disabled)

    def t_pdf_to_pptx_business_card_pairing_disabled():
        import fitz
        import zipfile

        src = tmp / "business_card_pair_ppt_off.pdf"
        doc = fitz.open()
        for text in ("FRONT SIDE", "BACK SIDE"):
            page = doc.new_page(width=252, height=144)
            page.insert_text((18, 34), text, fontsize=16)
        doc.save(str(src))

        out = tmp / "business_card_pair_off.pptx"
        pdf_tools.pdf_to_pptx(src, out, engine="fidelity", auto_pair_small_pages=False)
        assert_nonempty(out)
        with zipfile.ZipFile(out) as zf:
            presentation_xml = zf.read("ppt/presentation.xml").decode("utf-8", errors="replace")
            slide_names = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        if len(slide_names) != 2:
            raise ValueError(f"Expected 2 slides with pairing disabled, got {len(slide_names)}")
        if 'cx="6629400"' in presentation_xml:
            raise ValueError("Expected PowerPoint export with pairing disabled to keep single-card slide size")
    run("pdf_to_pptx (fidelity, pairing disabled)", t_pdf_to_pptx_business_card_pairing_disabled)

    passed = sum(1 for _, s, _, _ in results if s == "PASS")
    failed = sum(1 for _, s, _, _ in results if s == "FAIL")
    skipped = sum(1 for _, s, _, _ in results if s == "SKIP")
    total_time = sum(t for _, _, t, _ in results)

    print(f"\n{'-' * 56}")
    print(f"{BOLD}Results:{RESET}  "
          f"{GREEN}{passed} passed{RESET}  "
          f"{RED}{failed} failed{RESET}  "
          f"{YELLOW}{skipped} skipped{RESET}  "
          f"({total_time:.1f}s total)")

    if failed:
        print(f"\n{RED}Failed tests:{RESET}")
        for name, status, _, msg in results:
            if status == "FAIL":
                print(f"  • {name}: {msg[:200]}")
        print()
        sys.exit(1)
    else:
        print(f"\n{GREEN}All tests passed!{RESET}\n")
        sys.exit(0)


if __name__ == "__main__":
    main()
